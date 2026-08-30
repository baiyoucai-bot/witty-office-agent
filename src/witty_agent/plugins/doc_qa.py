"""本地文档/PPT 只读质检：空页、占位、过密要点、过小字号。不进内核循环。"""

from __future__ import annotations

import re
from pathlib import Path

from witty_agent import hooks
from witty_agent.prompts import get_prompt
from witty_agent.tools.registry import ToolSpec, register_tool

_PLACEHOLDER = re.compile(
    r"点击添加|Click to add|lorem ipsum|\bTODO\b|\bTBD\b|\bxxx\b|待补充|占位",
    re.IGNORECASE,
)
_TEXT_SUFFIX = {".md", ".txt", ".html"}
_MAX_BULLETS = 7
_MIN_FONT = 14
_LONG_PARA = 500


def _resolve(path: str) -> Path:
    target = Path(path).expanduser()
    if target.is_file():
        return target
    workspace = str(hooks.current_workspace or "").strip()
    if workspace and not target.is_absolute():
        alt = Path(workspace) / path
        if alt.is_file():
            return alt
    return target


def _inspect_pptx(target: Path) -> list[tuple[str, str]]:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError(get_prompt("pptx_missing")) from exc
    deck = Presentation(str(target))
    slides = list(deck.slides)
    if not slides:
        return [("-", get_prompt("doc_qa_no_slides"))]
    findings: list[tuple[str, str]] = []
    for index, slide in enumerate(slides, start=1):
        blobs: list[str] = []
        tiny = False
        for shape in slide.shapes:
            if getattr(shape, "name", "") == "witty-footer":
                continue
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = (para.text or "").strip()
                if text:
                    blobs.append(text)
                for run in para.runs:
                    size = run.font.size
                    if size is not None and float(size.pt) < _MIN_FONT:
                        tiny = True
        where = get_prompt("doc_qa_page", index=str(index))
        if not blobs:
            findings.append((where, get_prompt("doc_qa_empty_page")))
            continue
        if _PLACEHOLDER.search("\n".join(blobs)):
            findings.append((where, get_prompt("doc_qa_placeholder")))
        if len(blobs) - 1 > _MAX_BULLETS:
            findings.append((where, get_prompt("doc_qa_too_many", count=str(len(blobs) - 1))))
        if tiny:
            findings.append((where, get_prompt("doc_qa_tiny_font", size=str(_MIN_FONT))))
    return findings


def _inspect_text(target: Path) -> list[tuple[str, str]]:
    body = target.read_text(encoding="utf-8")
    if not body.strip():
        return [("-", get_prompt("doc_qa_empty_file"))]
    findings: list[tuple[str, str]] = []
    if not re.search(r"^#\s+\S", body, re.MULTILINE) and not re.search(r"<h[1-3]\b", body, re.IGNORECASE):
        findings.append(("-", get_prompt("doc_qa_no_heading")))
    if _PLACEHOLDER.search(body):
        findings.append(("-", get_prompt("doc_qa_placeholder")))
    for line in body.splitlines():
        if len(line.strip()) > _LONG_PARA:
            findings.append(("-", get_prompt("doc_qa_long_para", count=str(len(line.strip())))))
            break
    return findings


def doc_qa(path: str) -> str:
    """只读检查本地 PPTX 或 Markdown/文本的版式问题。

    Args:
        path: 本机文件路径
    """
    target = _resolve(path)
    if not target.is_file():
        return get_prompt("doc_qa_missing", path=str(target))
    suffix = target.suffix.lower()
    if suffix == ".pptx":
        findings = _inspect_pptx(target)
    elif suffix in _TEXT_SUFFIX:
        findings = _inspect_text(target)
    else:
        return get_prompt("doc_qa_unsupported", path=str(target), kind=suffix or "(none)")
    if not findings:
        return get_prompt("doc_qa_ok", path=str(target))
    rows = "\n".join(
        get_prompt("doc_qa_item", where=where, issue=issue) for where, issue in findings
    )
    return get_prompt("doc_qa_report", path=str(target), count=str(len(findings)), rows=rows)


register_tool(
    ToolSpec(
        name="doc_qa",
        description=get_prompt("tool_desc_doc_qa"),
        parameters={
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": get_prompt("doc_qa_param_path")},
            },
            "required": ["path"],
            "additionalProperties": False,
        },
        func=doc_qa,
    )
)
