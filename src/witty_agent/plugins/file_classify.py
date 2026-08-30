"""按调用方类型表做资料分类：路径证据优先、拆分件先合并、两轮判定。业务插件，不进内核循环。

判定口径全部在 config/prompts.toml 的 file_classify_* 里，本模块只负责扫描、分组、分批、落盘。
"""

from __future__ import annotations

import asyncio
import difflib
import hashlib
import json
import re
import time
from datetime import datetime
from collections.abc import Callable
from dataclasses import dataclass, field, replace
from pathlib import Path
from typing import Any

from witty_agent import hooks
from witty_agent.async_bridge import run_sync
from witty_agent.atomic_write import write_file_atomic
from witty_agent.logging import get_logger
from witty_agent.prompts import get_prompt
from witty_agent.tools.registry import ToolSpec, register_tool

logger = get_logger("file_classify")

_SKIP_DIRS = {
    ".git", ".svn", ".hg", "__pycache__", ".venv", "venv", "node_modules",
    ".idea", ".vscode", ".witty", ".ipynb_checkpoints",
}
_SKIP_NAMES = {".DS_Store", "Thumbs.db", "desktop.ini"}
_TEXT_SUFFIX = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".log", ".xml",
    ".htm", ".html", ".yaml", ".yml", ".ini", ".conf", ".rtf",
}
_IMAGE_SUFFIX = {
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp", ".heic", ".svg",
}
_OPAQUE_SUFFIX = {
    ".doc", ".xls", ".ppt", ".zip", ".rar", ".7z", ".tar", ".gz", ".dwg", ".dxf",
    ".exe", ".bin", ".mp3", ".mp4", ".avi", ".mov",
}

_ID_KEYS = ("id", "code", "category_id", "categoryid", "key", "value", "编号", "类型id", "分类id", "代码")
_NAME_KEYS = ("name", "category_name", "title", "label", "text", "名称", "类型", "分类", "类型名称", "分类名称")
_DESC_KEYS = ("description", "desc", "remark", "comment", "说明", "描述", "备注", "定义", "含义")
_PARENT_KEYS = ("parent", "parent_name", "group", "category", "父类", "上级", "所属", "大类", "一级分类")
_CHILD_KEYS = ("children", "items", "subs", "sub", "子类", "下级", "二级分类")

# 取主干里最后一段数字：1.png / 扫描件_01.pdf / 第3页.jpg / 报告-2-终稿.docx 都要认出来
_NUM_RUN = re.compile(r"^(.*?)(\d+)(\D*)$")
_TRIM = " \t-_.·、—～~()（）[]【】"
_WS = re.compile(r"[ \t\u3000]+")
_BLANK = re.compile(r"\n{3,}")
_UNITS = ("B", "KB", "MB", "GB")

# 组连贯性判定时，超大组只抽样这么多成员喂给模型，避免一个 500 页扫描件撑爆上下文
_GROUP_SAMPLE = 5
_EXCERPT_CACHE: dict[tuple[str, int, int], str] = {}
_CACHE_MAX = 4000


class _MissingDep(Exception):
    def __init__(self, dep: str) -> None:
        super().__init__(dep)
        self.dep = dep


@dataclass
class Category:
    id: str
    name: str
    desc: str = ""


@dataclass
class FileRec:
    path: Path
    rel: str
    size: int

    @property
    def name(self) -> str:
        return self.path.name

    @property
    def suffix(self) -> str:
        return self.path.suffix


@dataclass
class Unit:
    unit_id: str
    members: list[FileRec]
    group_id: str = ""

    @property
    def head(self) -> FileRec:
        return self.members[0]


@dataclass
class Verdict:
    unit_id: str
    category_id: str
    category_name: str
    confidence: float = 0.0
    need_content: bool = False
    evidence: list[str] = field(default_factory=list)
    reasoning: str = ""
    group_conflict: str = ""
    stage: str = "pass1"
    # 失败原因码，空串=成功。只在内部流转，结果行上翻译成 status/error 两个字段
    cause: str = ""


# 结果行只有两种结局：成功给结论，失败给原因（error 字段）。
# 没定论的单元不写行——半成品摆出去只会被当成结论误用。
ROW_OK = "ok"
ROW_FAILED = "failed"

# 内部失败原因码，决定 error 文案；不直接出现在结果行里
CAUSE_UNDECIDED = "undecided"  # 模型答了但判不出来，给了 _待分类
CAUSE_NO_VERDICT = "no_verdict"  # 各轮模型都没返回这个单元
CAUSE_BAD_JSON = "bad_json"  # 所在批次输出非法 JSON，修复/抢救/重试全失败


# ── 类型表 ────────────────────────────────────────────────────────────────


def _pick(entry: dict, keys: tuple[str, ...]) -> str:
    for key in keys:
        for actual in entry:
            if str(actual).casefold() == key.casefold():
                value = entry[actual]
                if isinstance(value, (str, int, float)) and str(value).strip():
                    return str(value).strip()
    return ""


def _flatten(entry: Any, parent: str, out: list[Category]) -> None:
    if isinstance(entry, str):
        text = entry.strip()
        if text:
            out.append(Category(id=text, name=text, desc=parent))
        return
    if isinstance(entry, list):
        for item in entry:
            _flatten(item, parent, out)
        return
    if not isinstance(entry, dict):
        return
    cid = _pick(entry, _ID_KEYS)
    name = _pick(entry, _NAME_KEYS) or cid
    cid = cid or name
    desc = _pick(entry, _DESC_KEYS)
    own_parent = _pick(entry, _PARENT_KEYS) or parent
    children: Any = None
    for key in _CHILD_KEYS:
        for actual in entry:
            if str(actual).casefold() == key.casefold() and isinstance(entry[actual], (list, dict)):
                children = entry[actual]
                break
        if children is not None:
            break
    if cid:
        parts = [part for part in (own_parent, desc) if part]
        out.append(Category(id=cid, name=name, desc=" / ".join(parts)))
    if children is not None:
        branch = " > ".join(part for part in (own_parent, name) if part)
        _flatten(children, branch, out)


def normalize_taxonomy(raw: Any) -> list[Category]:
    """把调用方给的类型表归一成 id/名称/说明。容忍数组、字典、嵌套子类、纯字符串列表。"""
    data = raw
    if isinstance(raw, str):
        text = raw.strip()
        # 只有不像 JSON 的短字符串才当路径试探，否则 stat 一个长 JSON 会抛 ENAMETOOLONG
        if text and text[0] not in "[{" and len(text) < 1024 and "\n" not in text:
            try:
                candidate = Path(text).expanduser()
                if candidate.is_file():
                    text = candidate.read_text(encoding="utf-8")
            except OSError:
                pass
        try:
            data = json.loads(text)
        except json.JSONDecodeError as exc:
            raise ValueError(str(exc)) from exc
    out: list[Category] = []
    if isinstance(data, dict):
        listish = None
        for key in ("categories", "types", "items", "data", "list", "类型", "分类", "类型表"):
            value = data.get(key)
            if isinstance(value, (list, dict)):
                listish = value
                break
        if listish is not None:
            _flatten(listish, "", out)
        else:
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    _flatten(value, str(key), out)
                else:
                    out.append(Category(id=str(key), name=str(value)))
    else:
        _flatten(data, "", out)
    seen: dict[str, Category] = {}
    for cat in out:
        if cat.id and cat.id not in seen:
            seen[cat.id] = cat
    return list(seen.values())


def _render_taxonomy(cats: list[Category]) -> str:
    rows = []
    for cat in cats:
        desc = get_prompt("file_classify_taxonomy_desc", text=cat.desc) if cat.desc else ""
        rows.append(get_prompt("file_classify_taxonomy_item", id=cat.id, name=cat.name, desc=desc))
    return "\n".join(rows)


# ── 扫描与候选分组 ────────────────────────────────────────────────────────


def _human_size(size: int) -> str:
    value = float(size)
    for unit in _UNITS:
        if value < 1024 or unit == _UNITS[-1]:
            return f"{value:.0f} {unit}" if unit == "B" else f"{value:.1f} {unit}"
        value /= 1024
    return f"{size} B"


def scan_files(root: Path) -> list[FileRec]:
    out: list[FileRec] = []
    for path in sorted(root.rglob("*")):
        if path.is_dir():
            continue
        if any(part in _SKIP_DIRS or part.startswith(".") for part in path.relative_to(root).parts[:-1]):
            continue
        if path.name in _SKIP_NAMES or path.name.startswith("."):
            continue
        try:
            size = path.stat().st_size
        except OSError:
            continue
        out.append(FileRec(path=path, rel=path.relative_to(root).as_posix(), size=size))
    return out


def _digest(prefix: str, parts: list[str]) -> str:
    raw = "\n".join(parts).encode("utf-8")
    return f"{prefix}_{hashlib.sha1(raw).hexdigest()[:10]}"


def candidate_groups(files: list[FileRec]) -> tuple[dict[str, list[FileRec]], list[FileRec]]:
    """同目录 + 同主干 + 仅编号不同 = 疑似同源。只给形式线索，是否真的同源由模型确认。"""
    buckets: dict[tuple[str, str, str, str], list[tuple[int, FileRec]]] = {}
    loose: list[FileRec] = []
    for rec in files:
        matched = _NUM_RUN.match(rec.path.stem)
        if not matched:
            loose.append(rec)
            continue
        prefix = matched.group(1).rstrip(_TRIM)
        tail = matched.group(3).lstrip(_TRIM)
        key = (str(rec.path.parent), prefix.casefold(), tail.casefold(), rec.suffix.casefold())
        buckets.setdefault(key, []).append((int(matched.group(2)), rec))
    groups: dict[str, list[FileRec]] = {}
    for members in buckets.values():
        if len(members) < 2:
            loose.extend(rec for _, rec in members)
            continue
        numbers = [num for num, _ in members]
        if max(numbers) - min(numbers) + 1 > 2 * len(members):
            loose.extend(rec for _, rec in members)
            continue
        ordered = [rec for _, rec in sorted(members, key=lambda item: item[0])]
        groups[_digest("g", [rec.rel for rec in ordered])] = ordered
    return groups, loose


# ── 正文抽取 ──────────────────────────────────────────────────────────────


def _squeeze(text: str) -> str:
    return _BLANK.sub("\n\n", _WS.sub(" ", text.replace("\r\n", "\n").replace("\r", "\n"))).strip()


def _pdf_text(path: Path, limit: int) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            from PyPDF2 import PdfReader  # type: ignore[no-redef]
        except ImportError as exc:
            raise _MissingDep("pypdf") from exc
    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
        if sum(len(part) for part in parts) >= limit:
            break
    return "\n".join(parts)


def _docx_text(path: Path, limit: int) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise _MissingDep("python-docx") from exc
    parts: list[str] = []
    for para in Document(str(path)).paragraphs:
        if para.text.strip():
            parts.append(para.text)
        if sum(len(part) for part in parts) >= limit:
            break
    return "\n".join(parts)


def _xlsx_text(path: Path, limit: int) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise _MissingDep("openpyxl") from exc
    book = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in book.worksheets:
        parts.append(f"[{sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell) for cell in row if cell is not None]
            if cells:
                parts.append(" | ".join(cells))
            if sum(len(part) for part in parts) >= limit:
                book.close()
                return "\n".join(parts)
    book.close()
    return "\n".join(parts)


def _pptx_text(path: Path, limit: int) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise _MissingDep("python-pptx") from exc
    parts: list[str] = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
        if sum(len(part) for part in parts) >= limit:
            break
    return "\n".join(parts)


def _sniff_text(path: Path, limit: int) -> str:
    blob = path.read_bytes()[: limit * 4]
    text = blob.decode("utf-8", errors="ignore")
    printable = sum(1 for ch in text if ch.isprintable() or ch in "\n\t")
    if not text or printable / max(len(text), 1) < 0.8:
        raise _MissingDep("")
    return text


def excerpt(rec: FileRec, limit: int) -> str:
    # 分组确认和第二轮判定会读到同一批文件，PDF/xlsx 解析一次就够
    key = (str(rec.path), rec.size, limit)
    hit = _EXCERPT_CACHE.get(key)
    if hit is not None:
        return hit
    text = _extract(rec, limit)
    if len(_EXCERPT_CACHE) < _CACHE_MAX:
        _EXCERPT_CACHE[key] = text
    return text


def _extract(rec: FileRec, limit: int) -> str:
    suffix = rec.suffix.casefold()
    kind = suffix or get_prompt("file_classify_kind_none")
    if suffix in _IMAGE_SUFFIX or suffix in _OPAQUE_SUFFIX:
        return get_prompt("file_classify_excerpt_binary", kind=kind)
    try:
        if suffix == ".pdf":
            text = _pdf_text(rec.path, limit)
        elif suffix == ".docx":
            text = _docx_text(rec.path, limit)
        elif suffix == ".xlsx":
            text = _xlsx_text(rec.path, limit)
        elif suffix == ".pptx":
            text = _pptx_text(rec.path, limit)
        elif suffix in _TEXT_SUFFIX:
            text = rec.path.read_text(encoding="utf-8", errors="replace")[: limit * 4]
        else:
            text = _sniff_text(rec.path, limit)
    except _MissingDep as exc:
        if not exc.dep:
            return get_prompt("file_classify_excerpt_binary", kind=kind)
        return get_prompt("file_classify_excerpt_dep", dep=exc.dep, kind=kind)
    except Exception as exc:  # 单个文件读不动不该让整批停下
        logger.warning("正文抽取失败 rel=%s err=%s", rec.rel, exc)
        return get_prompt("file_classify_excerpt_error", err=str(exc)[:120])
    clean = _squeeze(text)
    return clean[:limit] if clean else get_prompt("file_classify_excerpt_empty")


# ── 模型调用 ──────────────────────────────────────────────────────────────


# 资源池瞬时故障（超时/连接/429/5xx）后的默认重试间隔：等池子腾出资源，别原地打转
_RETRY_INTERVAL_SEC = 180.0
_RETRY_AFTER = re.compile(r"retry[-_ ]?after[:= ]+(\d+(?:\.\d+)?)", re.I)


def _retry_after_sec(text: str) -> float:
    """资源池在 Retry-After 里明确要求的等待秒数；没有则 0。"""
    hit = _RETRY_AFTER.search(text or "")
    return float(hit.group(1)) if hit else 0.0


class _PoolGate:
    """按资源池的节奏放行模型调用：任一批次撞墙就让整轮一起退避。

    并发批次会在同一时刻撞上打满的池子，各自独立退避等于对刚缓过来的池子反复
    冲击；共享一个「最早可发起时间」让所有在飞批次一起等，才是被池子调控。
    interval 为固定间隔（不做指数增长：目标是等资源腾出来，不是惩罚自己）。
    """

    def __init__(
        self,
        interval: float,
        *,
        on_wait: Callable[[str], None] | None = None,
        on_event: Callable[[dict], None] | None = None,
    ) -> None:
        self.interval = max(0.0, interval)
        self.on_wait = on_wait
        self.on_event = on_event
        self.waits = 0
        self.total_wait = 0.0
        self._not_before = 0.0

    def emit(self, payload: dict) -> None:
        """结构化事件：调用方据此在自己的接口里显示「正在等资源池」。

        文本进度会被下一批的进度消息覆盖，所以状态必须另有结构化通道。
        回调抛错不能拖垮分类。
        """
        if self.on_event is None:
            return
        try:
            self.on_event({**payload, "waits": self.waits, "total_wait_sec": round(self.total_wait, 1)})
        except Exception as exc:  # noqa: BLE001
            logger.warning("资源池事件上报失败 err=%s", exc)

    async def wait_turn(self) -> None:
        # 循环重取：等待期间别的批次又撞墙会把闸门推得更远
        while True:
            delay = self._not_before - time.monotonic()
            if delay <= 0:
                return
            await asyncio.sleep(delay)

    def penalize(self, retry_after: float = 0.0) -> float:
        """撞墙后推后闸门。池子明确要求等更久就听它的，但绝不早于配置间隔。"""
        delay = max(self.interval, retry_after)
        self._not_before = max(self._not_before, time.monotonic() + delay)
        self.waits += 1
        self.total_wait += delay
        return delay


async def _ask_available(
    system: str,
    user: str,
    *,
    max_tokens: int,
    timeout: int,
    think: str,
    gate: _PoolGate,
    max_attempts: int = 0,
) -> str:
    """等资源池放行再调模型；瞬时故障按固定间隔持续重试，鉴权/配额类立即失败。

    max_attempts=0 表示不限次数——资源池打满、网关超时都属于「等一会儿就好」，
    放弃只会让已烧掉的批次白费。鉴权错误和配额耗尽重试一万次也是错，直接抛。
    """
    from witty_agent.retry import is_retryable_error

    attempt = 0
    while True:
        await gate.wait_turn()
        try:
            reply = await _ask(system, user, max_tokens=max_tokens, timeout=timeout, think=think)
        except Exception as exc:
            if not is_retryable_error(exc):
                raise
            attempt += 1
            if max_attempts and attempt >= max_attempts:
                raise
            delay = gate.penalize(_retry_after_sec(str(exc)))
            logger.warning(
                "模型侧瞬时故障，%.0f 秒后重试（第 %d 次）err=%s", delay, attempt, str(exc)[:200]
            )
            text = get_prompt(
                "file_classify_pool_wait",
                delay=f"{delay:.0f}",
                n=str(attempt),
                err=str(exc)[:120],
            )
            if gate.on_wait is not None:
                gate.on_wait(text)
            gate.emit(
                {
                    "event": "pool_wait",
                    "attempt": attempt,
                    "delay_sec": round(delay, 1),
                    "error": str(exc)[:300],
                    "message": text,
                }
            )
            continue
        if attempt:
            gate.emit({"event": "pool_recovered", "attempt": attempt})
        return reply


async def _ask(system: str, user: str, *, max_tokens: int, timeout: int, think: str) -> str:
    from witty_agent.llm import OpenAICompatLLM
    from witty_agent.types import AgentContext, AgentMessage, ModelRef

    llm = OpenAICompatLLM(stream=False, timeout=timeout, max_tokens=max_tokens, retry_attempts=2)
    llm.think_level = think
    context = AgentContext(
        system_prompt=system,
        messages=[AgentMessage(role="user", content=user)],
        tools=[],
        workspace_dir="",
        model=ModelRef(provider="openai", model_id=llm.model_id),
        project_id="",
        agent_id="file-classify",
        session_id="file-classify",
    )
    message = await llm(context)
    if message.stop_reason == "error":
        raise RuntimeError(message.text() or "file_classify llm error")
    return message.text()


# 模型返回非法 JSON 时整批重试的次数；HTTP 层的可重试状态码在 llm.py 里另算
_JSON_RETRIES = 2
# 连续多少个批次重试耗尽仍无法解析，就判定模型不可用、中止整轮（防止对废模型烧完全部批次）
_BREAKER_LIMIT = 3


class BatchParseError(ValueError):
    """单批 JSON 重试耗尽。调用方降级本批继续跑，不该让它杀死整轮分类。"""


def _repair_json_text(text: str) -> str:
    """修复弱模型两类高频 JSON 笔误：字符串里未转义的内层引号、裸换行/制表符。

    判定引号是否为字符串终结符：其后首个非空白字符属于 ,:}] 或到文末才算终结，
    否则视为正文里的引号并转义。实测 qwq 在 evidence 里引用原文时会写出
    "content: "原文标题"" 这种结构，逐字符走一遍即可救回。
    """
    out: list[str] = []
    in_str = False
    escaped = False
    for i, ch in enumerate(text):
        if not in_str:
            if ch == '"':
                in_str = True
            out.append(ch)
            continue
        if escaped:
            out.append(ch)
            escaped = False
            continue
        if ch == "\\":
            out.append(ch)
            escaped = True
            continue
        if ch == "\n":
            out.append("\\n")
            continue
        if ch == "\t":
            out.append("\\t")
            continue
        if ch == '"':
            j = i + 1
            while j < len(text) and text[j] in " \t\r\n":
                j += 1
            if j >= len(text) or text[j] in ",:}]":
                in_str = False
                out.append(ch)
            else:
                out.append('\\"')
            continue
        out.append(ch)
    return "".join(out)


def _salvage_records(text: str) -> list[dict]:
    """从整体非法的回复里抢救完整的结果对象（带 unit_id/group_id 的 dict）。

    逐个 { 起点试 raw_decode：坏掉的只是局部（多为某条字符串没闭合），
    前后完整的记录仍能解出来，缺的那几条走既有漏判兜底。
    """
    decoder = json.JSONDecoder()
    records: list[dict] = []
    idx = 0
    while True:
        start = text.find("{", idx)
        if start < 0:
            break
        try:
            obj, end = decoder.raw_decode(text, start)
        except ValueError:
            idx = start + 1
            continue
        if isinstance(obj, dict) and ("unit_id" in obj or "group_id" in obj):
            records.append(obj)
            idx = end
        else:
            idx = start + 1
    return records


def _match_replies(expected: list[str], items: list[dict], key: str) -> dict[str, dict]:
    """把模型回的行按 id 对回本批的期望 id，容忍模型把 id 抄错一两个字符。

    弱模型常在回填 digest id 时抄错个别字符（实测 u_82fd436192 答成 u_82rd436192），
    精确匹配会把这些行当成漏判、让单元白白多走一轮。期望 id 是随机摘要、彼此差异
    极大，所以相似度兜底（cutoff 0.8 ≈ 容忍 12 位 id 错 2 字符）几乎无误配风险。
    """
    found: dict[str, dict] = {}
    leftovers: list[tuple[str, dict]] = []
    expected_set = set(expected)
    for item in items:
        raw = str(item.get(key, ""))
        if raw in expected_set and raw not in found:
            found[raw] = item
        else:
            leftovers.append((raw, item))
    for raw, item in leftovers:
        missing = [uid for uid in expected if uid not in found]
        if not missing or not raw:
            break
        close = difflib.get_close_matches(raw, missing, n=1, cutoff=0.8)
        if close:
            logger.info("修复模型抄错的 %s：%s -> %s", key, raw, close[0])
            found[close[0]] = item
    return found


async def _ask_json(
    system: str,
    user: str,
    *,
    max_tokens: int,
    timeout: int,
    think: str,
    record: Callable[[dict], None] | None = None,
    breaker: dict | None = None,
    gate: _PoolGate | None = None,
    retry_max_attempts: int = 0,
) -> dict:
    """调模型并解析 JSON；非法先修复/抢救，仍失败整批重试，重试耗尽抛 BatchParseError。

    record 每次真实调用（含失败的重试）各回一条：提示词原文、模型原始回复、耗时、
    解析结果（修复/抢救的会带 parse 字段）。它是执行过程的转录，异常只记 WARN 不中断。

    breaker 是整轮共享的熔断计数：单批重试耗尽只抛 BatchParseError（调用方降级本批），
    连续 _BREAKER_LIMIT 批都废才抛普通 ValueError 中止整轮——模型彻底不可用时别烧完全部批次。

    gate 在这一层之外：JSON 非法是模型内容问题（立刻重试有意义），资源池故障是
    基础设施问题（要等），两者的重试节奏必须分开。
    """
    for attempt in range(_JSON_RETRIES + 1):
        started = time.monotonic()
        if gate is None:
            raw = await _ask(system, user, max_tokens=max_tokens, timeout=timeout, think=think)
        else:
            raw = await _ask_available(
                system, user, max_tokens=max_tokens, timeout=timeout, think=think,
                gate=gate, max_attempts=retry_max_attempts,
            )
        entry = {
            "ts": datetime.now().astimezone().isoformat(timespec="seconds"),
            "attempt": attempt + 1,
            "duration_ms": int((time.monotonic() - started) * 1000),
            "system": system,
            "user": user,
            "reply": raw,
        }
        try:
            payload, parse_mode = _parse_lenient(raw)
        except ValueError as exc:
            entry["ok"] = False
            entry["error"] = str(exc)[:500]
            _record_call(record, entry)
            if attempt >= _JSON_RETRIES:
                if breaker is None:
                    raise
                breaker["streak"] += 1
                breaker["degraded"] += 1
                if breaker["streak"] >= _BREAKER_LIMIT:
                    raise ValueError(
                        get_prompt(
                            "file_classify_breaker_tripped",
                            n=str(breaker["streak"]),
                            err=str(exc)[:200],
                        )
                    ) from exc
                raise BatchParseError(str(exc)) from exc
            logger.warning(
                "模型返回非法 JSON，整批重试 %d/%d：%s",
                attempt + 1,
                _JSON_RETRIES,
                str(exc)[:200],
            )
            continue
        entry["ok"] = True
        if parse_mode != "ok":
            entry["parse"] = parse_mode
            logger.warning(
                "模型 JSON 非法但已%s继续", "修复" if parse_mode == "repaired" else "抢救部分记录"
            )
        _record_call(record, entry)
        if breaker is not None:
            breaker["streak"] = 0
        return payload
    raise RuntimeError("unreachable")


def _record_call(record: Callable[[dict], None] | None, entry: dict) -> None:
    if record is None:
        return
    try:
        record(entry)
    except Exception as exc:  # noqa: BLE001 转录失败不该拖垮分类本身
        logger.warning("执行过程转录失败 err=%s", exc)


async def _gather_limited(tasks: list, concurrency: int) -> list:
    """并发跑但卡上限；用 return_exceptions 让在飞的批次都落完盘再抛第一个错。"""
    gate = asyncio.Semaphore(max(1, concurrency))

    async def guarded(coro):
        async with gate:
            return await coro

    done = await asyncio.gather(*(guarded(task) for task in tasks), return_exceptions=True)
    for item in done:
        if isinstance(item, BaseException):
            raise item
    return done


async def _excerpts(recs: list[FileRec], limit: int) -> dict[str, str]:
    """正文抽取会解析 PDF/xlsx，是真阻塞，必须落到线程里，别卡事件循环。"""
    texts = await asyncio.gather(*(asyncio.to_thread(excerpt, rec, limit) for rec in recs))
    return {rec.rel: text for rec, text in zip(recs, texts)}


def _strip_fences(raw: str) -> str:
    text = (raw or "").strip()
    if "```" in text:
        start = text.find("```")
        chunk = text[start + 3 :]
        if chunk.lstrip().casefold().startswith("json"):
            chunk = chunk.lstrip()[4:]
        end = chunk.find("```")
        text = (chunk[:end] if end >= 0 else chunk).strip()
    if not text.startswith("{"):
        head, tail = text.find("{"), text.rfind("}")
        if head >= 0 and tail > head:
            text = text[head : tail + 1]
    return text


def _parse_json(raw: str) -> dict:
    text = _strip_fences(raw)
    try:
        parsed = json.loads(text)
    except json.JSONDecodeError as exc:
        raise ValueError(get_prompt("file_classify_bad_json", err=str(exc), head=raw[:200])) from exc
    if not isinstance(parsed, dict):
        raise ValueError(get_prompt("file_classify_bad_json", err="not an object", head=raw[:200]))
    return parsed


def _parse_lenient(raw: str) -> tuple[dict, str]:
    """严格解析失败就先修笔误再抢救记录。返回 (payload, 解析方式 ok/repaired/salvaged)。"""
    try:
        return _parse_json(raw), "ok"
    except ValueError as strict_err:
        text = _strip_fences(raw)
        repaired = _repair_json_text(text)
        if repaired != text:
            try:
                parsed = json.loads(repaired)
                if isinstance(parsed, dict):
                    return parsed, "repaired"
            except json.JSONDecodeError:
                pass
        records = _salvage_records(repaired)
        if records:
            payload = {
                "results": [item for item in records if "unit_id" in item],
                "groups": [item for item in records if "group_id" in item],
            }
            return payload, "salvaged"
        raise strict_err


def _batched(items: list, size: int) -> list[list]:
    return [items[i : i + size] for i in range(0, len(items), size)]


# ── 单元渲染 ──────────────────────────────────────────────────────────────


def _segments(rec: FileRec) -> str:
    parts = Path(rec.rel).parts[:-1]
    return " > ".join(parts) if parts else get_prompt("file_classify_root_segment")


def _members_line(unit: Unit) -> str:
    if len(unit.members) < 2:
        return ""
    names = ", ".join(rec.name for rec in unit.members)
    return get_prompt("file_classify_unit_members", count=str(len(unit.members)), names=names) + "\n"


def _title_line(text: str) -> str:
    """把正文开头压成一行标题行。抽不到文本时给占位说明，让模型知道核验不了。"""
    flat = " / ".join(part.strip() for part in text.splitlines() if part.strip())
    if not flat:
        return ""
    return get_prompt("file_classify_unit_title", text=flat) + "\n"


def _render_unit_pass1(unit: Unit, titles: dict[str, str]) -> str:
    head = unit.head
    return get_prompt(
        "file_classify_unit_pass1",
        unit_id=unit.unit_id,
        path=head.rel,
        segments=_segments(head),
        name=head.name,
        kind=head.suffix or get_prompt("file_classify_kind_none"),
        size=_human_size(sum(rec.size for rec in unit.members)),
        members=_members_line(unit),
        title=_title_line(titles.get(head.rel, "")),
    )


# 目录总览的截断上限：给模型全局观，但不能让超大项目把每批提示词撑爆
_TREE_MAX_DIRS = 200


def _render_tree(files: list[FileRec]) -> str:
    """整个项目的目录总览（每个目录一行，带直属文件数），按路径排序。

    这是模型的「全局观」：单元只带自己的完整路径，看不出项目整体怎么组织归档；
    总览让模型先知道有哪些标段、哪些业务环节目录，再判个体归属。
    """
    counts: dict[str, int] = {}
    for rec in files:
        folder = str(Path(rec.rel).parent)
        if folder == ".":
            folder = ""
        counts[folder] = counts.get(folder, 0) + 1
    folders = sorted(counts)
    lines = [
        get_prompt(
            "file_classify_tree_item",
            folder=folder or get_prompt("file_classify_tree_root"),
            count=str(counts[folder]),
        )
        for folder in folders[:_TREE_MAX_DIRS]
    ]
    if len(folders) > _TREE_MAX_DIRS:
        lines.append(
            get_prompt("file_classify_tree_truncated", count=str(len(folders) - _TREE_MAX_DIRS))
        )
    return "\n".join(lines)


_PROJECT_FIELDS = (("name", "name"), ("code", "code"), ("category", "category"))
# 容忍调用方直接把中文键名传进来，省一层映射
_PROJECT_ALIASES = {"项目名称": "name", "项目编号": "code", "项目分类": "category", "编号": "code"}
_PROJECT_VALUE_MAX = 200


def _render_project(project: dict[str, Any] | None) -> str:
    """项目背景段；调用方没给或全是空值就返回空串，让提示词里整段消失。

    作用不是多给点信息，而是让模型能分辨「本项目的过程文件」和「引用进来的外部
    项目材料」——投标目录下别的项目的合同是业绩证明，这个判断没有项目名就只能靠猜。
    """
    if not project:
        return ""
    picked: dict[str, str] = {}
    for key, value in project.items():
        field = _PROJECT_ALIASES.get(str(key).strip(), str(key).strip().lower())
        text = str(value).strip() if value is not None else ""
        if field in {name for name, _ in _PROJECT_FIELDS} and text:
            picked[field] = text[:_PROJECT_VALUE_MAX]
    if not picked:
        return ""
    items = [
        get_prompt("file_classify_project_item", label=get_prompt(f"file_classify_project_{field}"), value=picked[field])
        for field, _ in _PROJECT_FIELDS
        if field in picked
    ]
    return get_prompt("file_classify_project_block", items="\n".join(items))


def _unit_exhausted(unit: Unit, texts: dict[str, str], limit: int) -> bool:
    """采样成员的摘录都没顶到上限，说明正文已经读完，再加大也不会有新内容。"""
    return all(len(texts.get(rec.rel, "")) < limit for rec in unit.members[:_GROUP_SAMPLE])


def _render_unit_pass2(unit: Unit, prior: Verdict, texts: dict[str, str], limit: int) -> str:
    head = unit.head
    blocks = [
        get_prompt("file_classify_group_member", name=rec.name, excerpt=texts.get(rec.rel, ""))
        for rec in unit.members[:_GROUP_SAMPLE]
    ]
    note = (
        get_prompt("file_classify_excerpt_full_note")
        if _unit_exhausted(unit, texts, limit)
        else get_prompt("file_classify_excerpt_partial_note", chars=str(limit))
    )
    guess = f"{prior.category_id} {prior.category_name}".strip()
    return get_prompt(
        "file_classify_unit_pass2",
        unit_id=unit.unit_id,
        path=head.rel,
        segments=_segments(head),
        name=head.name,
        kind=head.suffix or get_prompt("file_classify_kind_none"),
        size=_human_size(sum(rec.size for rec in unit.members)),
        members=_members_line(unit),
        guess=guess,
        confidence=f"{prior.confidence:.2f}",
        reasoning=prior.reasoning,
        note=note,
        excerpt="\n".join(blocks),
    )


# ── 分组确认 ──────────────────────────────────────────────────────────────


async def _confirm_groups(
    groups: dict[str, list[FileRec]],
    *,
    batch: int,
    limit: int,
    max_tokens: int,
    timeout: int,
    think: str,
    concurrency: int,
    cached: dict[str, dict],
    on_batch: Callable[[str, dict], None] | None,
    record: Callable[[dict], None] | None = None,
    breaker: dict | None = None,
    gate: _PoolGate | None = None,
    retry_max_attempts: int = 0,
) -> dict[str, list[list[FileRec]]]:
    """返回 group_id -> 若干份资料，每份是一组成员。模型说拆就拆，说合就合。"""
    decided: dict[str, list[list[FileRec]]] = {}
    pending: list[tuple[str, list[FileRec]]] = []
    for gid, members in groups.items():
        hit = cached.get(gid)
        if hit is None:
            pending.append((gid, members))
            continue
        decided[gid] = _apply_group(members, hit)
    system = get_prompt("file_classify_group_system")

    async def one(chunk: list[tuple[str, list[FileRec]]]) -> None:
        sampled = [rec for _, members in chunk for rec in members[:_GROUP_SAMPLE]]
        texts = await _excerpts(sampled, limit)
        rendered = [
            get_prompt(
                "file_classify_group_item",
                group_id=gid,
                folder=str(Path(members[0].rel).parent),
                count=str(len(members)),
                members="\n".join(
                    get_prompt("file_classify_group_member", name=rec.name, excerpt=texts.get(rec.rel, ""))
                    for rec in members[:_GROUP_SAMPLE]
                ),
            )
            for gid, members in chunk
        ]
        user = get_prompt("file_classify_group_user", groups="\n\n".join(rendered))
        try:
            payload = await _ask_json(
                system, user, max_tokens=max_tokens, timeout=timeout, think=think,
                record=record, breaker=breaker, gate=gate, retry_max_attempts=retry_max_attempts,
            )
        except BatchParseError as exc:
            # 组确认降级 = 全按默认合并处理；合并策略保守，错了第二轮读正文还能纠
            logger.warning("组确认批次降级（非法 JSON），按默认合并：groups=%d err=%s", len(chunk), str(exc)[:200])
            payload = {"groups": []}
        by_id = _match_replies(
            [gid for gid, _ in chunk],
            [item for item in payload.get("groups", []) if isinstance(item, dict)],
            "group_id",
        )
        for gid, members in chunk:
            decision = by_id.get(gid) or {"merge": True, "reason": "", "confidence": 0.0}
            decided[gid] = _apply_group(members, decision)
            if on_batch is not None:
                on_batch(gid, decision)

    await _gather_limited([one(chunk) for chunk in _batched(pending, batch)], concurrency)
    return decided


def _apply_group(members: list[FileRec], record: dict) -> list[list[FileRec]]:
    by_rel = {rec.rel: rec for rec in members}
    by_name = {rec.name: rec for rec in members}
    if record.get("merge", True):
        return [members]
    parts: list[list[FileRec]] = []
    used: set[str] = set()
    for bucket in record.get("split") or []:
        if not isinstance(bucket, list):
            continue
        picked = []
        for item in bucket:
            rec = by_rel.get(str(item)) or by_name.get(Path(str(item)).name)
            if rec is not None and rec.rel not in used:
                picked.append(rec)
                used.add(rec.rel)
        if picked:
            parts.append(picked)
    for rec in members:
        if rec.rel not in used:
            parts.append([rec])
    return parts or [members]


# ── 主流程 ────────────────────────────────────────────────────────────────


def _default_out_dir(root: Path) -> Path:
    from witty_agent.layout import data_root

    return data_root() / "file_classify" / root.name


def _resolve_root(path: str) -> Path:
    target = Path(path).expanduser()
    if target.is_dir():
        return target.resolve()
    workspace = str(hooks.current_workspace or "").strip()
    if workspace and not target.is_absolute():
        alt = Path(workspace) / path
        if alt.is_dir():
            return alt.resolve()
    return target


def _load_jsonl(path: Path) -> dict[str, dict]:
    if not path.is_file():
        return {}
    out: dict[str, dict] = {}
    for line in path.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            record = json.loads(line)
        except json.JSONDecodeError:
            continue
        key = str(record.get("unit_id") or record.get("group_id") or "")
        if key:
            out[key] = record
    return out


def _append_jsonl(path: Path, records: list[dict]) -> None:
    if not records:
        return
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("a", encoding="utf-8") as stream:
        for record in records:
            stream.write(json.dumps(record, ensure_ascii=False) + "\n")


def _needs_second(verdict: Verdict, min_confidence: float) -> bool:
    """模型自己说要读正文，或它自己都没把握，就进第二轮。

    第一轮判「待分类」的也一律进第二轮：光凭路径和文件名就放弃太草率，
    读完正文仍毫无指向才允许收 _待分类——待分类必须是最后手段。
    """
    return (
        verdict.need_content
        or verdict.confidence < min_confidence
        or verdict.cause == CAUSE_UNDECIDED
    )


def _verdict_from(record: dict, unit: Unit, stage: str) -> Verdict:
    raw_conf = record.get("confidence", 0)
    try:
        confidence = float(raw_conf)
    except (TypeError, ValueError):
        confidence = 0.0
    evidence = record.get("evidence") or []
    if isinstance(evidence, str):
        evidence = [evidence]
    category_id = str(record.get("category_id") or get_prompt("file_classify_unassigned_id"))
    # 未定论的行也要有可显示的名称，否则调用方界面上是一片空白
    fallback_name = (
        get_prompt("file_classify_unassigned_name")
        if category_id == get_prompt("file_classify_unassigned_id")
        else ""
    )
    return Verdict(
        unit_id=unit.unit_id,
        category_id=category_id,
        category_name=str(record.get("category_name") or fallback_name),
        confidence=confidence,
        need_content=bool(record.get("need_content")),
        evidence=[str(item) for item in evidence if str(item).strip()],
        reasoning=str(record.get("reasoning") or ""),
        group_conflict=str(record.get("group_conflict") or ""),
        stage=stage,
        # 构造合成记录的地方会显式给 cause；模型自己回 _待分类 的算判不出来
        cause=str(record.get("cause") or "")
        or (
            CAUSE_UNDECIDED
            if category_id == get_prompt("file_classify_unassigned_id")
            else ""
        ),
    )


def _finalize_unresolved(verdict: Verdict, rounds: int, degraded: str) -> Verdict:
    """收口一个到最后一轮仍未拿到判定的单元：把终局原因补进 reasoning。

    沿用的是上一轮的结论，它的 reasoning 往往写着「已转入下一轮复判」——到这里
    已经没有下一轮了。调用方要靠这行字判断是模型没答（可重跑）还是输出非法
    （模型有问题），所以两种原因要分开写明，而不是笼统一句「未定论」。
    """
    why = (
        get_prompt("file_classify_unresolved_badjson", error=degraded)
        if degraded
        else get_prompt("file_classify_unresolved_missed")
    )
    note = get_prompt("file_classify_unresolved_final", rounds=str(rounds), why=why)
    prior = verdict.reasoning.strip()
    return replace(
        verdict,
        reasoning=f"{prior} {note}".strip(),
        # 终局原因覆盖沿用结论上的旧 cause：这一行到底为什么没定论，以这里为准
        cause=CAUSE_BAD_JSON if degraded else CAUSE_NO_VERDICT,
    )


async def aclassify_directory(
    root: str | Path,
    taxonomy: Any,
    *,
    out_dir: str | Path | None = None,
    limit: int = 0,
    concurrency: int = 4,
    pass1_batch: int = 15,
    pass2_batch: int = 6,
    group_batch: int = 5,
    excerpt_chars: int = 1200,
    title_chars: int = 160,
    max_tokens: int = 6000,
    timeout: int = 300,
    think: str = "off",
    min_confidence: float = 0.6,
    group_check: bool = True,
    content_rounds: int = 3,
    max_excerpt_chars: int = 20000,
    resume: bool = True,
    retry_interval: float = _RETRY_INTERVAL_SEC,
    retry_max_attempts: int = 0,
    project: dict[str, Any] | None = None,
    progress: Callable[[str], None] | None = None,
    on_result: Callable[[list[dict]], None] | None = None,
    on_retry: Callable[[dict], None] | None = None,
) -> dict[str, Any]:
    """扫描 root，按 taxonomy 给每个单元定一个类型。返回统计字典，结果落在 out_dir。

    批次并发跑，`concurrency` 卡住同时在飞的模型调用数。同步调用方用 `classify_directory`。
    on_result 每批判完调一次，只给已定论的行（status=ok/failed）；没判完的单元不出行。

    第一轮除路径和文件名外，还给每个单元头成员的正文开头 title_chars 字作「标题行」，
    用来核验文件名是否名实相符——文件名写错的文件靠它被察觉并转入第二轮。置 0 关闭。

    第二轮读正文由模型动态决定深度：首轮给每文件前 excerpt_chars 字，模型判「正文未完
    且不足以定论」（need_content=true）的单元，下一轮给 4 倍正文，直到读完全文、达到
    max_excerpt_chars 或耗尽 content_rounds 轮，最后一轮强制定论。

    project 是本项目的背景（`{"name": ..., "code": ..., "category": ...}`，也吃中文键名），
    写进每批提示词，让模型能分辨本项目自身的过程文件与引用进来的外部项目材料；
    不传则提示词里整段不出现。它只是背景参考，与路径、正文证据冲突时以后者为准。

    执行过程转录在 out_dir/calls.jsonl：每次模型调用一行（stage/round/提示词原文/
    原始回复/耗时/是否解析成功），非法 JSON 的重试也各占一行，跑的过程中即可读。

    模型侧瞬时故障（超时、连接断、429、5xx）按 retry_interval 秒固定间隔持续重试，
    retry_max_attempts=0 表示不限次数——资源池打满属于「等一会儿就好」，放弃只会
    让已烧掉的批次白费。整轮共享一个退避闸门：一个批次撞墙，所有在飞批次一起等，
    不去反复冲击刚缓过来的池子。鉴权失败、配额耗尽不重试，立即抛。

    等待状态有两条通道：progress 收到人话文本（会被后续进度覆盖），on_retry 收到
    结构化事件 {event: pool_wait|pool_recovered, attempt, delay_sec, error, waits,
    total_wait_sec}，调用方据此在自己的接口里长期展示「正在等模型资源」。
    """
    say = progress or (lambda _text: None)
    target = _resolve_root(str(root))
    if not target.is_dir():
        raise FileNotFoundError(get_prompt("file_classify_no_root", path=str(target)))
    try:
        cats = normalize_taxonomy(taxonomy)
    except ValueError as exc:
        raise ValueError(get_prompt("file_classify_bad_taxonomy", err=str(exc))) from exc
    if not cats:
        raise ValueError(get_prompt("file_classify_no_taxonomy"))
    files = scan_files(target)
    if not files:
        raise ValueError(get_prompt("file_classify_no_files", path=str(target)))

    out = Path(out_dir).expanduser() if out_dir else _default_out_dir(target)
    out.mkdir(parents=True, exist_ok=True)
    result_file, group_file, report_file = out / "results.jsonl", out / "groups.jsonl", out / "report.md"
    calls_file = out / "calls.jsonl"
    if not resume:
        result_file.unlink(missing_ok=True)
        group_file.unlink(missing_ok=True)
        calls_file.unlink(missing_ok=True)

    # 执行过程转录：每次模型调用（含非法 JSON 的重试）各落一行，供事后追查口径
    call_count = {"n": 0}
    # 非法 JSON 熔断：单批重试耗尽降级继续，连续 _BREAKER_LIMIT 批全废才中止整轮
    breaker = {"streak": 0, "degraded": 0}
    # 资源池闸门：整轮共享退避节奏，等待情况经 progress 报给调用方
    gate = _PoolGate(retry_interval, on_wait=say, on_event=on_retry)

    def _call_recorder(stage: str, round_no: int = 0) -> Callable[[dict], None]:
        def record(entry: dict) -> None:
            call_count["n"] += 1
            _append_jsonl(calls_file, [{"seq": call_count["n"], "stage": stage, "round": round_no, **entry}])

        return record

    groups, loose = candidate_groups(files)
    units: list[Unit] = []
    merged = 0
    if groups and group_check:
        say(f"候选拆分组 {len(groups)} 个，交模型确认连贯性")
        decided = await _confirm_groups(
            groups,
            batch=group_batch,
            limit=excerpt_chars,
            max_tokens=max_tokens,
            timeout=timeout,
            think=think,
            concurrency=concurrency,
            cached=_load_jsonl(group_file) if resume else {},
            on_batch=lambda gid, rec: _append_jsonl(group_file, [{**rec, "group_id": gid}]),
            record=_call_recorder("group"),
            breaker=breaker,
            gate=gate,
            retry_max_attempts=retry_max_attempts,
        )
        for gid, parts in decided.items():
            for members in parts:
                unit = Unit(unit_id=_digest("u", [rec.rel for rec in members]), members=members, group_id=gid)
                units.append(unit)
                if len(members) > 1:
                    merged += 1
    else:
        loose = list(files)
    units.extend(Unit(unit_id=_digest("u", [rec.rel]), members=[rec]) for rec in loose)
    units.sort(key=lambda item: item.head.rel)
    if limit > 0:
        units = units[:limit]

    # 结果文件里只有终局行（未定论的不落盘），所以有行就算已完成
    done = _load_jsonl(result_file) if resume else {}
    todo = [unit for unit in units if unit.unit_id not in done]
    if done:
        say(get_prompt("file_classify_resumed", count=str(len(units) - len(todo))))

    taxonomy_text = _render_taxonomy(cats)
    # 总览按全部扫描到的文件算（不受 limit 影响）：全局观就该是整个项目的
    tree_text = _render_tree(files)
    project_text = _render_project(project)
    system = get_prompt("file_classify_system")
    verdicts: dict[str, Verdict] = {}
    by_id = {unit.unit_id: unit for unit in units}

    pass1_chunks = _batched(todo, pass1_batch)
    counter = {"n": 0}

    async def run_pass1(chunk: list[Unit]) -> None:
        # 标题行只抽头成员的正文开头一小段：核验文件名名实相符（孤证不定案），不是读全文
        titles = (
            await _excerpts([unit.head for unit in chunk], title_chars) if title_chars > 0 else {}
        )
        user = get_prompt(
            "file_classify_user_pass1",
            project=project_text,
            taxonomy=taxonomy_text,
            tree=tree_text,
            units="\n\n".join(_render_unit_pass1(unit, titles) for unit in chunk),
            count=str(len(chunk)),
        )
        try:
            payload = await _ask_json(
                system, user, max_tokens=max_tokens, timeout=timeout, think=think,
                record=_call_recorder("pass1"), breaker=breaker,
                gate=gate, retry_max_attempts=retry_max_attempts,
            )
            degraded = ""
        except BatchParseError as exc:
            # 单批 JSON 重试耗尽：降级本批转正文轮复判，错误原因写进结果行，不杀整轮
            degraded = str(exc)[:300]
            logger.warning("第一轮批次降级（非法 JSON）units=%d err=%s", len(chunk), degraded[:200])
            payload = {"results": []}
        found = _match_replies(
            [unit.unit_id for unit in chunk],
            [item for item in payload.get("results", []) if isinstance(item, dict)],
            "unit_id",
        )
        for unit in chunk:
            record = found.get(unit.unit_id)
            if record is None:
                if degraded:
                    record = {
                        "need_content": True,
                        "reasoning": get_prompt("file_classify_batch_degraded", error=degraded),
                        "cause": CAUSE_BAD_JSON,
                    }
                else:
                    logger.warning("模型漏判单元 unit_id=%s", unit.unit_id)
                    record = {
                        "need_content": True,
                        "reasoning": get_prompt("file_classify_missing_verdict"),
                        "cause": CAUSE_NO_VERDICT,
                    }
            verdicts[unit.unit_id] = _verdict_from(record, unit, "pass1")
        # 每批落盘：几百个文件跑一半崩了，重跑要能接上，不能从头再烧一遍。
        # 只写已定论的；要进正文轮的单元这轮不出行，等它定论那轮再写
        settled = [
            _row(unit, verdicts[unit.unit_id], target)
            for unit in chunk
            if not _needs_second(verdicts[unit.unit_id], min_confidence)
        ]
        _append_jsonl(result_file, settled)
        _emit(on_result, settled, "pass1")
        counter["n"] += 1
        say(f"第一轮 {counter['n']}/{len(pass1_chunks)} 批完成：{len(chunk)} 个单元")

    if pass1_chunks:
        await _gather_limited([run_pass1(chunk) for chunk in pass1_chunks], concurrency)

    # 第二轮依赖第一轮的完整结果，两个阶段之间必须串行；阶段内部照样并发。
    # 读多少正文由模型逐轮决定：摘录不足以定论就 need_content=true，下一轮给 4 倍。
    pending = [by_id[uid] for uid, verdict in verdicts.items() if _needs_second(verdict, min_confidence)]
    pending.sort(key=lambda item: item.head.rel)
    total_rounds = max(1, content_rounds)

    for round_no in range(1, total_rounds + 1):
        if not pending:
            break
        limit_now = min(excerpt_chars * 4 ** (round_no - 1), max(excerpt_chars, max_excerpt_chars))
        # 上限顶死后再加轮次也拿不到新内容，直接当最后一轮
        last_round = round_no >= total_rounds or limit_now >= max_excerpt_chars
        round_chunks = _batched(pending, pass2_batch)
        counter["n"] = 0
        next_pending: list[Unit] = []

        async def run_pass2(chunk: list[Unit], *, limit_now: int = limit_now, last_round: bool = last_round, total: int = len(round_chunks), round_no: int = round_no) -> None:
            texts = await _excerpts([rec for unit in chunk for rec in unit.members[:_GROUP_SAMPLE]], limit_now)
            user = get_prompt(
                "file_classify_user_pass2",
                project=project_text,
                taxonomy=taxonomy_text,
                tree=tree_text,
                units="\n\n".join(
                    _render_unit_pass2(unit, verdicts[unit.unit_id], texts, limit_now) for unit in chunk
                ),
                count=str(len(chunk)),
                policy=get_prompt("file_classify_pass2_final" if last_round else "file_classify_pass2_more"),
            )
            try:
                payload = await _ask_json(
                    system, user, max_tokens=max_tokens, timeout=timeout, think=think,
                    record=_call_recorder("pass2", round_no), breaker=breaker,
                    gate=gate, retry_max_attempts=retry_max_attempts,
                )
                degraded = ""
            except BatchParseError as exc:
                degraded = str(exc)[:300]
                logger.warning(
                    "正文轮批次降级（非法 JSON）round=%d units=%d err=%s", round_no, len(chunk), degraded[:200]
                )
                payload = {"results": []}
            found = _match_replies(
                [unit.unit_id for unit in chunk],
                [item for item in payload.get("results", []) if isinstance(item, dict)],
                "unit_id",
            )
            settled = []
            for unit in chunk:
                record = found.get(unit.unit_id)
                if record is None:
                    if not degraded:
                        logger.warning("正文轮漏判单元 unit_id=%s round=%d", unit.unit_id, round_no)
                    if last_round:
                        # 最后一轮也漏，就以现有结论收口，不许有单元没有终局行。
                        # 沿用的结论里写着「已转入下一轮复判」，到这儿已经是终局，
                        # 必须补一句真实原因，否则调用方看到的理由是过期的。
                        verdicts[unit.unit_id] = _finalize_unresolved(
                            verdicts[unit.unit_id], round_no, degraded
                        )
                        settled.append(_row(unit, verdicts[unit.unit_id], target))
                    else:
                        next_pending.append(unit)
                    continue
                verdicts[unit.unit_id] = _verdict_from(record, unit, "pass2")
                verdict = verdicts[unit.unit_id]
                needs_more = (
                    verdict.need_content
                    and not last_round
                    and not _unit_exhausted(unit, texts, limit_now)
                )
                # 还要继续读正文的这轮不出行，免得半成品被实时消费方落库
                if needs_more:
                    next_pending.append(unit)
                else:
                    settled.append(_row(unit, verdict, target))
            _append_jsonl(result_file, settled)
            _emit(on_result, settled, "pass2")
            counter["n"] += 1
            say(f"正文第 {round_no} 轮 {counter['n']}/{total} 批完成：{len(chunk)} 个单元（每文件至多 {limit_now} 字）")

        await _gather_limited([run_pass2(chunk) for chunk in round_chunks], concurrency)
        if last_round:
            break
        next_pending.sort(key=lambda item: item.head.rel)
        pending = next_pending

    # 每个单元只在定论那一刻写一行；读回按 unit_id 去重只是兜底
    rows = list(_load_jsonl(result_file).values())
    write_file_atomic(report_file, _render_report(rows, cats, target))

    # tally 只数成功行：把没判成的算进类型分布，等于把失败报成成功
    tally: dict[str, int] = {}
    failed = 0
    for row in rows:
        if row.get("status") == ROW_FAILED:
            failed += 1
        else:
            tally[str(row.get("category_id"))] = tally.get(str(row.get("category_id")), 0) + 1
    summary = {
        "root": str(target),
        "categories": len(cats),
        "files": len(files),
        "units": len(units),
        "merged": merged,
        "pass1": sum(1 for row in rows if row.get("stage") == "pass1"),
        "pass2": sum(1 for row in rows if row.get("stage") == "pass2"),
        # 成功数就是 units_ok，不用调用方自己减；失败原因在每行的 error 字段
        "units_ok": len(rows) - failed,
        "units_failed": failed,
        "result_file": str(result_file),
        "report_file": str(report_file),
        "calls_file": str(calls_file),
        "model_calls": call_count["n"],
        "degraded_batches": breaker["degraded"],
        "pool_waits": gate.waits,
        "pool_wait_sec": round(gate.total_wait, 1),
        "tally": tally,
    }
    logger.info(
        "分类完成 root=%s units=%s ok=%s failed=%s",
        target,
        summary["units"],
        summary["units_ok"],
        summary["units_failed"],
    )
    return summary


def classify_directory(root: str | Path, taxonomy: Any, **kwargs: Any) -> dict[str, Any]:
    """`aclassify_directory` 的同步包装，给脚本和库调用方用。参数完全一致。

    已经在事件循环里的调用方直接 await `aclassify_directory`，别调本函数。
    """
    return run_sync(aclassify_directory(root, taxonomy, **kwargs), entry="aclassify_directory")


_CAUSE_PROMPTS = {
    CAUSE_UNDECIDED: "file_classify_error_undecided",
    CAUSE_NO_VERDICT: "file_classify_error_no_verdict",
    CAUSE_BAD_JSON: "file_classify_error_bad_json",
}


def _row(unit: Unit, verdict: Verdict, root: Path) -> dict[str, Any]:
    """只对已定论的单元产出行：没判完的不写，半成品摆出去只会被当结论用。"""
    failed = bool(verdict.cause)
    return {
        "unit_id": verdict.unit_id,
        "category_id": verdict.category_id,
        "category_name": verdict.category_name,
        "confidence": round(verdict.confidence, 3),
        "stage": verdict.stage,
        # 每行两个判定字段就够了：status 成败、error 失败原因（成功为 null）。
        # 别让调用方去 reasoning 里找字样，那是给人读的，不是判定接口。
        "status": ROW_FAILED if failed else ROW_OK,
        "error": get_prompt(_CAUSE_PROMPTS[verdict.cause]) if failed else None,
        "evidence": verdict.evidence,
        "reasoning": verdict.reasoning,
        "group_conflict": verdict.group_conflict,
        "group_id": unit.group_id,
        "members": [rec.rel for rec in unit.members],
        "root": str(root),
    }


def _emit(on_result: Callable[[list[dict]], None] | None, rows: list[dict], stage: str) -> None:
    if on_result is None or not rows:
        return
    try:
        on_result(rows)
    except Exception as exc:  # 下游回调崩了不该让整轮白跑，结果已落盘、可回放
        logger.warning("on_result 回调失败 stage=%s count=%s err=%s", stage, len(rows), exc)


def _render_report(rows: list[dict], cats: list[Category], root: Path) -> str:
    names = {cat.id: cat.name for cat in cats}
    buckets: dict[str, list[dict]] = {}
    for row in rows:
        buckets.setdefault(str(row.get("category_id")), []).append(row)
    lines = [get_prompt("file_classify_md_header", root=str(root), units=str(len(rows))), ""]
    for cid in sorted(buckets, key=lambda key: (-len(buckets[key]), key)):
        label = names.get(cid) or str(buckets[cid][0].get("category_name") or cid)
        lines.append(get_prompt("file_classify_md_section", name=label, id=cid, count=str(len(buckets[cid]))))
        lines.append("")
        for row in sorted(buckets[cid], key=lambda item: (item.get("members") or [""])[0]):
            members = row.get("members") or [""]
            title = (
                members[0]
                if len(members) == 1
                else get_prompt("file_classify_md_merged_title", first=members[0], count=str(len(members)))
            )
            failed_note = (
                get_prompt("file_classify_md_failed", error=str(row.get("error") or ""))
                if row.get("status") == ROW_FAILED
                else ""
            )
            lines.append(
                get_prompt(
                    "file_classify_md_unit",
                    title=title,
                    confidence=str(row.get("confidence")),
                    stage=str(row.get("stage")),
                    status=failed_note,
                )
            )
            if row.get("evidence"):
                lines.append(get_prompt("file_classify_md_evidence", text="; ".join(row["evidence"])))
            if row.get("reasoning"):
                lines.append(get_prompt("file_classify_md_reason", text=str(row["reasoning"])))
            if row.get("group_conflict"):
                lines.append(get_prompt("file_classify_md_conflict", text=str(row["group_conflict"])))
        lines.append("")
    return "\n".join(lines)


# ── 工具入口 ──────────────────────────────────────────────────────────────


async def classify_files(root: str, taxonomy: str, out_dir: str = "", limit: int = 0) -> str:
    # 工具在内核循环里被 await（loop.py 认 __await__），不能走同步包装去 asyncio.run
    summary = await aclassify_directory(root, taxonomy, out_dir=out_dir or None, limit=limit)
    names: dict[str, str] = {}
    for row in _load_jsonl(Path(summary["result_file"])).values():
        names.setdefault(str(row.get("category_id")), str(row.get("category_name") or ""))
    top = sorted(summary["tally"].items(), key=lambda item: -item[1])[:12]
    rows = "\n".join(
        get_prompt("file_classify_report_row", name=names.get(cid, cid), id=cid, count=str(count))
        for cid, count in top
    )
    return get_prompt(
        "file_classify_report",
        root=summary["root"],
        categories=str(summary["categories"]),
        files=str(summary["files"]),
        units=str(summary["units"]),
        merged=str(summary["merged"]),
        pass1=str(summary["pass1"]),
        pass2=str(summary["pass2"]),
        ok=str(summary["units_ok"]),
        failed=str(summary["units_failed"]),
        result_file=summary["result_file"],
        report_file=summary["report_file"],
        top=str(len(top)),
        rows=rows,
    )


register_tool(
    ToolSpec(
        name="classify_files",
        description=get_prompt("tool_desc_classify_files"),
        parameters={
            "type": "object",
            "properties": {
                "root": {"type": "string", "description": get_prompt("classify_files_param_root")},
                "taxonomy": {"type": "string", "description": get_prompt("classify_files_param_taxonomy")},
                "out_dir": {"type": "string", "description": get_prompt("classify_files_param_out_dir")},
                "limit": {"type": "integer", "description": get_prompt("classify_files_param_limit")},
            },
            "required": ["root", "taxonomy"],
            "additionalProperties": False,
        },
        func=classify_files,
    )
)
