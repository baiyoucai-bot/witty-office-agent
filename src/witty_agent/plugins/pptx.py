"""离线 PPTX：本地版式（字号/配色/页脚），不走公网模板。"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from witty_agent.logging import get_logger
from witty_agent.plugins.pptx_kit import (
    Deck,
    Slide,
    append_html,
    append_slide,
    append_slides,
    format_issues,
    insert_slide,
    lint_deck,
    parse_deck,
    parse_html,
    parse_slide,
    replace_slide,
    resolve_theme,
    themes_text,
    write_html,
    write_pptx,
)
from witty_agent.plugins.pptx_kit.render import atomic_save
from witty_agent.prompts import get_prompt
from witty_agent.tools.registry import ToolSpec, register_tool

logger = get_logger("plugins.pptx")

_INK = (0x1C, 0x1C, 0x22)
_FONT = "Microsoft YaHei"


def _pptx():
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Emu, Inches, Pt
    except ImportError as exc:  # pragma: no cover - 依赖未装时
        raise RuntimeError(get_prompt("pptx_missing")) from exc
    return Presentation, RGBColor, MSO_SHAPE, MSO_ANCHOR, PP_ALIGN, Emu, Inches, Pt


def _rgb(mod, triple):
    return mod(*triple)


def _run_style(frame):
    size = None
    color = None
    bold = None
    font = None
    align = None
    if frame.paragraphs:
        align = frame.paragraphs[0].alignment
    for para in frame.paragraphs:
        for run in para.runs:
            if run.font.name:
                font = run.font.name
            if run.font.size:
                size = run.font.size
            if run.font.bold is not None:
                bold = run.font.bold
            try:
                color = run.font.color.rgb
            except (AttributeError, TypeError, ValueError):
                pass
            if size is not None or color is not None:
                return size, color, bold, font, align
    return size, color, bold, font, align


def _write(
    frame,
    text: str,
    *,
    pt_mod,
    rgb_mod,
    align,
    size: int,
    color,
    bold: bool = False,
    font: str = _FONT,
):
    existing_size, existing_color, existing_bold, existing_font, existing_align = _run_style(frame)
    frame.clear()
    para = frame.paragraphs[0]
    para.alignment = existing_align or align
    run = para.add_run()
    run.text = text
    run.font.name = existing_font or font
    run.font.size = existing_size or pt_mod(size)
    run.font.bold = existing_bold if existing_bold is not None else bold
    if existing_color is not None:
        run.font.color.rgb = existing_color
    else:
        run.font.color.rgb = _rgb(rgb_mod, color)


def _bullets_of(raw: str) -> list[str]:
    return [item.strip().lstrip("-•· ").strip() for item in raw.splitlines() if item.strip()][:7]


def _fill_bullets(frame, lines: list[str], *, pt_mod, rgb_mod, align) -> None:
    size, color, _bold, font, existing_align = _run_style(frame)
    frame.clear()
    frame.word_wrap = True
    if not lines:
        return
    for index, line in enumerate(lines):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = existing_align or align.LEFT
        para.level = 0
        para.space_after = pt_mod(10)
        run = para.add_run()
        run.text = line
        run.font.name = font or _FONT
        run.font.size = size or pt_mod(20)
        if color is not None:
            run.font.color.rgb = color
        else:
            run.font.color.rgb = _rgb(rgb_mod, _INK)


def pptx_create(path: str, title: str, subtitle: str = "") -> str:
    deck = Deck(title=title, slides=[Slide(kind="cover", title=title, subtitle=subtitle)])
    return write_pptx(deck, path)


def pptx_add_slide(path: str, title: str, bullets: str) -> str:
    return append_slide(path, Slide(kind="bullets", title=title, items=_bullets_of(bullets)))


def _named(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name and shape.has_text_frame:
            return shape
    return None


def _content_body(slide):
    title = _named(slide, "witty-title")
    body = _named(slide, "witty-body")
    if title is not None:
        return title, body
    boxes = [
        shape
        for shape in slide.shapes
        if shape.has_text_frame and shape.name not in {"witty-footer", "witty-page", "witty-theme"}
    ]
    if not boxes:
        return None, None
    boxes.sort(key=lambda item: item.top)
    return boxes[0], boxes[1] if len(boxes) > 1 else None


def pptx_edit_slide(path: str, index: int, title: str = "", bullets: str = "") -> str:
    Presentation, RGBColor, _shape, _anchor, PP_ALIGN, _emu, Inches, Pt = _pptx()
    target = Path(path).expanduser()
    deck = Presentation(str(target))
    slides = list(deck.slides)
    if index < 1 or index > len(slides):
        raise ValueError(get_prompt("pptx_bad_index", index=str(index), total=str(len(slides))))
    slide = slides[index - 1]
    head, body = _content_body(slide)
    if title and head is not None:
        _write(
            head.text_frame,
            title,
            pt_mod=Pt,
            rgb_mod=RGBColor,
            align=PP_ALIGN.LEFT,
            size=26,
            color=_INK,
            bold=True,
        )
        if index == 1:
            deck.core_properties.title = title
    if bullets.strip() and body is not None:
        _fill_bullets(body.text_frame, _bullets_of(bullets), pt_mod=Pt, rgb_mod=RGBColor, align=PP_ALIGN)
    atomic_save(deck, target)
    return str(target)


def pptx_add_picture(path: str, image_path: str, title: str = "") -> str:
    image = Path(image_path).expanduser()
    if not image.is_file():
        raise FileNotFoundError(get_prompt("pptx_image_missing", path=str(image)))
    return append_slide(path, Slide(kind="picture", title=title, image=str(image)))


def pptx_outline(path: str) -> str:
    Presentation, _rgb, _shape, _anchor, _align, _emu, _inches, _pt = _pptx()
    deck = Presentation(str(Path(path).expanduser()))
    pages: list[str] = []
    for index, slide in enumerate(deck.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.name in {"witty-footer", "witty-page", "witty-theme"}:
                continue
            if getattr(shape, "has_chart", False):
                label = "[图表]"
                try:
                    chart = shape.chart
                    if chart.has_title:
                        title_text = chart.chart_title.text_frame.text.strip()
                        if title_text:
                            label = f"[图表] {title_text}"
                except (AttributeError, ValueError, TypeError):
                    pass
                texts.append(label)
                continue
            if shape.has_text_frame:
                blob = "\n".join(para.text.strip() for para in shape.text_frame.paragraphs if para.text.strip())
                if blob:
                    texts.append(blob)
        pages.append(f"## {index}\n" + ("\n".join(texts) or get_prompt("pptx_empty_page")))
    return "\n\n".join(pages) or get_prompt("pptx_empty")


def _read_deck_file(target: Path) -> Deck:
    text = target.read_text(encoding="utf-8")
    if target.suffix.lower() in {".html", ".htm"}:
        return parse_html(text)
    return parse_deck(text)


def _load_deck(deck: str, deck_path: str) -> Deck:
    file_arg = (deck_path or "").strip()
    if file_arg:
        target = Path(file_arg).expanduser()
        if not target.is_file():
            raise FileNotFoundError(get_prompt("pptx_html_missing", path=str(target)))
        return _read_deck_file(target)
    raw = (deck or "").strip()
    if not raw:
        raise ValueError(get_prompt("pptx_bad_json"))
    stripped = raw.lstrip()
    if stripped.startswith("<"):
        return parse_html(raw)
    if stripped.startswith("{") or stripped.startswith("["):
        return parse_deck(raw)
    target = Path(raw).expanduser()
    try:
        is_file = target.is_file()
    except OSError:
        is_file = False
    if is_file:
        return _read_deck_file(target)
    return parse_deck(raw)


def _apply_theme(payload: Deck, theme: str, theme_overrides: str) -> Deck:
    if theme.strip():
        payload.theme = theme.strip()
    extra = (theme_overrides or "").strip()
    if extra:
        try:
            parsed = json.loads(extra)
        except json.JSONDecodeError as exc:
            raise ValueError(get_prompt("pptx_bad_json")) from exc
        if not isinstance(parsed, dict):
            raise ValueError(get_prompt("pptx_bad_json"))
        payload.theme_overrides.update({str(key): str(value) for key, value in parsed.items()})
    return payload


def pptx_themes() -> str:
    return themes_text()


def pptx_render(
    path: str,
    deck: str = "",
    deck_path: str = "",
    theme: str = "",
    theme_overrides: str = "",
    preview: bool = True,
) -> str:
    payload = _apply_theme(_load_deck(deck, deck_path), theme, theme_overrides)
    resolved = resolve_theme(payload.theme, payload.theme_overrides)
    issues = lint_deck(payload, resolved)
    pptx_path = write_pptx(payload, path, resolved)
    html_path = ""
    if preview:
        html_path = write_html(payload, str(Path(path).expanduser().with_suffix(".html")), resolved)
    logger.info("整稿渲染 pptx=%s html=%s pages=%s issues=%s", pptx_path, html_path or "-", len(payload.slides), len(issues))
    return _render_message(pptx_path, payload, resolved, html_path, issues)


def _render_message(pptx_path: str, payload: Deck, resolved, html_path: str, issues) -> str:
    # 封面/章节的宏版式是引擎排的模板页，不提醒；只有内容页用宏才算偷懒。
    macro_pages = sum(1 for slide in payload.slides if not slide.boxes and slide.kind not in {"cover", "section"})
    msg = get_prompt(
        "pptx_render_ok",
        pptx=pptx_path,
        pages=str(len(payload.slides)),
        theme=resolved.id,
        html=html_path or get_prompt("pptx_no_preview"),
    )
    if macro_pages:
        msg += " " + get_prompt("pptx_macro_note", n=str(macro_pages)).lstrip()
    if issues:
        msg += "\n" + get_prompt("pptx_lint_note", count=str(len(issues)), issues=format_issues(issues)).lstrip()
    return msg


def pptx_check(deck: str = "", deck_path: str = "", theme: str = "", theme_overrides: str = "") -> str:
    payload = _apply_theme(_load_deck(deck, deck_path), theme, theme_overrides)
    resolved = resolve_theme(payload.theme, payload.theme_overrides)
    issues = lint_deck(payload, resolved)
    if not issues:
        return get_prompt("pptx_lint_clean", pages=str(len(payload.slides)))
    return get_prompt("pptx_lint_report", count=str(len(issues)), issues=format_issues(issues))


def pptx_snapshot(deck: str = "", deck_path: str = "", out: str = "", theme: str = "", theme_overrides: str = "") -> str:
    """整稿画成 PNG 联络表：交稿前自己看一眼版式。纯 Pillow，不开浏览器。"""
    from witty_agent.plugins.pptx_kit.raster import render_deck_png

    payload = _apply_theme(_load_deck(deck, deck_path), theme, theme_overrides)
    resolved = resolve_theme(payload.theme, payload.theme_overrides)
    target = (out or "").strip()
    if not target:
        source = (deck_path or "").strip()
        if not source:
            raise ValueError(get_prompt("pptx_snapshot_need_out"))
        target = str(Path(source).expanduser().with_suffix(".png"))
    path = render_deck_png(payload, target, resolved)
    logger.info("光栅快照 path=%s pages=%s theme=%s", path, len(payload.slides), resolved.id)
    return get_prompt("pptx_snapshot_ok", path=path, pages=str(len(payload.slides)))


def pptx_replace_slide(
    path: str,
    index: int,
    slide: str,
    theme: str = "",
) -> str:
    from witty_agent.plugins.pptx_kit.lint import lint_slide
    from witty_agent.plugins.pptx_kit.render import load_stored_theme

    try:
        raw = json.loads(slide)
    except json.JSONDecodeError as exc:
        raise ValueError(get_prompt("pptx_bad_json")) from exc
    if not isinstance(raw, dict):
        raise ValueError(get_prompt("pptx_bad_json"))
    spec = parse_slide(raw)
    target = replace_slide(path, spec, index, theme)
    Presentation, *_rest = _pptx()
    stored = load_stored_theme(Presentation(str(Path(target).expanduser())))
    resolved = resolve_theme(theme) if theme.strip() else stored or resolve_theme("custom")
    issues = lint_slide(spec, resolved, index)
    msg = get_prompt("pptx_replace_ok", path=target, index=str(index))
    if issues:
        msg += "\n" + get_prompt("pptx_lint_note", count=str(len(issues)), issues=format_issues(issues)).lstrip()
    return msg


def _skip_shape(name: str) -> bool:
    return name in {"witty-footer", "witty-page", "witty-theme"}


def pptx_list_boxes(path: str, index: int = 0) -> str:
    Presentation, *_rest = _pptx()
    target = Path(path).expanduser()
    deck = Presentation(str(target))
    slides = list(deck.slides)
    if index < 0 or index > len(slides):
        raise ValueError(get_prompt("pptx_bad_index", index=str(index), total=str(len(slides))))
    chosen = list(enumerate(slides, start=1)) if index == 0 else [(index, slides[index - 1])]
    pages: list[str] = []
    for page, slide in chosen:
        lines = [f"## {page}"]
        for shape in slide.shapes:
            if _skip_shape(shape.name or ""):
                continue
            label = shape.name or "(未命名)"
            kind = "chart" if getattr(shape, "has_chart", False) else "table" if shape.has_table else "text" if shape.has_text_frame else "shape"
            preview = ""
            if shape.has_text_frame:
                preview = " ".join(para.text.strip() for para in shape.text_frame.paragraphs if para.text.strip())[:40]
            if preview:
                lines.append(f"- {label}  {kind}  {preview}")
            else:
                lines.append(f"- {label}  {kind}")
        pages.append("\n".join(lines))
    return "\n\n".join(pages) or get_prompt("pptx_empty")


def _find_named(slide, name: str):
    wanted = (name or "").strip()
    for shape in slide.shapes:
        if shape.name == wanted:
            return shape
    folded = wanted.casefold()
    for shape in slide.shapes:
        if (shape.name or "").casefold() == folded:
            return shape
    return None


def pptx_edit_box(
    path: str,
    index: int,
    name: str,
    text: str = "",
    items: str = "",
    fill: str = "",
    color: str = "",
    size: int = 0,
) -> str:
    from witty_agent.plugins.pptx_kit.themes import parse_hex

    Presentation, RGBColor, _shape, _anchor, PP_ALIGN, _emu, Inches, Pt = _pptx()
    target = Path(path).expanduser()
    deck = Presentation(str(target))
    slides = list(deck.slides)
    if index < 1 or index > len(slides):
        raise ValueError(get_prompt("pptx_bad_index", index=str(index), total=str(len(slides))))
    shape = _find_named(slides[index - 1], name)
    if shape is None:
        raise ValueError(get_prompt("pptx_box_missing", name=name, index=str(index)))
    if fill.strip():
        try:
            rgb = parse_hex(fill)
            paint = shape.fill
            paint.solid()
            paint.fore_color.rgb = RGBColor(*rgb)
        except (AttributeError, TypeError, ValueError):
            pass
    if shape.has_text_frame and (text or items.strip()):
        content = text if text else ""
        if items.strip():
            _fill_bullets(shape.text_frame, _bullets_of(items), pt_mod=Pt, rgb_mod=RGBColor, align=PP_ALIGN)
        elif text:
            _write(
                shape.text_frame,
                content,
                pt_mod=Pt,
                rgb_mod=RGBColor,
                align=PP_ALIGN.LEFT,
                size=size or 18,
                color=_INK,
                bold=False,
            )
        if color.strip() or size:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if color.strip():
                        try:
                            run.font.color.rgb = RGBColor(*parse_hex(color))
                        except (AttributeError, TypeError, ValueError):
                            pass
                    if size:
                        run.font.size = Pt(size)
    if index == 1 and name == "witty-title" and text:
        deck.core_properties.title = text
    atomic_save(deck, target)
    return get_prompt("pptx_edit_box_ok", path=str(target), index=str(index), name=name)


def pptx_add_page(path: str, slide: str, index: int = 0) -> str:
    try:
        raw = json.loads(slide)
    except json.JSONDecodeError as exc:
        raise ValueError(get_prompt("pptx_bad_json")) from exc
    if not isinstance(raw, dict):
        raise ValueError(get_prompt("pptx_bad_json"))
    spec = parse_slide(raw)
    if index and index > 0:
        target = insert_slide(path, spec, index)
    else:
        target = append_slide(path, spec)
    return get_prompt("pptx_add_page_ok", path=target, index=str(index or "末"))


def pptx_add_pages(path: str, slides: str, theme: str = "", preview: bool = True) -> str:
    """一批页追加到末尾：长稿分块生成走这条，不要一页一次调 pptx_add_page。"""
    from witty_agent.plugins.pptx_kit.lint import lint_slide
    from witty_agent.plugins.pptx_kit.render import load_stored_theme

    try:
        raw = json.loads(slides)
    except json.JSONDecodeError as exc:
        raise ValueError(get_prompt("pptx_bad_slides")) from exc
    if isinstance(raw, dict):
        raw = raw.get("slides", raw)
    if not isinstance(raw, list) or not raw or not all(isinstance(item, dict) for item in raw):
        raise ValueError(get_prompt("pptx_bad_slides"))
    specs = [parse_slide(item) for item in raw]

    Presentation, *_rest = _pptx()
    source = Path(path).expanduser()
    if not source.is_file():
        raise FileNotFoundError(get_prompt("pptx_deck_missing", path=str(source)))
    before = len(list(Presentation(str(source)).slides))
    target = append_slides(path, specs, theme)

    stored = load_stored_theme(Presentation(str(Path(target).expanduser())))
    resolved = resolve_theme(theme) if theme.strip() else stored or resolve_theme("custom")
    issues: list[Any] = []
    for offset, spec in enumerate(specs):
        issues.extend(lint_slide(spec, resolved, before + offset + 1))

    html_path = ""
    if preview:
        html_path = append_html(str(Path(target).with_suffix(".html")), specs, resolved)
    logger.info(
        "分块追加 pptx=%s html=%s added=%s total=%s issues=%s",
        target,
        html_path or "-",
        len(specs),
        before + len(specs),
        len(issues),
    )
    msg = get_prompt(
        "pptx_add_pages_ok",
        path=target,
        added=str(len(specs)),
        first=str(before + 1),
        last=str(before + len(specs)),
        html=html_path or get_prompt("pptx_no_preview"),
    )
    if issues:
        msg += "\n" + get_prompt("pptx_lint_note", count=str(len(issues)), issues=format_issues(issues)).lstrip()
    return msg


def pptx_from_html(
    path: str,
    html_path: str,
    theme: str = "",
    theme_overrides: str = "",
    preview: bool = False,
) -> str:
    source = Path(html_path).expanduser()
    if not source.is_file():
        raise FileNotFoundError(get_prompt("pptx_html_missing", path=str(source)))
    payload = _apply_theme(parse_html(source.read_text(encoding="utf-8")), theme, theme_overrides)
    return pptx_render(
        path,
        deck=json.dumps(
            {
                "title": payload.title,
                "theme": payload.theme,
                "theme_overrides": payload.theme_overrides,
                "footer": payload.footer,
                "slides": [_slide_dict(item) for item in payload.slides],
            },
            ensure_ascii=False,
        ),
        preview=preview,
    )


def _slide_dict(item: Slide) -> dict[str, Any]:
    payload: dict[str, Any] = {"kind": item.kind, "title": item.title}
    if item.bg:
        payload["bg"] = item.bg
    if item.boxes:
        payload["boxes"] = [_box_dict(box) for box in item.boxes]
    if item.subtitle:
        payload["subtitle"] = item.subtitle
    if item.kicker:
        payload["kicker"] = item.kicker
    if item.meta:
        payload["meta"] = item.meta
    if item.items:
        payload["items"] = item.items
    if item.left or item.left_title:
        payload["left_title"] = item.left_title
        payload["left"] = item.left
    if item.right or item.right_title:
        payload["right_title"] = item.right_title
        payload["right"] = item.right
    if item.metrics:
        payload["metrics"] = [{"label": m.label, "value": m.value, "note": m.note} for m in item.metrics]
    if item.cards:
        payload["cards"] = [{"title": c.title, "body": c.body} for c in item.cards]
    if item.headers:
        payload["headers"] = item.headers
    if item.rows:
        payload["rows"] = item.rows
    if item.steps:
        payload["steps"] = [{"title": s.title, "body": s.body} for s in item.steps]
    if item.quote:
        payload["quote"] = item.quote
    if item.by:
        payload["by"] = item.by
    if item.image:
        payload["image"] = item.image
    if item.notes:
        payload["notes"] = item.notes
    return payload


def _box_dict(box: Any) -> dict[str, Any]:
    payload: dict[str, Any] = {
        "kind": box.kind,
        "x": box.x,
        "y": box.y,
        "w": box.w,
        "h": box.h,
    }
    if box.text:
        payload["text"] = box.text
    if box.items:
        payload["items"] = box.items
    if box.fill:
        payload["fill"] = box.fill
    if box.color:
        payload["color"] = box.color
    if box.size:
        payload["size"] = box.size
    if box.bold:
        payload["bold"] = True
    if box.italic:
        payload["italic"] = True
    if box.align and box.align != "left":
        payload["align"] = box.align
    if box.name:
        payload["name"] = box.name
    if box.font:
        payload["font"] = box.font
    if box.headers:
        payload["headers"] = box.headers
    if box.rows:
        payload["rows"] = box.rows
    if box.image:
        payload["image"] = box.image
    if box.kind == "chart" or box.chart:
        payload["chart"] = box.chart or "column"
    if box.categories:
        payload["categories"] = box.categories
    if box.series:
        payload["series"] = [{"name": item.name, "values": item.values} for item in box.series]
    if box.colors:
        payload["colors"] = box.colors
    if box.shadow is not None:
        payload["shadow"] = box.shadow
    if box.radius:
        payload["radius"] = box.radius
    return payload


def _spec(name: str, func: Any, properties: dict[str, Any], required: list[str] | None = None) -> None:
    parameters: dict[str, Any] = {
        "type": "object",
        "properties": properties,
        "additionalProperties": False,
    }
    if required:
        parameters["required"] = required
    register_tool(
        ToolSpec(
            name=name,
            description=get_prompt(f"tool_desc_{name}"),
            parameters=parameters,
            func=func,
        )
    )


_spec(
    "pptx_create",
    pptx_create,
    {
        "path": {"type": "string", "description": get_prompt("pptx_param_path")},
        "title": {"type": "string", "description": get_prompt("pptx_param_title")},
        "subtitle": {"type": "string", "description": get_prompt("pptx_param_subtitle")},
    },
    ["path", "title"],
)
_spec(
    "pptx_add_slide",
    pptx_add_slide,
    {
        "path": {"type": "string", "description": get_prompt("pptx_param_path")},
        "title": {"type": "string", "description": get_prompt("pptx_param_title")},
        "bullets": {"type": "string", "description": get_prompt("pptx_param_bullets")},
    },
    ["path", "title", "bullets"],
)
_spec(
    "pptx_edit_slide",
    pptx_edit_slide,
    {
        "path": {"type": "string", "description": get_prompt("pptx_param_path")},
        "index": {"type": "integer", "description": get_prompt("pptx_param_index")},
        "title": {"type": "string", "description": get_prompt("pptx_param_title")},
        "bullets": {"type": "string", "description": get_prompt("pptx_param_bullets")},
    },
    ["path", "index"],
)
_spec(
    "pptx_add_picture",
    pptx_add_picture,
    {
        "path": {"type": "string", "description": get_prompt("pptx_param_path")},
        "image_path": {"type": "string", "description": get_prompt("pptx_param_image")},
        "title": {"type": "string", "description": get_prompt("pptx_param_title")},
    },
    ["path", "image_path"],
)
_spec(
    "pptx_outline",
    pptx_outline,
    {"path": {"type": "string", "description": get_prompt("pptx_param_path")}},
    ["path"],
)
_spec("pptx_themes", pptx_themes, {})
_spec(
    "pptx_render",
    pptx_render,
    {
        "path": {"type": "string", "description": get_prompt("pptx_param_path")},
        "deck": {"type": "string", "description": get_prompt("pptx_param_deck")},
        "deck_path": {"type": "string", "description": get_prompt("pptx_param_deck_path")},
        "theme": {"type": "string", "description": get_prompt("pptx_param_theme")},
        "theme_overrides": {"type": "string", "description": get_prompt("pptx_param_overrides")},
        "preview": {"type": "boolean", "description": get_prompt("pptx_param_preview")},
    },
    ["path"],
)
_spec(
    "pptx_from_html",
    pptx_from_html,
    {
        "path": {"type": "string", "description": get_prompt("pptx_param_path")},
        "html_path": {"type": "string", "description": get_prompt("pptx_param_html")},
        "theme": {"type": "string", "description": get_prompt("pptx_param_theme")},
        "theme_overrides": {"type": "string", "description": get_prompt("pptx_param_overrides")},
        "preview": {"type": "boolean", "description": get_prompt("pptx_param_preview")},
    },
    ["path", "html_path"],
)
_spec(
    "pptx_check",
    pptx_check,
    {
        "deck": {"type": "string", "description": get_prompt("pptx_param_deck")},
        "deck_path": {"type": "string", "description": get_prompt("pptx_param_deck_path")},
        "theme": {"type": "string", "description": get_prompt("pptx_param_theme")},
        "theme_overrides": {"type": "string", "description": get_prompt("pptx_param_overrides")},
    },
)
_spec(
    "pptx_snapshot",
    pptx_snapshot,
    {
        "deck": {"type": "string", "description": get_prompt("pptx_param_deck")},
        "deck_path": {"type": "string", "description": get_prompt("pptx_param_deck_path")},
        "out": {"type": "string", "description": get_prompt("pptx_param_out")},
        "theme": {"type": "string", "description": get_prompt("pptx_param_theme")},
        "theme_overrides": {"type": "string", "description": get_prompt("pptx_param_overrides")},
    },
)
_spec(
    "pptx_replace_slide",
    pptx_replace_slide,
    {
        "path": {"type": "string", "description": get_prompt("pptx_param_path")},
        "index": {"type": "integer", "description": get_prompt("pptx_param_index")},
        "slide": {"type": "string", "description": get_prompt("pptx_param_slide")},
        "theme": {"type": "string", "description": get_prompt("pptx_param_theme")},
    },
    ["path", "index", "slide"],
)
_spec(
    "pptx_list_boxes",
    pptx_list_boxes,
    {
        "path": {"type": "string", "description": get_prompt("pptx_param_path")},
        "index": {"type": "integer", "description": get_prompt("pptx_param_index_or_all")},
    },
    ["path"],
)
_spec(
    "pptx_edit_box",
    pptx_edit_box,
    {
        "path": {"type": "string", "description": get_prompt("pptx_param_path")},
        "index": {"type": "integer", "description": get_prompt("pptx_param_index")},
        "name": {"type": "string", "description": get_prompt("pptx_param_box_name")},
        "text": {"type": "string", "description": get_prompt("pptx_param_box_text")},
        "items": {"type": "string", "description": get_prompt("pptx_param_bullets")},
        "fill": {"type": "string", "description": get_prompt("pptx_param_box_fill")},
        "color": {"type": "string", "description": get_prompt("pptx_param_box_color")},
        "size": {"type": "integer", "description": get_prompt("pptx_param_box_size")},
    },
    ["path", "index", "name"],
)
_spec(
    "pptx_add_page",
    pptx_add_page,
    {
        "path": {"type": "string", "description": get_prompt("pptx_param_path")},
        "slide": {"type": "string", "description": get_prompt("pptx_param_slide")},
        "index": {"type": "integer", "description": get_prompt("pptx_param_insert_index")},
    },
    ["path", "slide"],
)
_spec(
    "pptx_add_pages",
    pptx_add_pages,
    {
        "path": {"type": "string", "description": get_prompt("pptx_param_path")},
        "slides": {"type": "string", "description": get_prompt("pptx_param_slides")},
        "theme": {"type": "string", "description": get_prompt("pptx_param_theme")},
        "preview": {"type": "boolean", "description": get_prompt("pptx_param_preview")},
    },
    ["path", "slides"],
)
