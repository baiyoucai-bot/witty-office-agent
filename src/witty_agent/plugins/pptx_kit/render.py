"""把语义稿写成原生 PPTX 形状。文字/表格可在 WPS、Office 里改。"""

from __future__ import annotations

import json
from dataclasses import replace
from pathlib import Path
from typing import Any

from witty_agent.logging import get_logger
from witty_agent.plugins.pptx_kit.chrome import (
    BAND_BOTTOM,
    BAND_H,
    COVER_BAND_H,
    COVER_BAND_Y,
    COVER_DIVIDER_X,
    COVER_GOLD_H,
    COVER_KICKER_PT,
    COVER_META_PT,
    COVER_SUB_PT,
    COVER_TEXT_X,
    COVER_TITLE_PT,
    LOGO_NAME,
    RULE_H,
    RULE_LEFT,
    SECTION_BAND_H,
    SECTION_BAND_Y,
    SECTION_BLOCK_TINT,
    SECTION_BLOCK_W,
    SECTION_SUB_PT,
    SECTION_TEXT_GAP,
    SECTION_TITLE_PT,
    cover_logo_rect,
    emblem_asset,
    emblem_rect,
    is_dark,
    logo_asset,
    logo_ground,
    logo_rect,
    section_mark_pt,
    shade,
    wants_band,
    wants_logo,
)
from witty_agent.plugins.pptx_kit.metrics import (
    BULLET_INDENT,
    bullet_gap,
    line_factor,
    table_row_height,
    text_height,
)
from witty_agent.plugins.pptx_kit.schema import CANVAS_H, CANVAS_W, SHAPE_KINDS, Box, Slide
from witty_agent.plugins.pptx_kit.shapes import normalize_point
from witty_agent.plugins.pptx_kit.themes import (
    Theme,
    color_of,
    contrast_ratio,
    resolve_installed_font,
    resolve_theme,
    theme_from_payload,
    theme_payload,
)
from witty_agent.prompts import get_prompt

_THEME_MARK = "WITTY_THEME="

logger = get_logger("plugins.pptx")

MARGIN_X = 0.62
MARGIN_TOP = 0.42
CONTENT_MAX_Y = 6.70
FOOTER_TOP = 6.88
EMU_PER_INCH = 914400
_COVER_KINDS = frozenset({"cover", "section"})


def _deps():
    try:
        from lxml import etree
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.oxml.ns import qn
        from pptx.util import Emu, Inches, Pt
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError(get_prompt("pptx_missing")) from exc
    return Presentation, RGBColor, MSO_SHAPE, MSO_CONNECTOR, MSO_ANCHOR, PP_ALIGN, Emu, Inches, Pt, qn, etree


# 矢量盒子 → PowerPoint 内置形状。带朝向的用 point 选具体那一个。
_ARROW_BY_POINT = {
    "right": "RIGHT_ARROW",
    "left": "LEFT_ARROW",
    "up": "UP_ARROW",
    "down": "DOWN_ARROW",
}
_MSO_BY_KIND = {
    "oval": "OVAL",
    "diamond": "DIAMOND",
    "triangle": "ISOSCELES_TRIANGLE",
    "chevron": "CHEVRON",
    "pentagon": "PENTAGON",
    "rect": "RECTANGLE",
    "round": "ROUNDED_RECTANGLE",
}


def _mso_shape(shapes, kind: str, point: str):
    """形状枚举。chevron / pentagon 在 PowerPoint 里只朝右，其它朝向靠旋转，
    这里先只支持朝右——朝左的流程带在幻灯片里几乎不用，硬转会把里面的字也转过去。"""
    if kind == "arrow":
        return getattr(shapes, _ARROW_BY_POINT[normalize_point(point)])
    return getattr(shapes, _MSO_BY_KIND.get(kind, "RECTANGLE"))


def _blank(deck):
    layouts = deck.slide_layouts
    return layouts[6] if len(layouts) > 6 else layouts[0]


def _set_wide(deck) -> None:
    # 自己取 Inches，不从 _mods() 按下标拿——加一个依赖就会把下标全打乱。
    try:
        from pptx.util import Inches
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError(get_prompt("pptx_missing")) from exc
    deck.slide_width = Inches(CANVAS_W)
    deck.slide_height = Inches(CANVAS_H)


def _latin_only(text: str) -> bool:
    return bool(text) and not any("\u4e00" <= ch <= "\u9fff" for ch in text)


def _apply_font(
    run,
    font: str,
    size: int,
    color,
    *,
    bold: bool,
    italic: bool,
    rgb_mod,
    pt_mod,
    qn,
    etree,
    tracking: int = 0,
) -> None:
    run.font.size = pt_mod(size)
    run.font.bold = bold
    run.font.italic = italic and _latin_only(run.text or "")
    run.font.color.rgb = rgb_mod(*color)
    run.font.name = font
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = r_pr.find(qn(tag))
        if el is None:
            el = etree.SubElement(r_pr, qn(tag))
        el.set("typeface", font)
    if tracking:
        r_pr.set("spc", str(tracking))


def _apply_shadow(shape, qn, etree) -> None:
    """极浅的一层影。重影在投影仪上是脏的。"""
    sp_pr = shape._element.spPr
    effect = sp_pr.find(qn("a:effectLst"))
    if effect is None:
        effect = etree.SubElement(sp_pr, qn("a:effectLst"))
    for child in list(effect):
        if child.tag == qn("a:outerShdw"):
            effect.remove(child)
    shdw = etree.SubElement(effect, qn("a:outerShdw"))
    shdw.set("blurRad", "34000")
    shdw.set("dist", "17000")
    shdw.set("dir", "5400000")
    shdw.set("algn", "t")
    shdw.set("rotWithShape", "0")
    color = etree.SubElement(shdw, qn("a:srgbClr"))
    color.set("val", "1A2428")
    alpha = etree.SubElement(color, qn("a:alpha"))
    alpha.set("val", "11000")


def _fit_text(body_pr, qn, etree) -> None:
    """字多了缩字号，不要让形状自己长胖把版面顶乱。"""
    for tag in ("a:spAutoFit", "a:noAutofit", "a:normAutofit"):
        found = body_pr.find(qn(tag))
        if found is not None:
            body_pr.remove(found)
    etree.SubElement(body_pr, qn("a:normAutofit"))


_BULLET_GLYPHS = ("·", "•", "▪", "◦", "－", "*")


def _bullet_text(line: str) -> str:
    text = str(line).strip()
    if text[:1] in _BULLET_GLYPHS:
        text = text[1:].lstrip()
    elif text[:2] == "- ":
        text = text[2:].lstrip()
    return text


def _apply_bullet(para, qn, etree, *, marl: int, color=None) -> None:
    """真项目符号 + 悬挂缩进。折行会对齐到文字，不会顶到符号底下。"""
    p_pr = para._p.get_or_add_pPr()
    p_pr.set("marL", str(marl))
    p_pr.set("indent", str(-marl))
    order = ("a:buClr", "a:buFont", "a:buChar")
    for tag in (*order, "a:buNone", "a:buAutoNum"):
        found = p_pr.find(qn(tag))
        if found is not None:
            p_pr.remove(found)
    if color is not None:
        bu_clr = etree.SubElement(p_pr, qn("a:buClr"))
        srgb = etree.SubElement(bu_clr, qn("a:srgbClr"))
        srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
    bu_font = etree.SubElement(p_pr, qn("a:buFont"))
    bu_font.set("typeface", "Arial")
    bu_char = etree.SubElement(p_pr, qn("a:buChar"))
    bu_char.set("char", "•")


_LN_SIDES = ("a:lnL", "a:lnR", "a:lnT", "a:lnB")


def _cell_border(cell, qn, etree, color, *, side: str = "a:lnB", width: int = 9525) -> None:
    """给单元格描一条细线。a:ln* 必须排在填充之前，否则 WPS 直接判文件损坏。"""
    tc_pr = cell._tc.get_or_add_tcPr()
    found = tc_pr.find(qn(side))
    if found is not None:
        tc_pr.remove(found)
    ln = etree.Element(qn(side))
    ln.set("w", str(width))
    ln.set("cap", "flat")
    ln.set("cmpd", "sng")
    ln.set("algn", "ctr")
    fill = etree.SubElement(ln, qn("a:solidFill"))
    srgb = etree.SubElement(fill, qn("a:srgbClr"))
    srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
    tags = {child.tag for child in tc_pr}
    pos = sum(1 for earlier in _LN_SIDES[: _LN_SIDES.index(side)] if qn(earlier) in tags)
    tc_pr.insert(pos, ln)


def _align(pp_align, name: str):
    if name == "center":
        return pp_align.CENTER
    if name == "right":
        return pp_align.RIGHT
    return pp_align.LEFT


class Painter:
    def __init__(self, slide, theme: Theme, mods) -> None:
        (
            self.rgb,
            self.shapes,
            self.connectors,
            self.anchor,
            self.pp_align,
            self.inches,
            self.pt,
            self.qn,
            self.etree,
        ) = mods
        self.slide = slide
        self.theme = theme
        self.bg = theme.bg

    def paint_bg(self, color) -> None:
        fill = self.slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.rgb(*color)
        self.bg = (color[0], color[1], color[2])

    def shape(
        self,
        kind,
        left,
        top,
        width,
        height,
        color,
        *,
        name: str = "",
        rounded: bool = False,
        radius: float = 0.12,
        shadow: bool = False,
        stroke=None,
        stroke_w: float = 0.0,
    ):
        enum = self.shapes.ROUNDED_RECTANGLE if rounded else kind
        box = self.slide.shapes.add_shape(
            enum,
            self.inches(left),
            self.inches(top),
            self.inches(width),
            self.inches(height),
        )
        if stroke is not None and stroke_w > 0:
            box.line.color.rgb = self.rgb(*stroke)
            box.line.width = self.pt(stroke_w)
        else:
            box.line.fill.background()
        fill = box.fill
        if color is None:
            fill.background()
        else:
            fill.solid()
            fill.fore_color.rgb = self.rgb(*color)
        if name:
            box.name = name
        if rounded:
            try:
                box.adjustments[0] = radius if radius > 0 else 0.12
            except (AttributeError, IndexError, ValueError):
                pass
        if shadow:
            _apply_shadow(box, self.qn, self.etree)
        return box

    def vector(
        self,
        kind: str,
        left,
        top,
        width,
        height,
        color,
        *,
        name: str = "",
        point: str = "right",
        shadow: bool = False,
        stroke=None,
        stroke_w: float = 0.0,
    ):
        """画一个原生矢量形状。用 MSO 内置形状而不是自由曲线，WPS 里才拖得动、改得了色。"""
        enum = _mso_shape(self.shapes, kind, point)
        return self.shape(
            enum,
            left,
            top,
            width,
            height,
            color,
            name=name,
            shadow=shadow,
            stroke=stroke,
            stroke_w=stroke_w,
        )

    def connector(self, left, top, width, height, color, *, name: str = "", point: str = "right", weight: float = 1.5):
        """带箭头的细连接线。架构图里连两个框用它，比粗箭头干净。"""
        facing = normalize_point(point)
        if facing in {"left", "right"}:
            mid = top + height / 2
            x1, y1, x2, y2 = left, mid, left + width, mid
        else:
            mid = left + width / 2
            x1, y1, x2, y2 = mid, top, mid, top + height
        if facing in {"left", "up"}:
            x1, y1, x2, y2 = x2, y2, x1, y1
        line = self.slide.shapes.add_connector(
            self.connectors.STRAIGHT,
            self.inches(x1),
            self.inches(y1),
            self.inches(x2),
            self.inches(y2),
        )
        if name:
            line.name = name
        line.line.color.rgb = self.rgb(*color)
        line.line.width = self.pt(weight)
        tail = self.etree.SubElement(line.line._get_or_add_ln(), self.qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
        return line

    def rect(self, left, top, width, height, color, *, name: str = "", shadow: bool = False):
        return self.shape(self.shapes.RECTANGLE, left, top, width, height, color, name=name, shadow=shadow)

    def card(self, left, top, width, height, color, *, name: str = "", radius: float = 0.12, shadow: bool = True):
        return self.shape(
            self.shapes.ROUNDED_RECTANGLE,
            left,
            top,
            width,
            height,
            color,
            name=name,
            rounded=True,
            radius=radius,
            shadow=shadow,
        )

    def text(
        self,
        left,
        top,
        width,
        height,
        content: str,
        *,
        size: int,
        color,
        bold: bool = False,
        italic: bool = False,
        align: str = "left",
        name: str = "",
        anchor: str = "top",
        font: str | None = None,
    ):
        box = self.slide.shapes.add_textbox(
            self.inches(left),
            self.inches(top),
            self.inches(width),
            self.inches(height),
        )
        if name:
            box.name = name
        self.write_text(
            box.text_frame,
            content,
            size=size,
            color=color,
            bold=bold,
            italic=italic,
            align=align,
            anchor=anchor,
            font=font,
            name=name,
        )
        return box

    def write_text(
        self,
        frame,
        content: str,
        *,
        size: int,
        color,
        bold: bool = False,
        italic: bool = False,
        align: str = "left",
        anchor: str = "top",
        font: str | None = None,
        name: str = "",
        pad: float = 0.06,
    ) -> None:
        """往一个 text_frame 里排字。文本框和形状内文字共用，免得两处样式各走各的。"""
        frame.word_wrap = True
        frame.clear()
        frame.margin_left = self.inches(pad)
        frame.margin_right = self.inches(pad)
        frame.margin_top = self.inches(0.04)
        frame.margin_bottom = self.inches(0.04)
        body_pr = frame._element.bodyPr
        body_pr.set("anchor", {"top": "t", "middle": "ctr", "bottom": "b"}.get(anchor, "t"))
        _fit_text(body_pr, self.qn, self.etree)
        tracking = 180 if (size <= 13 and bold) or (name or "").endswith("kicker") else 0
        chunks = str(content if content else " ").replace("\r\n", "\n").split("\n")
        for index, chunk in enumerate(chunks):
            para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            para.alignment = _align(self.pp_align, align)
            try:
                para.line_spacing = line_factor(size)
            except (AttributeError, TypeError, ValueError):
                pass
            if index:
                para.space_before = self.pt(max(4.0, size * 0.30))
            run = para.add_run()
            run.text = chunk if chunk.strip() else " "
            _apply_font(
                run,
                font or self.theme.font,
                size,
                color,
                bold=bold,
                italic=italic,
                rgb_mod=self.rgb,
                pt_mod=self.pt,
                qn=self.qn,
                etree=self.etree,
                tracking=tracking,
            )

    def bullets(
        self,
        left,
        top,
        width,
        height,
        lines: list[str],
        *,
        size: int = 18,
        color=None,
        name: str = "witty-body",
        bullet_color=None,
    ):
        box = self.slide.shapes.add_textbox(
            self.inches(left),
            self.inches(top),
            self.inches(width),
            self.inches(height),
        )
        box.name = name
        frame = box.text_frame
        frame.word_wrap = True
        frame.clear()
        frame.margin_left = self.inches(0.06)
        frame.margin_right = self.inches(0.06)
        frame.margin_top = self.inches(0.03)
        frame.margin_bottom = self.inches(0.03)
        _fit_text(frame._element.bodyPr, self.qn, self.etree)
        ink = color or self.theme.ink
        clean = [_bullet_text(item) for item in lines if str(item).strip()]
        clean = [item for item in clean if item]
        if not clean:
            return box
        marl = int(BULLET_INDENT * EMU_PER_INCH)
        gap = bullet_gap(size)
        for index, line in enumerate(clean):
            para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            para.alignment = self.pp_align.LEFT
            para.level = 0
            try:
                para.line_spacing = line_factor(size)
            except (AttributeError, TypeError, ValueError):
                pass
            if index:
                para.space_before = self.pt(gap)
            run = para.add_run()
            run.text = line
            _apply_font(
                run,
                self.theme.font,
                size,
                ink,
                bold=False,
                italic=False,
                rgb_mod=self.rgb,
                pt_mod=self.pt,
                qn=self.qn,
                etree=self.etree,
            )
            _apply_bullet(para, self.qn, self.etree, marl=marl, color=bullet_color)
        return box

    def footer(self, title: str, page: int) -> None:
        label = f"{title}  ·  {page}" if title else str(page)
        tint = _muted_on(self.theme, self.bg)
        self.text(
            MARGIN_X,
            FOOTER_TOP,
            10.6,
            0.32,
            label,
            size=10,
            color=tint,
            name="witty-footer",
        )
        self.text(
            11.4,
            FOOTER_TOP,
            1.3,
            0.32,
            f"{page:02d}",
            size=10,
            color=tint,
            align="right",
            name="witty-page",
        )

    def band_chrome(self) -> None:
        """企业模板的顶部白带 + 压线。宏页和 boxes 页共用，模板才不会缺件。"""
        self.rect(0, 0, CANVAS_W, BAND_H, self.theme.surface, name="witty-band")
        self.rect(0, BAND_H, RULE_LEFT, RULE_H, self.theme.accent2, name="witty-band-rule")
        self.rect(
            RULE_LEFT,
            BAND_H,
            CANVAS_W - RULE_LEFT,
            RULE_H,
            self.theme.line,
            name="witty-band-line",
        )

    def chrome(self, *, dark: bool = False) -> None:
        """内页顶栏。band 是企业模板的白带压线，bar 是通用的左侧竖条。"""
        style = self.theme.chrome
        if style == "band":
            if not dark:
                self.band_chrome()
        elif style == "bar" and not dark:
            self.rect(0, 0, CANVAS_W, 0.08, self.theme.accent2, name="witty-top")
            self.rect(0, 0, 0.12, CANVAS_H, self.theme.accent2, name="witty-bar")
        self.logo(dark=dark)

    def logo(self, *, dark: bool = False, cover: bool = False) -> None:
        """右上角摆标识。深底自动换反白版，封面放大一号，缺图就静默跳过。"""
        path = logo_asset(self.theme, dark=dark)
        if not path:
            return
        x, y, w, _h = cover_logo_rect(self.theme) if cover else logo_rect(self.theme)
        self.picture(path, x, y, w, name=LOGO_NAME)

    def picture(self, path, left, top, width, *, name: str = ""):
        """只给宽，高按图自身比例走，永远不拉变形。"""
        pic = self.slide.shapes.add_picture(
            str(path),
            self.inches(left),
            self.inches(top),
            width=self.inches(width),
        )
        if name:
            pic.name = name
        return pic

    def notes(self, text: str) -> None:
        if not text.strip():
            return
        notes = self.slide.notes_slide
        notes.notes_text_frame.text = text.strip()


def _mods():
    _P, RGBColor, MSO_SHAPE, MSO_CONNECTOR, MSO_ANCHOR, PP_ALIGN, _emu, Inches, Pt, qn, etree = _deps()
    return RGBColor, MSO_SHAPE, MSO_CONNECTOR, MSO_ANCHOR, PP_ALIGN, Inches, Pt, qn, etree


def _on(theme: Theme, bg) -> tuple[int, int, int]:
    # 深底用白字、浅底用正文色，避免金条/红条上的字看不清。
    return (255, 255, 255) if is_dark(bg) else theme.ink


def _muted_on(theme: Theme, bg) -> tuple[int, int, int]:
    # 次要字（页脚、注解）跟着底色换：深底用封面次色，浅底用常规灰。
    return theme.cover_muted if is_dark(bg) else theme.muted


def _band_header(p: Painter, spec: Slide) -> float:
    """白带下面起标题。压线已经交代了层级，这里不再画短横线。"""
    y = BAND_BOTTOM + 0.26
    if spec.kicker:
        p.text(
            MARGIN_X,
            y,
            12.0,
            0.30,
            spec.kicker,
            size=12,
            color=p.theme.accent,
            bold=True,
            name="witty-kicker",
        )
        y += 0.34
    p.text(
        MARGIN_X,
        y,
        12.0,
        0.56,
        spec.title or " ",
        size=26,
        color=p.theme.accent2,
        bold=True,
        name="witty-title",
    )
    return y + 0.64


def _header_block(p: Painter, spec: Slide, *, top: float = MARGIN_TOP) -> float:
    if p.theme.chrome == "band":
        return _band_header(p, spec)
    y = top
    if spec.kicker:
        p.text(
            MARGIN_X,
            y,
            12.0,
            0.32,
            spec.kicker,
            size=12,
            color=p.theme.accent,
            bold=True,
            name="witty-kicker",
        )
        y += 0.34
    title = spec.title or " "
    p.text(
        MARGIN_X,
        y,
        12.0,
        0.62,
        title,
        size=26,
        color=p.theme.ink,
        bold=True,
        name="witty-title",
    )
    y += 0.68
    p.rect(MARGIN_X, y, 1.35, 0.045, p.theme.accent, name="witty-rule")
    return y + 0.18


def _cover_grid(p: Painter, spec: Slide, page: int, footer: str) -> None:
    """grid 封面：浅底 + 通栏青绿带，主题带徽标时带内左徽标、右标题，带下压一条金线。

    对照原稿：封面顶部不压白带；主题带标识时只在右上角放大一号的标识。
    """
    theme = p.theme
    p.paint_bg(theme.bg)
    p.logo(cover=True)
    band_y, band_h = COVER_BAND_Y, COVER_BAND_H
    p.rect(0, band_y, CANVAS_W, band_h, theme.cover_bg, name="witty-band-main")
    p.rect(0, band_y + band_h, CANVAS_W, COVER_GOLD_H, theme.bar, name="witty-band-gold")
    text_x = MARGIN_X
    emblem = emblem_asset(theme, dark=True)
    if emblem:
        ex, ey, ew, _eh = emblem_rect(theme)
        p.picture(emblem, ex, ey, ew, name="witty-emblem")
        p.rect(COVER_DIVIDER_X, band_y + 0.66, 0.02, band_h - 1.32, theme.cover_muted, name="witty-cover-line")
        text_x = COVER_TEXT_X
    width = CANVAS_W - text_x - MARGIN_X
    # 标题占几行是现算的：一行的封面不能留出两行的空档，两行的也不能压到副标题上。
    title = spec.title or " "
    title_h = max(text_height(title, width, COVER_TITLE_PT, theme.font), 0.62)
    subtitle_h = text_height(spec.subtitle, width, COVER_SUB_PT, theme.font) if spec.subtitle else 0.0
    block_h = title_h + (subtitle_h + 0.30 if subtitle_h else 0.0)
    title_y = band_y + (band_h - block_h) / 2 + (0.16 if spec.kicker else 0.0)
    if spec.kicker:
        p.text(
            text_x,
            title_y - 0.46,
            width,
            0.36,
            spec.kicker,
            size=COVER_KICKER_PT,
            color=theme.cover_muted,
            name="witty-kicker",
        )
    p.text(
        text_x, title_y, width, title_h, title, size=COVER_TITLE_PT, color=theme.cover_ink, bold=True, name="witty-title"
    )
    if spec.subtitle:
        p.text(
            text_x,
            title_y + title_h + 0.30,
            width,
            subtitle_h,
            spec.subtitle,
            size=COVER_SUB_PT,
            color=theme.cover_muted,
            name="witty-body",
        )
    if spec.meta:
        p.text(MARGIN_X, 6.02, 12.0, 0.40, spec.meta, size=COVER_META_PT, color=theme.muted, name="witty-meta")
    p.footer(footer, page)


def _cover(p: Painter, spec: Slide, page: int, footer: str) -> None:
    theme = p.theme
    style = theme.cover
    if style == "grid":
        _cover_grid(p, spec, page, footer)
        return
    p.paint_bg(theme.cover_bg if style in {"bar", "split"} else theme.bg)
    if style == "bar":
        p.paint_bg(theme.cover_bg)
        p.rect(0, 0, 0.16, CANVAS_H, theme.bar, name="witty-bar")
        p.text(MARGIN_X, 1.55, 11.8, 0.36, spec.kicker or " ", size=14, color=theme.cover_muted, name="witty-kicker")
        p.rect(MARGIN_X, 2.05, 2.15, 0.055, theme.bar, name="witty-rule")
        p.text(MARGIN_X, 2.25, 11.8, 1.35, spec.title or " ", size=40, color=theme.cover_ink, bold=True, name="witty-title")
        if spec.subtitle:
            p.text(MARGIN_X, 3.75, 11.8, 1.1, spec.subtitle, size=18, color=theme.cover_muted, name="witty-body")
        if spec.meta:
            p.text(MARGIN_X, 6.15, 11.8, 0.4, spec.meta, size=13, color=theme.cover_muted, name="witty-meta")
    elif style == "split":
        p.paint_bg(theme.bg)
        p.rect(0, 0, 5.15, CANVAS_H, theme.cover_bg, name="witty-bar")
        p.text(0.55, 2.15, 4.3, 0.4, spec.kicker or " ", size=13, color=theme.cover_muted, name="witty-kicker")
        p.text(0.55, 2.6, 4.3, 2.2, spec.title or " ", size=32, color=theme.cover_ink, bold=True, name="witty-title")
        if spec.subtitle:
            p.text(5.7, 2.7, 6.8, 2.4, spec.subtitle, size=18, color=theme.ink, name="witty-body")
        if spec.meta:
            p.text(5.7, 6.15, 6.8, 0.4, spec.meta, size=13, color=theme.muted, name="witty-meta")
    elif style == "band":
        p.paint_bg(theme.bg)
        p.rect(0, 0, CANVAS_W, 1.42, theme.accent, name="witty-bar")
        p.text(MARGIN_X, 0.48, 12.0, 0.55, spec.kicker or footer, size=16, color=_on(theme, theme.accent), name="witty-kicker")
        p.text(MARGIN_X, 2.15, 12.0, 1.4, spec.title or " ", size=40, color=theme.ink, bold=True, name="witty-title")
        if spec.subtitle:
            p.text(MARGIN_X, 3.7, 12.0, 1.1, spec.subtitle, size=18, color=theme.muted, name="witty-body")
        if spec.meta:
            p.text(MARGIN_X, 6.15, 12.0, 0.4, spec.meta, size=13, color=theme.muted, name="witty-meta")
    elif style == "mark":
        p.paint_bg(theme.bg)
        p.rect(MARGIN_X, 2.05, 0.55, 0.55, theme.accent, name="witty-bar")
        p.text(MARGIN_X, 1.35, 12.0, 0.4, spec.kicker or " ", size=13, color=theme.accent, name="witty-kicker")
        p.text(MARGIN_X, 2.75, 12.0, 1.5, spec.title or " ", size=42, color=theme.ink, bold=True, name="witty-title")
        if spec.subtitle:
            p.text(MARGIN_X, 4.4, 12.0, 1.1, spec.subtitle, size=18, color=theme.muted, name="witty-body")
        if spec.meta:
            p.text(MARGIN_X, 6.15, 12.0, 0.4, spec.meta, size=13, color=theme.muted, name="witty-meta")
    else:
        p.paint_bg(theme.bg)
        p.text(MARGIN_X, 1.7, 12.0, 0.4, spec.kicker or " ", size=13, color=theme.accent, name="witty-kicker")
        p.text(MARGIN_X, 2.25, 12.0, 1.5, spec.title or " ", size=40, color=theme.ink, bold=True, name="witty-title")
        p.rect(MARGIN_X, 3.9, 1.6, 0.05, theme.accent, name="witty-rule")
        if spec.subtitle:
            p.text(MARGIN_X, 4.15, 12.0, 1.1, spec.subtitle, size=18, color=theme.muted, name="witty-body")
        if spec.meta:
            p.text(MARGIN_X, 6.15, 12.0, 0.4, spec.meta, size=13, color=theme.muted, name="witty-meta")
    p.footer(footer, page)


def _section_band(p: Painter, spec: Slide, page: int, footer: str) -> None:
    """grid 章节页：通栏青绿带，左边切一块浅一档的同色放大号数，右边白标题。"""
    theme = p.theme
    p.paint_bg(theme.bg)
    p.logo()
    band_y, band_h, block_w = SECTION_BAND_Y, SECTION_BAND_H, SECTION_BLOCK_W
    p.rect(0, band_y, CANVAS_W, band_h, theme.cover_bg, name="witty-band-main")
    p.rect(0, band_y, block_w, band_h, shade(theme.cover_bg, SECTION_BLOCK_TINT), name="witty-band-no")
    mark = (spec.kicker or f"{page:02d}").strip()
    p.text(
        0,
        band_y,
        block_w,
        band_h,
        mark,
        size=section_mark_pt(mark),
        color=theme.cover_ink,
        bold=True,
        align="center",
        anchor="middle",
        name="witty-kicker",
    )
    p.text(
        block_w + SECTION_TEXT_GAP,
        band_y,
        CANVAS_W - block_w - SECTION_TEXT_GAP - MARGIN_X,
        band_h,
        spec.title or " ",
        size=SECTION_TITLE_PT,
        color=theme.cover_ink,
        bold=True,
        anchor="middle",
        name="witty-title",
    )
    if spec.subtitle:
        p.text(
            MARGIN_X,
            band_y + band_h + 0.44,
            12.0,
            0.80,
            spec.subtitle,
            size=SECTION_SUB_PT,
            color=theme.muted,
            name="witty-body",
        )
    p.footer(footer, page)


def _section(p: Painter, spec: Slide, page: int, footer: str) -> None:
    if p.theme.chrome == "band":
        _section_band(p, spec, page, footer)
        return
    p.paint_bg(p.theme.cover_bg)
    p.rect(0, 0, 0.16, CANVAS_H, p.theme.bar, name="witty-bar")
    p.text(MARGIN_X, 2.35, 12.0, 0.4, spec.kicker or " ", size=14, color=p.theme.cover_muted, name="witty-kicker")
    p.text(MARGIN_X, 2.85, 12.0, 1.6, spec.title or " ", size=36, color=p.theme.cover_ink, bold=True, name="witty-title")
    if spec.subtitle:
        p.text(MARGIN_X, 4.6, 12.0, 0.8, spec.subtitle, size=16, color=p.theme.cover_muted, name="witty-body")
    p.footer(footer, page)


def _bullets(p: Painter, spec: Slide, page: int, footer: str) -> None:
    p.paint_bg(p.theme.bg)
    p.chrome()
    y = _header_block(p, spec)
    p.bullets(MARGIN_X, y, 12.05, CONTENT_MAX_Y - y, spec.items, size=20, bullet_color=p.theme.accent)
    p.footer(footer, page)


def _two_col(p: Painter, spec: Slide, page: int, footer: str) -> None:
    p.paint_bg(p.theme.bg)
    p.chrome()
    y = _header_block(p, spec)
    height = CONTENT_MAX_Y - y
    gap = 0.28
    width = (12.05 - gap) / 2
    p.card(MARGIN_X, y, width, height, p.theme.card, name="witty-card-left")
    p.card(MARGIN_X + width + gap, y, width, height, p.theme.card, name="witty-card-right")
    p.text(MARGIN_X + 0.22, y + 0.18, width - 0.4, 0.4, spec.left_title or " ", size=16, color=p.theme.accent, bold=True)
    p.bullets(MARGIN_X + 0.18, y + 0.62, width - 0.36, height - 0.8, spec.left, size=16, name="witty-body", bullet_color=p.theme.accent)
    p.text(MARGIN_X + width + gap + 0.22, y + 0.18, width - 0.4, 0.4, spec.right_title or " ", size=16, color=p.theme.accent, bold=True)
    p.bullets(
        MARGIN_X + width + gap + 0.18,
        y + 0.62,
        width - 0.36,
        height - 0.8,
        spec.right,
        size=16,
        name="witty-right",
        bullet_color=p.theme.accent,
    )
    p.footer(footer, page)


def _kpi(p: Painter, spec: Slide, page: int, footer: str) -> None:
    p.paint_bg(p.theme.bg)
    p.chrome()
    y = _header_block(p, spec)
    metrics = spec.metrics[:4]
    if not metrics:
        p.footer(footer, page)
        return
    cols, rows = (2, 2) if len(metrics) == 4 else (len(metrics), 1)
    gap = 0.22
    area_h = CONTENT_MAX_Y - y
    cell_w = (12.05 - gap * (cols - 1)) / cols
    cell_h = (area_h - gap * (rows - 1)) / rows
    # 一行排的时候卡片别拉到通栏高：数字页留白比撑满好看。
    if rows == 1:
        cell_h = min(cell_h, 2.5)
    for index, metric in enumerate(metrics):
        row, col = divmod(index, cols)
        left = MARGIN_X + col * (cell_w + gap)
        top = y + row * (cell_h + gap)
        p.card(left, top, cell_w, cell_h, p.theme.card, name=f"witty-card-{index}")
        p.rect(left + 0.28, top + 0.3, 0.34, 0.05, p.theme.accent, name=f"witty-kpi-rule-{index}")
        p.text(left + 0.28, top + 0.48, cell_w - 0.56, 0.34, metric.label or " ", size=13, color=p.theme.muted)
        p.text(
            left + 0.28,
            top + 0.88,
            cell_w - 0.56,
            0.9,
            metric.value or " ",
            size=32,
            color=p.theme.ink,
            bold=True,
            name=f"witty-kpi-{index}",
        )
        if metric.note:
            p.text(left + 0.28, top + 1.82, cell_w - 0.56, 0.42, metric.note, size=12, color=p.theme.accent)
    p.footer(footer, page)


def _cards(p: Painter, spec: Slide, page: int, footer: str) -> None:
    p.paint_bg(p.theme.bg)
    p.chrome()
    y = _header_block(p, spec)
    cards = spec.cards[:3]
    if not cards:
        p.footer(footer, page)
        return
    gap = 0.24
    width = (12.05 - gap * (len(cards) - 1)) / len(cards)
    height = CONTENT_MAX_Y - y
    for index, card in enumerate(cards):
        left = MARGIN_X + index * (width + gap)
        p.card(left, y, width, height, p.theme.card, name=f"witty-card-{index}")
        # 色条画在卡内，不贴边：贴边的直角条会从圆角里探出尖角。
        p.rect(left + 0.28, y + 0.32, 0.42, 0.06, p.theme.accent, name=f"witty-card-rule-{index}")
        p.text(left + 0.28, y + 0.54, width - 0.56, 0.68, card.title or " ", size=18, color=p.theme.ink, bold=True)
        p.text(
            left + 0.28,
            y + 1.3,
            width - 0.56,
            height - 1.62,
            card.body,
            size=14,
            color=p.theme.muted,
            name=f"witty-card-body-{index}",
        )
    p.footer(footer, page)


def _fill_table(p: Painter, table, headers: list[str], body_rows: list[list[str]], cols: int, size: int) -> None:
    """按内容填格并统一样式。行高按字号给，别让引擎把最后一行拉长。"""
    unit = table_row_height(size)
    table.first_row = False
    table.horz_banding = False
    for col, header in enumerate(headers[:cols]):
        cell = table.cell(0, col)
        cell.text = header
        _style_cell(p, cell, header=True, size=size)
    for r_index, row in enumerate(body_rows, start=1):
        for col in range(cols):
            cell = table.cell(r_index, col)
            cell.text = row[col] if col < len(row) else ""
            _style_cell(p, cell, header=False, stripe=r_index % 2 == 0, size=size)
    for index, row in enumerate(table.rows):
        row.height = p.inches(unit * (1.18 if index == 0 else 1.0))


def _table(p: Painter, spec: Slide, page: int, footer: str) -> None:
    p.paint_bg(p.theme.bg)
    p.chrome()
    y = _header_block(p, spec)
    headers = spec.headers or (spec.rows[0] if spec.rows else [])
    rows = spec.rows
    if spec.headers:
        body_rows = rows
    elif rows:
        headers = rows[0]
        body_rows = rows[1:]
    else:
        p.footer(footer, page)
        return
    cols = max(len(headers), 1)
    table_rows = 1 + len(body_rows)
    size = 14 if table_rows <= 7 else 12
    height = min(CONTENT_MAX_Y - y, table_row_height(size) * (table_rows + 0.18))
    table_shape = p.slide.shapes.add_table(
        table_rows,
        cols,
        p.inches(MARGIN_X),
        p.inches(y),
        p.inches(12.05),
        p.inches(height),
    )
    table_shape.name = "witty-table"
    _fill_table(p, table_shape.table, headers, body_rows, cols, size)
    p.footer(footer, page)


def _style_cell(p: Painter, cell, *, header: bool, stripe: bool = False, size: int = 13) -> None:
    fill = cell.fill
    fill.solid()
    if header:
        fill.fore_color.rgb = p.rgb(*p.theme.accent2)
        color = _on(p.theme, p.theme.accent2)
    elif stripe:
        fill.fore_color.rgb = p.rgb(*p.theme.surface)
        color = p.theme.ink
    else:
        fill.fore_color.rgb = p.rgb(*p.theme.card)
        color = p.theme.ink
    _cell_border(cell, p.qn, p.etree, p.theme.accent2 if header else p.theme.line, side="a:lnB")
    frame = cell.text_frame
    frame.word_wrap = True
    frame.margin_left = p.inches(0.14)
    frame.margin_right = p.inches(0.14)
    frame.margin_top = p.inches(0.06)
    frame.margin_bottom = p.inches(0.06)
    try:
        cell.vertical_anchor = p.anchor.MIDDLE
    except (AttributeError, TypeError, ValueError):
        pass
    for para in frame.paragraphs:
        try:
            para.line_spacing = 1.14
        except (AttributeError, TypeError, ValueError):
            pass
        for run in para.runs:
            _apply_font(
                run,
                p.theme.font,
                size + 1 if header else size,
                color,
                bold=header,
                italic=False,
                rgb_mod=p.rgb,
                pt_mod=p.pt,
                qn=p.qn,
                etree=p.etree,
            )


def _process(p: Painter, spec: Slide, page: int, footer: str) -> None:
    p.paint_bg(p.theme.bg)
    p.chrome()
    y = _header_block(p, spec)
    steps = spec.steps[:5]
    if not steps:
        p.footer(footer, page)
        return
    count = len(steps)
    width = 12.05
    gap = 0.2
    cell = (width - gap * (count - 1)) / count
    cy = y + 0.55
    if count > 1:
        p.rect(MARGIN_X + cell / 2, cy + 0.18, width - cell, 0.045, p.theme.line, name="witty-process-line")
    for index, step in enumerate(steps):
        left = MARGIN_X + index * (cell + gap)
        p.card(left + cell / 2 - 0.22, cy, 0.44, 0.44, p.theme.accent, name=f"witty-step-dot-{index}")
        p.text(
            left + cell / 2 - 0.22,
            cy,
            0.44,
            0.44,
            str(index + 1),
            size=14,
            color=_on(p.theme, p.theme.accent),
            bold=True,
            align="center",
            anchor="middle",
        )
        p.text(left, cy + 0.62, cell, 0.5, step.title or " ", size=15, color=p.theme.ink, bold=True, align="center")
        if step.body:
            p.text(left, cy + 1.15, cell, 2.4, step.body, size=13, color=p.theme.muted, align="center")
    p.footer(footer, page)


def _compare(p: Painter, spec: Slide, page: int, footer: str) -> None:
    p.paint_bg(p.theme.bg)
    p.chrome()
    y = _header_block(p, spec)
    height = CONTENT_MAX_Y - y
    gap = 0.28
    width = (12.05 - gap) / 2
    for index, (title, items, accent, name) in enumerate(
        (
            (spec.left_title, spec.left, p.theme.accent2, "witty-body"),
            (spec.right_title, spec.right, p.theme.accent, "witty-right"),
        )
    ):
        left = MARGIN_X + index * (width + gap)
        p.card(left, y, width, height, p.theme.card, name=f"witty-card-{'left' if index == 0 else 'right'}")
        # 标题条不贴卡边：直角条压在圆角卡上会露尖角，改成卡内一小段色块。
        p.rect(left + 0.26, y + 0.3, 0.42, 0.06, accent, name=f"witty-compare-rule-{index}")
        p.text(left + 0.26, y + 0.52, width - 0.52, 0.42, title or " ", size=16, color=accent, bold=True)
        p.bullets(left + 0.22, y + 1.06, width - 0.44, height - 1.24, items, size=16, name=name, bullet_color=accent)
    p.footer(footer, page)


def _quote(p: Painter, spec: Slide, page: int, footer: str) -> None:
    p.paint_bg(p.theme.bg)
    p.chrome()
    p.text(MARGIN_X, 1.35, 12.0, 1.1, "“", size=72, color=p.theme.accent, name="witty-quote-mark")
    body = spec.quote or spec.title or " "
    p.text(
        MARGIN_X,
        2.4,
        12.0,
        2.4,
        body,
        size=26,
        color=p.theme.ink,
        italic=_latin_only(body),
        name="witty-title",
    )
    if spec.by:
        p.rect(MARGIN_X, 5.02, 0.9, 0.05, p.theme.accent, name="witty-rule")
        p.text(MARGIN_X, 5.2, 12.0, 0.45, spec.by, size=14, color=p.theme.muted, name="witty-body")
    p.footer(footer, page)


def _closing(p: Painter, spec: Slide, page: int, footer: str) -> None:
    p.paint_bg(p.theme.bg)
    p.chrome()
    y = _header_block(p, spec)
    p.bullets(MARGIN_X, y, 12.05, CONTENT_MAX_Y - y - 0.7, spec.items, size=20, bullet_color=p.theme.accent)
    p.rect(MARGIN_X, CONTENT_MAX_Y - 0.12, 12.05, 0.08, p.theme.accent, name="witty-cta")
    p.footer(footer, page)


def _picture(p: Painter, spec: Slide, page: int, footer: str) -> None:
    p.paint_bg(p.theme.bg)
    p.chrome()
    band = p.theme.chrome == "band"
    top = (BAND_BOTTOM + 0.22) if band else 0.5
    if spec.title:
        p.text(MARGIN_X, top, 12.0, 0.5, spec.title, size=22, color=p.theme.accent2 if band else p.theme.ink, bold=True, name="witty-title")
        top += 0.73
    image = Path(spec.image).expanduser()
    if not image.is_file():
        raise FileNotFoundError(get_prompt("pptx_image_missing", path=str(image)))
    p.picture(image, MARGIN_X, top, 12.05, name="witty-image")
    p.footer(footer, page)


def _box_color(p: Painter, value: str, fallback) -> tuple[int, int, int]:
    return color_of(p.theme, value, fallback)


def _inside(outer: Box, inner: Box) -> bool:
    return (
        outer.x - 0.05 <= inner.x
        and outer.y - 0.05 <= inner.y
        and outer.x + outer.w + 0.05 >= inner.x + inner.w
        and outer.y + outer.h + 0.05 >= inner.y + inner.h
    )


def _behind(p: Painter, box: Box, earlier: list[Box]) -> tuple[int, int, int]:
    """这个盒子实际压在什么颜色上：最近一个套住它的色块，否则页底。"""
    for other in reversed(earlier):
        if other.kind in {"rect", "round"} and other.fill and _inside(other, box):
            return _box_color(p, other.fill, p.bg)
    return p.bg


def _bullet_tint(p: Painter, box: Box, earlier: list[Box]) -> tuple[int, int, int] | None:
    """符号用强调色提气；和身后底色贴太近就退回跟字同色。"""
    bg = _box_color(p, box.fill, p.bg) if box.fill else _behind(p, box, earlier)
    accent = p.theme.accent
    return accent if contrast_ratio(accent, bg) >= 2.2 else None


def _chart_type_enum(name: str):
    from pptx.enum.chart import XL_CHART_TYPE

    return {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }.get(name, XL_CHART_TYPE.COLUMN_CLUSTERED)


def _chart_palette(p: Painter, box: Box) -> list[tuple[int, int, int]]:
    if box.colors:
        return [_box_color(p, item, p.theme.accent) for item in box.colors]
    return [p.theme.accent, p.theme.accent2, p.theme.bar, p.theme.ink, p.theme.muted]


def _chart_text(p: Painter, chart) -> None:
    """图表默认是黑字宋体小字。统一成主题字体和次要字色。"""
    try:
        font = chart.font
    except (AttributeError, ValueError):
        return
    try:
        font.size = p.pt(11)
        font.name = p.theme.font
        font.color.rgb = p.rgb(*p.theme.muted)
    except (AttributeError, TypeError, ValueError):
        pass


def _chart_axes(p: Painter, chart, horizontal: bool) -> None:
    """留一组浅横线做读数参考，另一向的网格全关。"""
    try:
        value_axis = chart.value_axis
    except (AttributeError, ValueError):
        value_axis = None
    try:
        category_axis = chart.category_axis
    except (AttributeError, ValueError):
        category_axis = None
    if value_axis is not None:
        try:
            value_axis.has_major_gridlines = not horizontal
            if not horizontal:
                value_axis.major_gridlines.format.line.color.rgb = p.rgb(*p.theme.line)
                value_axis.major_gridlines.format.line.width = p.pt(0.6)
        except (AttributeError, TypeError, ValueError):
            pass
        try:
            value_axis.format.line.color.rgb = p.rgb(*p.theme.line)
            value_axis.has_minor_gridlines = False
        except (AttributeError, TypeError, ValueError):
            pass
    if category_axis is not None:
        try:
            category_axis.has_major_gridlines = horizontal
            if horizontal:
                category_axis.major_gridlines.format.line.color.rgb = p.rgb(*p.theme.line)
            category_axis.has_minor_gridlines = False
            category_axis.format.line.color.rgb = p.rgb(*p.theme.line)
        except (AttributeError, TypeError, ValueError):
            pass


def _chart_bars(p: Painter, chart) -> None:
    """柱子瘦一点，别糊成一片。"""
    try:
        plot = chart.plots[0]
    except (AttributeError, IndexError, ValueError):
        return
    try:
        plot.gap_width = 90
    except (AttributeError, TypeError, ValueError):
        pass
    try:
        plot.overlap = -8
    except (AttributeError, TypeError, ValueError):
        pass


def _chart_slices(p: Painter, chart) -> None:
    """饼/环标签：深色大字 + 数值和百分比。

    以前写死白字 11pt、只显示百分比。浅黄/浅绿扇区上几乎看不见
    （调研满意度 PPT 第 7 页用户反馈「数字都看不清楚」）。深墨字在
    浅扇区清晰，深色扇区也能读；字号 14 在 16:9 图上够辨认。
    """
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        labels = plot.data_labels
        labels.show_percentage = True
        labels.show_value = True
        labels.number_format = "0%"
        labels.number_format_is_linked = False
        labels.font.size = p.pt(14)
        labels.font.bold = True
        labels.font.name = p.theme.font
        labels.font.color.rgb = p.rgb(18, 63, 60)
    except (AttributeError, IndexError, TypeError, ValueError):
        pass


def _paint_chart(p: Painter, box: Box, name: str) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_LEGEND_POSITION

    categories = box.categories
    series = box.series
    if not categories or not series:
        return
    data = CategoryChartData()
    data.categories = categories
    width = len(categories)
    for item in series:
        values = list(item.values[:width])
        if len(values) < width:
            values.extend([0.0] * (width - len(values)))
        data.add_series(item.name or get_prompt("pptx_series_default"), tuple(values))
    frame = p.slide.shapes.add_chart(
        _chart_type_enum(box.chart or "column"),
        p.inches(box.x),
        p.inches(box.y),
        p.inches(box.w),
        p.inches(box.h),
        data,
    )
    frame.name = name
    chart = frame.chart
    kind = box.chart or "column"
    round_kind = kind in {"pie", "doughnut"}
    # 单系列柱线图的图例只是重复标题，白占地方。
    chart.has_legend = len(series) > 1 or (round_kind and len(categories) > 1)
    if chart.has_legend:
        chart.legend.include_in_layout = False
        try:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        except (AttributeError, ValueError):
            pass
    _chart_text(p, chart)
    if box.text:
        chart.has_title = True
        try:
            title_frame = chart.chart_title.text_frame
            title_frame.paragraphs[0].text = box.text
            title_run = title_frame.paragraphs[0].runs[0]
            title_run.font.size = p.pt(13)
            title_run.font.bold = True
            title_run.font.name = p.theme.font
            title_run.font.color.rgb = p.rgb(*p.theme.ink)
        except (AttributeError, IndexError, ValueError):
            pass
    else:
        chart.has_title = False
    if round_kind:
        _chart_slices(p, chart)
    else:
        _chart_axes(p, chart, horizontal=kind == "bar")
        _chart_bars(p, chart)
    palette = _chart_palette(p, box)
    try:
        plot_series = list(chart.series)
    except (AttributeError, ValueError):
        plot_series = []
    if round_kind and plot_series:
        try:
            points = list(plot_series[0].points)
        except (AttributeError, ValueError):
            points = []
        for index, point in enumerate(points):
            try:
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = p.rgb(*palette[index % len(palette)])
                point.format.line.color.rgb = p.rgb(*p.theme.surface)
                point.format.line.width = p.pt(1.5)
            except (AttributeError, TypeError, ValueError):
                continue
        return
    for index, item in enumerate(plot_series):
        color = palette[index % len(palette)]
        if kind == "line":
            try:
                item.format.line.color.rgb = p.rgb(*color)
                item.format.line.width = p.pt(2.25)
                item.smooth = False
            except (AttributeError, TypeError, ValueError):
                pass
            continue
        try:
            fill = item.format.fill
            fill.solid()
            fill.fore_color.rgb = p.rgb(*color)
        except (AttributeError, TypeError, ValueError):
            pass
        try:
            item.format.line.fill.background()
        except (AttributeError, TypeError, ValueError):
            pass


def _paint_boxes(p: Painter, spec: Slide, page: int, footer: str) -> None:
    p.paint_bg(_box_color(p, spec.bg, p.theme.bg))
    # band 主题的内容页先垫顶部白带压线，boxes 画在上面——自排页和宏页是同一套模板。
    if wants_band(p.theme, spec.kind, spec.boxes, p.bg):
        p.band_chrome()
    for index, box in enumerate(spec.boxes):
        name = box.name or f"witty-box-{index}"
        fill = _box_color(p, box.fill, p.theme.card) if box.fill else None
        ink = _box_color(p, box.color, p.theme.ink)
        stroke = _box_color(p, box.stroke, p.theme.accent) if box.stroke else None
        if box.kind == "line":
            if box.point:
                # 带朝向的 line 是连接线，画细线加箭头；不带朝向就还是一条色条。
                p.connector(
                    box.x,
                    box.y,
                    box.w,
                    box.h,
                    fill or stroke or p.theme.accent,
                    name=name,
                    point=box.point,
                    weight=box.stroke_w or 1.5,
                )
                continue
            p.rect(box.x, box.y, box.w, box.h, fill or p.theme.accent, name=name, shadow=bool(box.shadow))
            continue
        if box.kind in SHAPE_KINDS:
            use_shadow = box.shadow if box.shadow is not None else box.kind == "round"
            # 描了边就允许无填充（架构图的空心框）；没描边还没填色才回落卡片底色。
            body = fill if fill is not None else (None if stroke is not None else p.theme.card)
            if box.kind == "round":
                shape = p.card(
                    box.x, box.y, box.w, box.h, body, name=name, radius=box.radius or 0.12, shadow=use_shadow
                )
                if stroke is not None and box.stroke_w > 0:
                    shape.line.color.rgb = p.rgb(*stroke)
                    shape.line.width = p.pt(box.stroke_w)
            elif box.kind == "rect":
                shape = p.shape(
                    p.shapes.RECTANGLE, box.x, box.y, box.w, box.h, body,
                    name=name, shadow=use_shadow, stroke=stroke, stroke_w=box.stroke_w,
                )
            else:
                shape = p.vector(
                    box.kind, box.x, box.y, box.w, box.h, body,
                    name=name, point=box.point or "right", shadow=use_shadow,
                    stroke=stroke, stroke_w=box.stroke_w,
                )
            if box.text:
                # 形状里的字默认居中——流程带、圆点、菱形靠边写都难看。
                p.write_text(
                    shape.text_frame,
                    box.text,
                    size=box.size,
                    color=ink,
                    bold=box.bold,
                    italic=box.italic,
                    align=box.align if box.align != "left" else "center",
                    anchor="middle",
                    font=box.font or None,
                    name=name,
                )
            continue
        if box.kind == "text":
            p.text(
                box.x,
                box.y,
                box.w,
                box.h,
                box.text or " ",
                size=box.size,
                color=ink,
                bold=box.bold,
                italic=box.italic,
                align=box.align,
                name=name,
                anchor=box.anchor,
                font=box.font or None,
            )
            continue
        if box.kind == "bullets":
            p.bullets(
                box.x,
                box.y,
                box.w,
                box.h,
                box.items or ([box.text] if box.text else []),
                size=box.size,
                color=ink,
                name=name,
                bullet_color=_bullet_tint(p, box, spec.boxes[:index]),
            )
            continue
        if box.kind == "table":
            headers = box.headers or (box.rows[0] if box.rows else [])
            body_rows = box.rows if box.headers else box.rows[1:]
            cols = max(len(headers), 1)
            table_rows = 1 + len(body_rows)
            if table_rows < 1 or not headers:
                continue
            # 表格字号自己压一档：正文 18 直接进单元格会顶满。
            size = max(11, min(box.size, 15 if table_rows <= 7 else 13))
            height = min(box.h, table_row_height(size) * (table_rows + 0.18))
            table_shape = p.slide.shapes.add_table(
                table_rows,
                cols,
                p.inches(box.x),
                p.inches(box.y),
                p.inches(box.w),
                p.inches(height),
            )
            table_shape.name = name
            _fill_table(p, table_shape.table, headers, body_rows, cols, size)
            continue
        if box.kind == "chart":
            _paint_chart(p, box, name)
            continue
        if box.kind == "image":
            image = Path(box.image).expanduser()
            if not image.is_file():
                raise FileNotFoundError(get_prompt("pptx_image_missing", path=str(image)))
            p.picture(image, box.x, box.y, box.w, name=name)
    if wants_logo(p.theme, spec.boxes):
        p.logo(dark=is_dark(logo_ground(p.theme, spec.boxes, p.bg)), cover=spec.kind == "cover")
    if not any(box.name == "witty-footer" for box in spec.boxes):
        p.footer(footer, page)


_DISPATCH = {
    "cover": _cover,
    "section": _section,
    "bullets": _bullets,
    "two_col": _two_col,
    "kpi": _kpi,
    "cards": _cards,
    "table": _table,
    "process": _process,
    "compare": _compare,
    "quote": _quote,
    "closing": _closing,
    "picture": _picture,
}


def _paint_slide(deck, spec: Slide, theme: Theme, page: int, footer: str) -> None:
    slide = deck.slides.add_slide(_blank(deck))
    painter = Painter(slide, theme, _mods())
    if spec.boxes:
        _paint_boxes(painter, spec, page, footer)
    else:
        fn = _DISPATCH.get(spec.kind)
        if fn is None:
            raise ValueError(get_prompt("pptx_bad_kind", kind=spec.kind or "(空)"))
        fn(painter, spec, page, footer)
    painter.notes(spec.notes)


def _existing_bg(slide, theme: Theme) -> tuple[int, int, int]:
    """从已有页里猜底色：整幅色块优先，其次页面填充，最后主题底。"""
    for shape in slide.shapes:
        try:
            width = shape.width / EMU_PER_INCH
            height = shape.height / EMU_PER_INCH
        except (AttributeError, TypeError):
            continue
        if width < CANVAS_W - 0.12 or height < CANVAS_H - 0.12:
            continue
        try:
            fill = shape.fill
            if fill.type is not None and fill.fore_color.type is not None:
                rgb = fill.fore_color.rgb
                return (rgb[0], rgb[1], rgb[2])
        except (AttributeError, TypeError, ValueError):
            continue
    return theme.bg


def _refresh_footers(deck, theme: Theme, footer: str) -> None:
    painter_mods = _mods()
    for index, slide in enumerate(deck.slides, start=1):
        found = None
        for shape in slide.shapes:
            if shape.name == "witty-footer" and shape.has_text_frame:
                found = shape
                break
        if found is None:
            continue
        label = f"{footer}  ·  {index}" if footer else str(index)
        frame = found.text_frame
        if frame.paragraphs and frame.paragraphs[0].runs:
            frame.paragraphs[0].runs[0].text = label
            continue
        tint = _muted_on(theme, _existing_bg(slide, theme))
        Painter(slide, theme, painter_mods).text(
            MARGIN_X, FOOTER_TOP, 10.6, 0.32, label, size=10, color=tint, name="witty-footer"
        )


def _theme_shape(pres):
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.name == "witty-theme" and shape.has_text_frame:
                return shape
    return None


def _store_theme(pres, theme: Theme) -> None:
    blob = json.dumps(theme_payload(theme), ensure_ascii=False, separators=(",", ":"))
    found = _theme_shape(pres)
    if found is not None:
        found.text_frame.text = blob
        return
    slides = list(pres.slides)
    if not slides:
        return
    Painter(slides[0], theme, _mods()).text(-0.25, -0.25, 0.2, 0.2, blob, size=8, color=theme.bg, name="witty-theme")


def load_stored_theme(pres) -> Theme | None:
    found = _theme_shape(pres)
    if found is None:
        return None
    raw = "\n".join(para.text for para in found.text_frame.paragraphs).strip()
    if raw.startswith(_THEME_MARK):
        raw = raw[len(_THEME_MARK) :].strip()
    try:
        payload = json.loads(raw)
    except json.JSONDecodeError:
        return None
    if not isinstance(payload, dict):
        return None
    return theme_from_payload({str(key): str(value) for key, value in payload.items()})


def _ready_theme(theme: Theme) -> Theme:
    font = resolve_installed_font(theme.font)
    return theme if font == theme.font else replace(theme, font=font)


def _delete_slide(pres, index: int) -> None:
    sld_id_lst = pres.slides._sldIdLst
    sld_id = sld_id_lst[index]
    pres.part.drop_rel(sld_id.rId)
    sld_id_lst.remove(sld_id)


def _move_slide(pres, old_index: int, new_index: int) -> None:
    if old_index == new_index:
        return
    sld_id_lst = pres.slides._sldIdLst
    sld_id = sld_id_lst[old_index]
    sld_id_lst.remove(sld_id)
    sld_id_lst.insert(new_index, sld_id)


def atomic_save(pres, target: Path) -> None:
    """先写临时文件再原子换名。分块生成写到几十页时进程被杀，
    不能拿半个损坏的 zip 盖掉已经排好的稿子——原文件要么旧要么新，不会烂。"""
    tmp = target.with_name(target.name + ".tmp~")
    try:
        pres.save(str(tmp))
        tmp.replace(target)
    finally:
        tmp.unlink(missing_ok=True)


def write_pptx(deck, path: str, theme: Theme | None = None) -> str:
    Presentation, *_rest = _deps()
    resolved = _ready_theme(theme or resolve_theme(deck.theme, deck.theme_overrides))
    pres = Presentation()
    _set_wide(pres)
    pres.core_properties.title = deck.title
    footer = deck.footer or deck.title
    for index, spec in enumerate(deck.slides, start=1):
        _paint_slide(pres, spec, resolved, index, footer)
    _store_theme(pres, resolved)
    target = Path(path).expanduser()
    target.parent.mkdir(parents=True, exist_ok=True)
    atomic_save(pres, target)
    logger.info("写入PPTX path=%s pages=%s theme=%s", target, len(deck.slides), resolved.id)
    return str(target)


def insert_slide(path: str, spec: Slide, index: int, theme_id: str = "") -> str:
    Presentation, *_rest = _deps()
    target = Path(path).expanduser()
    pres = Presentation(str(target))
    total = len(list(pres.slides))
    if index < 1 or index > total + 1:
        raise ValueError(get_prompt("pptx_bad_index", index=str(index), total=str(total)))
    if index == total + 1:
        return append_slide(path, spec, theme_id)
    append_slide(path, spec, theme_id)
    pres = Presentation(str(target))
    last = len(list(pres.slides)) - 1
    _move_slide(pres, last, index - 1)
    theme = load_stored_theme(pres) or _ready_theme(resolve_theme(theme_id or "custom"))
    footer = str(getattr(pres.core_properties, "title", "") or "")
    _refresh_footers(pres, theme, footer)
    atomic_save(pres, target)
    return str(target)


def append_slides(path: str, specs: list[Slide], theme_id: str = "") -> str:
    """一次追加多页：开一次文件、画完所有页、存一次。

    长稿是分批生成的，一页一次 append_slide 就要开关文件 N 次，页数上去之后
    光是重复解析 XML 就够慢的。批量走这条路。
    """
    Presentation, *_rest = _deps()
    target = Path(path).expanduser()
    pres = Presentation(str(target))
    _set_wide(pres)
    if theme_id.strip():
        theme = _ready_theme(resolve_theme(theme_id))
    else:
        theme = load_stored_theme(pres) or _ready_theme(resolve_theme("custom"))
    footer = str(getattr(pres.core_properties, "title", "") or "")
    page = len(list(pres.slides))
    for spec in specs:
        page += 1
        _paint_slide(pres, spec, theme, page, footer)
    _refresh_footers(pres, theme, footer)
    _store_theme(pres, theme)
    atomic_save(pres, target)
    logger.info("追加PPTX path=%s added=%s total=%s theme=%s", target, len(specs), page, theme.id)
    return str(target)


def append_slide(path: str, spec: Slide, theme_id: str = "") -> str:
    return append_slides(path, [spec], theme_id)


def replace_slide(path: str, spec: Slide, index: int, theme_id: str = "") -> str:
    Presentation, *_rest = _deps()
    target = Path(path).expanduser()
    pres = Presentation(str(target))
    slides = list(pres.slides)
    if index < 1 or index > len(slides):
        raise ValueError(get_prompt("pptx_bad_index", index=str(index), total=str(len(slides))))
    if theme_id.strip():
        theme = _ready_theme(resolve_theme(theme_id))
    else:
        theme = load_stored_theme(pres) or _ready_theme(resolve_theme("custom"))
    footer = str(getattr(pres.core_properties, "title", "") or "")
    _delete_slide(pres, index - 1)
    _paint_slide(pres, spec, theme, index, footer)
    last = len(list(pres.slides)) - 1
    _move_slide(pres, last, index - 1)
    _refresh_footers(pres, theme, footer)
    _store_theme(pres, theme)
    atomic_save(pres, target)
    return str(target)


def shape_fill_rgb(shape: Any) -> tuple[int, int, int] | None:
    try:
        color = shape.fill.fore_color.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    return int(color[0]), int(color[1]), int(color[2])
