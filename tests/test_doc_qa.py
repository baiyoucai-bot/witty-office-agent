from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from witty_agent.kernel_surface import KERNEL_TOOLS
from witty_agent.loop import READONLY_TOOLS
from witty_agent.plugins.doc_qa import doc_qa
from witty_agent.plugins.pptx import pptx_add_slide, pptx_create
from witty_agent.skills import match_relevant_skills
from witty_agent.tools import list_tools
from witty_agent.tool_surface import select_advertised_names


class DocQaPluginTests(unittest.TestCase):
    def test_not_a_kernel_tool(self) -> None:
        self.assertNotIn("doc_qa", KERNEL_TOOLS)
        names = {item.name for item in list_tools()}
        self.assertIn("doc_qa", names)
        self.assertIn("doc_qa", READONLY_TOOLS)

    def test_missing_and_unsupported(self) -> None:
        self.assertIn("找不到文件", doc_qa("/tmp/no-such-deck.pptx"))
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "notes.csv"
            path.write_text("a,b\n", encoding="utf-8")
            self.assertIn("暂不质检", doc_qa(str(path)))

    def test_flags_pptx_and_markdown(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            deck = str(Path(tmp) / "deck.pptx")
            pptx_create(deck, "封面", "内网")
            pptx_add_slide(deck, "要点", "1\n2\n3\n4\n5\n6\n7")
            from pptx import Presentation

            packed = Presentation(deck)
            body = [
                shape
                for shape in packed.slides[1].shapes
                if shape.has_text_frame and shape.name != "witty-footer"
            ][-1]
            extra = body.text_frame.add_paragraph()
            extra.text = "8"
            packed.save(deck)
            pptx_add_slide(deck, "占位", "TODO 待补充")
            report = doc_qa(deck)
            self.assertIn("问题", report)
            self.assertIn("要点过多", report)
            self.assertIn("占位句", report)
            clean = Path(tmp) / "ok.md"
            clean.write_text("# 标题\n\n一段话。\n", encoding="utf-8")
            self.assertIn("通过", doc_qa(str(clean)))
            messy = Path(tmp) / "draft.md"
            messy.write_text("TODO 待补充\n", encoding="utf-8")
            text = doc_qa(str(messy))
            self.assertIn("没有标题", text)
            self.assertIn("占位句", text)

    def test_skill_matches_qa_prompt(self) -> None:
        names = [item.name for item in match_relevant_skills("质检这份PPT版式")]
        self.assertEqual(names[:1], ["doc-qa"])
        self.assertNotIn("doc-qa", [item.name for item in match_relevant_skills("做个PPT")])

    def test_advertised_on_qa_prompt(self) -> None:
        names = select_advertised_names("质检这份PPT版式", ["doc_qa", "pptx_outline", "read"])
        self.assertIn("doc_qa", names)
        idle = select_advertised_names("review the auth module", ["doc_qa", "read"])
        self.assertNotIn("doc_qa", idle)
