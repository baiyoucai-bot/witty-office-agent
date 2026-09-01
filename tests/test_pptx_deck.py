from __future__ import annotations

import json
import re
import tempfile
import unittest
from pathlib import Path

from witty_agent.loop import READONLY_TOOLS
from witty_agent.plan_mode import MUTATING_TOOLS
from witty_agent.plugins.pptx import (
    pptx_add_page,
    pptx_add_pages,
    pptx_add_slide,
    pptx_check,
    pptx_create,
    pptx_edit_box,
    pptx_edit_slide,
    pptx_from_html,
    pptx_list_boxes,
    pptx_outline,
    pptx_render,
    pptx_replace_slide,
    pptx_themes,
)
from witty_agent.plugins.pptx_kit import parse_deck, parse_html, write_pptx
from witty_agent.plugins.pptx_kit.lint import lint_deck
from witty_agent.plugins.pptx_kit.render import load_stored_theme, shape_fill_rgb

from witty_agent.tools import list_tools


def _sample_deck() -> dict:
    return {
        "title": "年中汇报",
        "theme": "grid-navy",
        "footer": "内部",
        "slides": [
            {
                "kind": "cover",
                "kicker": "2026",
                "title": "年中汇报",
                "subtitle": "建设进度",
                "meta": "内部",
            },
            {"kind": "section", "kicker": "01", "title": "进展"},
            {"kind": "bullets", "kicker": "要点", "title": "三件事", "items": ["调研", "试点", "推广"]},
            {
                "kind": "two_col",
                "title": "分工",
                "left_title": "建设",
                "left": ["变电站"],
                "right_title": "运营",
                "right": ["巡检"],
            },
            {
                "kind": "kpi",
                "title": "指标",
                "metrics": [{"label": "完成率", "value": "86%", "note": "同比+12pt"}],
            },
            {"kind": "cards", "title": "路径", "cards": [{"title": "短", "body": "先试点"}]},
            {
                "kind": "table",
                "title": "对照",
                "headers": ["项", "值"],
                "rows": [["工期", "12 月"], ["投资", "待核实"]],
            },
            {"kind": "process", "title": "步骤", "steps": [{"title": "调研", "body": "现场"}]},
            {
                "kind": "compare",
                "title": "取舍",
                "left_title": "现状",
                "left": ["慢"],
                "right_title": "目标",
                "right": ["准"],
            },
            {"kind": "quote", "quote": "安全第一", "by": "现场"},
            {"kind": "closing", "title": "下一步", "items": ["发纪要"]},
        ],
    }


class PptxDeckTests(unittest.TestCase):
    def test_new_tools_are_pluggable(self) -> None:
        names = {item.name for item in list_tools()}
        for name in (
            "pptx_render",
            "pptx_from_html",
            "pptx_themes",
            "pptx_check",
            "pptx_replace_slide",
            "pptx_list_boxes",
            "pptx_edit_box",
            "pptx_add_page",
        ):
            self.assertIn(name, names)
        self.assertIn("pptx_themes", READONLY_TOOLS)
        self.assertIn("pptx_check", READONLY_TOOLS)
        self.assertIn("pptx_list_boxes", READONLY_TOOLS)
        self.assertIn("pptx_render", MUTATING_TOOLS)
        self.assertIn("pptx_from_html", MUTATING_TOOLS)
        self.assertIn("pptx_replace_slide", MUTATING_TOOLS)
        self.assertNotIn("pptx_themes", MUTATING_TOOLS)

    def test_render_native_shapes_and_table(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = str(Path(tmp) / "deck.pptx")
            text = pptx_render(path, deck=json.dumps(_sample_deck(), ensure_ascii=False))
            self.assertIn(path, text)
            html = Path(tmp) / "deck.html"
            self.assertTrue(html.is_file())
            self.assertIn("data-witty-deck", html.read_text(encoding="utf-8"))
            from pptx import Presentation

            pres = Presentation(path)
            slides = list(pres.slides)
            self.assertEqual(len(slides), 11)
            titles = []
            for slide in slides:
                named = {shape.name: shape for shape in slide.shapes}
                self.assertIn("witty-footer", named)
                if "witty-title" in named:
                    titles.append(named["witty-title"].text_frame.text)
            self.assertIn("年中汇报", titles)
            self.assertIn("三件事", titles)
            table_slide = slides[6]
            self.assertTrue(any(shape.has_table for shape in table_slide.shapes))
            pictures = [shape for slide in slides for shape in slide.shapes if shape.shape_type is not None and shape.has_text_frame]
            self.assertGreater(len(pictures), 10)
            outline = pptx_outline(path)
            self.assertIn("完成率", outline)
            self.assertIn("安全第一", outline)

    def test_theme_override_changes_bar(self) -> None:
        payload = _sample_deck()
        payload["slides"] = [payload["slides"][0]]
        payload["theme_overrides"] = {"bar": "#C41E3A"}
        with tempfile.TemporaryDirectory() as tmp:
            path = str(Path(tmp) / "red.pptx")
            write_pptx(parse_deck(payload), path)
            from pptx import Presentation

            slide = Presentation(path).slides[0]
            bar = next(shape for shape in slide.shapes if shape.name == "witty-bar")
            self.assertEqual(shape_fill_rgb(bar), (0xC4, 0x1E, 0x3A))

    def test_html_roundtrip_editable(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            src = str(Path(tmp) / "src.pptx")
            pptx_render(src, deck=json.dumps(_sample_deck(), ensure_ascii=False))
            html_path = str(Path(tmp) / "src.html")
            parsed = parse_html(Path(html_path).read_text(encoding="utf-8"))
            self.assertEqual(parsed.theme, "grid-navy")
            self.assertEqual([item.kind for item in parsed.slides][:4], ["cover", "section", "bullets", "two_col"])
            self.assertEqual(parsed.slides[2].items, ["调研", "试点", "推广"])
            out = str(Path(tmp) / "back.pptx")
            pptx_from_html(out, html_path)
            self.assertIn("三件事", pptx_outline(out))

    def test_bad_kind_and_unrecognized_html(self) -> None:
        with self.assertRaises(ValueError):
            parse_deck({"title": "x", "slides": [{"kind": "canvas", "title": "no"}]})
        with self.assertRaises(ValueError):
            parse_html("<html><body><p>hello</p></body></html>")
        catalog = pptx_themes()
        self.assertIn("grid-navy", catalog)
        self.assertIn("swiss-red", catalog)
        self.assertIn("grid", catalog)
        self.assertIn("青绿金线企业风", catalog)

    def test_grid_theme_is_teal_gold(self) -> None:
        from witty_agent.plugins.pptx_kit.themes import match_theme_id, resolve_theme

        self.assertEqual(match_theme_id("青绿"), "grid")
        self.assertEqual(match_theme_id("审计培训"), "grid")
        self.assertEqual(match_theme_id("数字化审计"), "grid")
        theme = resolve_theme("grid")
        self.assertEqual(theme.accent, (0x01, 0x70, 0x6C))
        self.assertEqual(theme.bar, (0xFF, 0xC2, 0x0A))
        self.assertEqual(theme.cover_bg, (0x01, 0x70, 0x6C))
        self.assertEqual(theme.chrome, "band")
        # 开源版不带内置标识：logo/emblem 为空
        self.assertEqual(theme.logo, "")
        self.assertEqual(theme.emblem, "")
        with tempfile.TemporaryDirectory() as tmp:
            path = str(Path(tmp) / "grid.pptx")
            write_pptx(
                parse_deck(
                    {
                        "title": "青绿",
                        "theme": "青绿",
                        "slides": [
                            {"kind": "cover", "title": "数字化审计", "kicker": "企业汇报"},
                            {"kind": "bullets", "title": "背景", "items": ["一", "二"]},
                            {
                                "kind": "custom",
                                "title": "架构",
                                "boxes": [{"kind": "text", "x": 0.62, "y": 2.0, "w": 6.0, "h": 0.5, "text": "自定义"}],
                            },
                        ],
                    }
                ),
                path,
            )
            from pptx import Presentation

            slides = Presentation(path).slides
            gold = next(shape for shape in slides[0].shapes if shape.name == "witty-band-gold")
            self.assertEqual(shape_fill_rgb(gold), (0xFF, 0xC2, 0x0A))
            # 主题不带标识：任何页都不该出现标识或徽标图片
            for index in range(3):
                names = {shape.name for shape in slides[index].shapes}
                self.assertNotIn("witty-logo", names, f"第 {index + 1} 页不该有标识")
            self.assertNotIn("witty-emblem", {shape.name for shape in slides[0].shapes})

    def test_grid_template_band_logo_and_lint(self) -> None:
        """grid 模板底线：boxes 内容页自动垫白带；开源版主题不带标识，页面不出标识；
        主题声明了标识但缺图时静默降级不炸；主题带标识时内容压标识区由 lint 报出；
        封面章节宏不算 no_boxes。"""
        from witty_agent.plugins.pptx_kit.lint import lint_slide
        from witty_agent.plugins.pptx_kit.schema import parse_slide as parse_one
        from witty_agent.plugins.pptx_kit.themes import resolve_theme

        deck = parse_deck(
            {
                "title": "青绿",
                "theme": "青绿",
                "slides": [
                    {"kind": "cover", "title": "数字化审计"},
                    {"kind": "section", "kicker": "01", "title": "总体"},
                    {
                        "kind": "custom",
                        "title": "要点",
                        "boxes": [
                            # 通栏标题压进右上标识区：主题带标识时 lint 报 logo_zone
                            {"kind": "text", "x": 0.7, "y": 0.3, "w": 12.0, "h": 0.6, "text": "要点", "size": 28, "bold": True, "name": "witty-title"},
                            {"kind": "bullets", "x": 0.7, "y": 1.5, "w": 11.9, "h": 4.6, "items": ["一", "二"], "name": "witty-body"},
                        ],
                    },
                    {
                        "kind": "custom",
                        "title": "深底页",
                        "bg": "cover_bg",
                        "boxes": [{"kind": "text", "x": 0.7, "y": 2.0, "w": 8.0, "h": 0.8, "text": "深底", "size": 30, "bold": True, "color": "cover_ink", "name": "witty-title"}],
                    },
                ],
            }
        )
        theme = resolve_theme("grid")
        with tempfile.TemporaryDirectory() as tmp:
            path = str(Path(tmp) / "grid-band.pptx")
            write_pptx(deck, path, theme)
            from pptx import Presentation

            slides = list(Presentation(path).slides)
            names = [{shape.name for shape in slide.shapes} for slide in slides]
            # 封面章节不压白带，内容 boxes 页自动垫白带；深底页不垫
            self.assertNotIn("witty-band", names[0])
            self.assertNotIn("witty-band", names[1])
            self.assertIn("witty-band", names[2])
            self.assertIn("witty-band-rule", names[2])
            self.assertNotIn("witty-band", names[3])
            # 开源版 grid 不带内置标识：每页都不该有标识
            for index, page_names in enumerate(names):
                self.assertNotIn("witty-logo", page_names, f"第 {index + 1} 页不该有标识")
            # 主题声明了标识但缺图：渲染静默跳过、不炸（预期降级）
            themed = resolve_theme("grid", {"logo": "acme-mark", "emblem": "acme-emblem"})
            degraded = str(Path(tmp) / "grid-degraded.pptx")
            write_pptx(deck, degraded, themed)
            for slide in Presentation(degraded).slides:
                self.assertNotIn("witty-logo", {shape.name for shape in slide.shapes})
        # lint：主题带标识才报 logo_zone（挪内容而不是丢标识）；无标识主题不报；
        # 封面章节宏不报 no_boxes；内容宏仍报
        codes = {item.code for item in lint_slide(deck.slides[2], theme, 3)}
        self.assertNotIn("logo_zone", codes)
        themed = resolve_theme("grid", {"logo": "acme-mark"})
        themed_codes = {item.code for item in lint_slide(deck.slides[2], themed, 3)}
        self.assertIn("logo_zone", themed_codes)
        self.assertEqual([item.code for item in lint_slide(deck.slides[0], theme, 1)], [])
        macro = parse_one({"kind": "bullets", "title": "偷懒页", "items": ["一"]})
        self.assertIn("no_boxes", {item.code for item in lint_slide(macro, theme, 9)})

    def test_user_supplied_logo_resolves_from_home_and_path(self) -> None:
        """开源版包里不带企业标识，用户得能自己把标识放上去——两条路都要通：
        `$WITTY_HOME/brand/<名字>.png` 按名字取，或者 logo 直接写文件路径。
        深底页要的 `-white` 反白版两种写法都得找得到（路径写法要按后缀前插）。
        裸名字仍然禁止分隔符和点开头，不然就成了任意路径读取。
        """
        import os

        from witty_agent.plugins.pptx_kit import assets
        from witty_agent.plugins.pptx_kit.themes import resolve_theme

        def mark(path: Path, color: tuple[int, int, int]) -> None:
            from PIL import Image

            Image.new("RGB", (600, 200), color).save(path)

        deck = parse_deck(
            {
                "title": "青绿",
                "theme": "青绿",
                "slides": [
                    {"kind": "cover", "title": "数字化审计"},
                    {
                        "kind": "custom",
                        "title": "要点",
                        "boxes": [
                            {"kind": "text", "x": 0.7, "y": 1.5, "w": 8.0, "h": 0.6, "text": "要点", "size": 28, "bold": True, "name": "witty-title"},
                        ],
                    },
                ],
            }
        )
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            brand = root / "home" / "brand"
            brand.mkdir(parents=True)
            mark(brand / "acme-mark.png", (1, 112, 108))
            mark(brand / "acme-mark-white.png", (255, 255, 255))
            loose = root / "loose.png"
            mark(loose, (255, 194, 10))
            mark(root / "loose-white.png", (255, 255, 255))

            previous = os.environ.get("WITTY_HOME")
            os.environ["WITTY_HOME"] = str(root / "home")
            assets.clear_asset_cache()
            try:
                # 按名字取用户目录里的图，反白版同名 -white。
                # 期望值也要 resolve：WITTY_HOME 走 data_root() 解析过符号链接
                # （macOS 上 /var → /private/var），比字符串会假红。
                self.assertEqual(assets.asset_path("acme-mark"), str((brand / "acme-mark.png").resolve()))
                self.assertEqual(assets.asset_path("acme-mark-white"), str((brand / "acme-mark-white.png").resolve()))
                self.assertEqual(assets.asset_size("acme-mark"), (600, 200))
                self.assertIn("acme-mark", assets.asset_names())
                # 路径写法：反白版按后缀前插，不是拼在 .png 后面
                self.assertEqual(assets.asset_path(str(loose)), str(loose))
                self.assertEqual(assets.asset_path(f"{loose}-white"), str(root / "loose-white.png"))
                # 裸名字不许带分隔符/点开头；路径写法必须是 .png 且真的存在
                self.assertEqual(assets.asset_path("../../../etc/passwd"), "")
                self.assertEqual(assets.asset_path(".secret"), "")
                self.assertEqual(assets.asset_path(str(root / "nope")), "")
                self.assertEqual(assets.asset_path(str(root / "missing.png")), "")

                themed = resolve_theme("grid", {"logo": "acme-mark", "emblem": str(loose)})
                path = str(root / "branded.pptx")
                write_pptx(deck, path, themed)
                from pptx import Presentation

                slides = list(Presentation(path).slides)
                # 用户底线：标识每页都在，封面另有大徽标
                for index, slide in enumerate(slides):
                    self.assertIn("witty-logo", {shape.name for shape in slide.shapes}, f"第 {index + 1} 页缺标识")
                self.assertIn("witty-emblem", {shape.name for shape in slides[0].shapes})
                # 预览是单文件，图得内联进去
                from witty_agent.plugins.pptx_kit.preview import render_html

                self.assertIn("data:image/png;base64,", render_html(deck, themed))
            finally:
                if previous is None:
                    os.environ.pop("WITTY_HOME", None)
                else:
                    os.environ["WITTY_HOME"] = previous
                assets.clear_asset_cache()

    def test_add_pages_continues_numbering_theme_and_preview(self) -> None:
        """长稿分块的接缝：页码/页脚接着排，主题（含用户标识）从文件里读回来不用重传，
        预览把新页拼在末尾且旧页字节不动，自检按真实页码报而不是从 1 数。
        """
        import os

        from witty_agent.plugins.pptx_kit import assets

        def mark(path: Path) -> None:
            from PIL import Image

            Image.new("RGB", (600, 200), (1, 112, 108)).save(path)

        first = {
            "title": "数字化审计",
            "theme": "青绿",
            "slides": [
                {"kind": "cover", "title": "新一代数字化审计平台"},
                {"kind": "section", "kicker": "01", "title": "平台概览"},
            ],
        }
        batch = [
            {
                "kind": "custom",
                "title": "建设目标",
                "boxes": [
                    {"kind": "text", "x": 0.62, "y": 1.40, "w": 8.0, "h": 0.62, "text": "建设目标", "size": 28, "bold": True, "name": "witty-title"},
                    {"kind": "bullets", "x": 0.62, "y": 2.30, "w": 11.9, "h": 3.9, "items": ["数据贯通", "模型驱动"], "name": "witty-body"},
                ],
            },
            # 故意留一页宏：自检要报在第 4 页（真实页码），不能报成第 2 页
            {"kind": "bullets", "title": "偷懒页", "items": ["一"]},
        ]
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            brand = root / "home" / "brand"
            brand.mkdir(parents=True)
            mark(brand / "acme-mark.png")
            previous = os.environ.get("WITTY_HOME")
            os.environ["WITTY_HOME"] = str(root / "home")
            assets.clear_asset_cache()
            try:
                path = str(root / "long.pptx")
                pptx_render(
                    path,
                    deck=json.dumps(first, ensure_ascii=False),
                    theme_overrides=json.dumps({"logo": "acme-mark"}),
                )
                html = Path(path).with_suffix(".html")
                before = html.read_text(encoding="utf-8")
                head = before[: before.rfind("</article>")]

                # 第二批不传主题：标识和配色都得从文件里记的主题读回来
                message = pptx_add_pages(path, json.dumps(batch, ensure_ascii=False))
                self.assertIn("第 3–4 页", message)
                self.assertIn("P4", message)
                self.assertNotIn("P2", message)

                from pptx import Presentation

                pres = Presentation(path)
                slides = list(pres.slides)
                self.assertEqual(len(slides), 4)
                stored = load_stored_theme(pres)
                self.assertIsNotNone(stored)
                self.assertEqual(stored.logo, "acme-mark")
                for index, slide in enumerate(slides, 1):
                    shapes = {shape.name for shape in slide.shapes}
                    self.assertIn("witty-logo", shapes, f"第 {index} 页缺标识")
                    page = next(shape for shape in slide.shapes if shape.name == "witty-page")
                    self.assertEqual(page.text_frame.text, f"{index:02d}")
                    footer = next(shape for shape in slide.shapes if shape.name == "witty-footer")
                    self.assertTrue(footer.text_frame.text.endswith(str(index)))
                    self.assertIn("数字化审计", footer.text_frame.text)

                after = html.read_text(encoding="utf-8")
                self.assertEqual(after.count("<section class='slide"), 4)
                # 旧页保持字节不动：分块到第几批都不会把前面画坏
                self.assertTrue(after.startswith(head))
            finally:
                if previous is None:
                    os.environ.pop("WITTY_HOME", None)
                else:
                    os.environ["WITTY_HOME"] = previous
                assets.clear_asset_cache()

    def test_preview_band_under_boxes_and_legacy_upgrade(self) -> None:
        """预览与成稿同模板：内容 boxes 页出白带且垫底、标识常驻；
        旧版单页预览（滚轮翻不了页那代）续页时整稿重写成滚动版式。"""
        from witty_agent.plugins.pptx_kit.preview import append_html, render_html
        from witty_agent.plugins.pptx_kit.schema import parse_slide as parse_one
        from witty_agent.plugins.pptx_kit.themes import resolve_theme

        deck = parse_deck(
            {
                "title": "青绿",
                "theme": "青绿",
                "slides": [
                    {"kind": "cover", "title": "数字化审计"},
                    {
                        "kind": "custom",
                        "title": "要点",
                        "boxes": [{"kind": "text", "x": 0.7, "y": 0.3, "w": 10.0, "h": 0.6, "text": "要点", "size": 28, "bold": True, "name": "witty-title"}],
                    },
                ],
            }
        )
        html = render_html(deck)
        sections = re.findall(r"<section[^>]*>", html)
        self.assertNotIn("class='chrome'", html.split("<section")[1])  # 封面不压白带
        self.assertIn("class='chrome'", html.split("<section")[2])  # 内容页压白带
        self.assertIn(".slide.free .chrome { z-index: 0; }", html)
        self.assertEqual(len(sections), 2)
        # 旧版单页 HTML：没有 aspect-ratio，是 display:none 翻页那代
        legacy = (
            "<!DOCTYPE html><html lang='zh-CN' data-witty-deck='1' data-theme='grid' "
            "data-title='旧稿' data-footer='旧稿'><head><style>"
            ".deck { width: 100vw; height: 100vh; overflow: hidden; }"
            ".slide { display: none; }"
            "</style></head><body><article class='deck'>"
            "<section class='slide free is-on' data-kind='cover'>"
            "<div class='box' data-box='text' data-x='1' data-y='1' data-w='6' data-h='1' "
            "data-size='30' data-name='witty-title'>旧封面</div></section>"
            "</article></body></html>"
        )
        with tempfile.TemporaryDirectory() as tmp:
            target = Path(tmp) / "legacy.html"
            target.write_text(legacy, encoding="utf-8")
            extra = parse_one(
                {
                    "kind": "custom",
                    "title": "新页",
                    "boxes": [{"kind": "text", "x": 0.7, "y": 1.5, "w": 8.0, "h": 0.6, "text": "新页", "size": 24, "bold": True, "name": "witty-title"}],
                }
            )
            out = append_html(str(target), [extra], resolve_theme("grid"))
            self.assertEqual(out, str(target))
            rebuilt = target.read_text(encoding="utf-8")
            self.assertIn("aspect-ratio", rebuilt)
            self.assertEqual(rebuilt.count("<section class='slide"), 2)
            self.assertNotIn(".slide { display: none; }", rebuilt)
            self.assertNotIn("is-on", rebuilt)

    def test_preview_resolves_theme_tokens_to_hex(self) -> None:
        """预览里的 fill/color/bg 必须是解析后的色值。

        写成 `background:accent2` 浏览器认不出来，会静默丢掉整条声明——盒子透明、
        字回落成默认墨色，预览和成稿两个样。data-* 属性得留原始写法，回读要用。
        """
        from witty_agent.plugins.pptx_kit.preview import render_html

        deck = parse_deck(
            {
                "title": "青绿",
                "theme": "青绿",
                "slides": [
                    {"kind": "section", "title": "第一章"},
                    {
                        "kind": "custom",
                        "title": "架构",
                        "bg": "accent2",
                        "boxes": [
                            {
                                "kind": "round",
                                "x": 0.62,
                                "y": 2.0,
                                "w": 6.0,
                                "h": 1.0,
                                "fill": "accent2",
                            },
                            {
                                "kind": "text",
                                "x": 0.92,
                                "y": 2.2,
                                "w": 5.4,
                                "h": 0.5,
                                "text": "基础设施层",
                                "color": "bar",
                            },
                        ],
                    },
                ],
            }
        )
        html = render_html(deck)
        self.assertIn("background:#076655;", html)
        self.assertIn("color:#ffc20a;", html)
        self.assertIn("style='background:#076655'", html)
        # 原始色名只留在 data-* 上，供 parse_html 回读
        self.assertIn("data-fill='accent2'", html)
        self.assertIn("data-color='bar'", html)
        self.assertNotIn("background:accent2", html)
        self.assertNotIn("color:bar;", html)
        # 每个 var(--x) 都得在 :root 里定义过，否则声明会被整条丢掉
        used = set(re.findall(r"var\(--([a-z0-9-]+)", html))
        declared = set(re.findall(r"^\s*--([a-z0-9-]+):", html, re.MULTILINE))
        self.assertEqual(used - declared, set(), "预览引用了未定义的 CSS 变量")
        # 章节页号数没写就顶页码，和 render._section_band 一致
        self.assertIn(">01</p>", html)

    def test_dynamic_theme_and_boxes(self) -> None:
        payload = {
            "title": "现场稿",
            "theme": {
                "bg": "#F4F1EA",
                "ink": "#1B1714",
                "accent": "#C45C26",
                "cover_bg": "#1B1714",
                "cover_ink": "#F4F1EA",
            },
            "slides": [
                {
                    "bg": "#1B1714",
                    "boxes": [
                        {
                            "kind": "rect",
                            "x": 0,
                            "y": 0,
                            "w": 0.2,
                            "h": 7.5,
                            "fill": "#C45C26",
                            "name": "witty-bar",
                        },
                        {
                            "kind": "text",
                            "x": 0.7,
                            "y": 2.2,
                            "w": 12,
                            "h": 1.2,
                            "text": "现场标题",
                            "size": 40,
                            "color": "#F4F1EA",
                            "bold": True,
                            "name": "witty-title",
                        },
                    ],
                }
            ],
        }
        with tempfile.TemporaryDirectory() as tmp:
            path = str(Path(tmp) / "live.pptx")
            text = pptx_render(path, deck=json.dumps(payload, ensure_ascii=False))
            self.assertIn("custom", text)
            from pptx import Presentation

            slide = Presentation(path).slides[0]
            named = {shape.name: shape for shape in slide.shapes}
            self.assertIn("witty-title", named)
            self.assertIn("现场标题", named["witty-title"].text_frame.text)
            self.assertEqual(shape_fill_rgb(named["witty-bar"]), (0xC4, 0x5C, 0x26))
            html = Path(tmp) / "live.html"
            parsed = parse_html(html.read_text(encoding="utf-8"))
            self.assertTrue(parsed.slides[0].boxes)
            self.assertEqual(parsed.slides[0].boxes[1].text, "现场标题")

    def test_native_chart_and_html_roundtrip(self) -> None:
        payload = {
            "title": "图表稿",
            "theme": {"accent": "#C45C26", "bg": "#F3EFE6", "ink": "#1A2428"},
            "slides": [
                {
                    "bg": "#F3EFE6",
                    "boxes": [
                        {
                            "kind": "text",
                            "x": 0.62,
                            "y": 0.4,
                            "w": 12,
                            "h": 0.6,
                            "text": "完成率",
                            "size": 26,
                            "color": "#C45C26",
                            "bold": True,
                            "name": "witty-title",
                        },
                        {
                            "kind": "chart",
                            "chart": "column",
                            "x": 0.62,
                            "y": 1.2,
                            "w": 12,
                            "h": 5,
                            "text": "季度完成率",
                            "categories": ["一季度", "二季度", "三季度"],
                            "series": [{"name": "完成率", "values": [72, 81, 86]}],
                            "colors": ["#C45C26"],
                            "name": "witty-chart",
                        },
                    ],
                }
            ],
        }
        with tempfile.TemporaryDirectory() as tmp:
            path = str(Path(tmp) / "chart.pptx")
            text = pptx_render(path, deck=json.dumps(payload, ensure_ascii=False))
            self.assertIn(path, text)
            self.assertNotIn("只用了 kind 宏", text)
            from pptx import Presentation
            from pptx.enum.chart import XL_CHART_TYPE

            slide = Presentation(path).slides[0]
            chart_shape = next(shape for shape in slide.shapes if getattr(shape, "has_chart", False))
            self.assertEqual(chart_shape.name, "witty-chart")
            self.assertEqual(chart_shape.chart.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
            outline = pptx_outline(path)
            self.assertIn("完成率", outline)
            self.assertIn("[图表]", outline)
            html = Path(tmp) / "chart.html"
            parsed = parse_html(html.read_text(encoding="utf-8"))
            chart_box = next(box for box in parsed.slides[0].boxes if box.kind == "chart")
            self.assertEqual(chart_box.chart, "column")
            self.assertEqual(chart_box.categories, ["一季度", "二季度", "三季度"])
            self.assertEqual(chart_box.series[0].values, [72.0, 81.0, 86.0])
            out = str(Path(tmp) / "back.pptx")
            pptx_from_html(out, str(html))
            self.assertTrue(any(getattr(shape, "has_chart", False) for shape in Presentation(out).slides[0].shapes))

    def test_pie_doughnut_labels_are_dark_and_readable(self) -> None:
        """饼/环默认标签必须是深色大字，不能再写死白字 11pt（浅扇区上看不见）。"""
        payload = {
            "title": "饼图稿",
            "theme": "grid",
            "slides": [
                {
                    "boxes": [
                        {
                            "kind": "text",
                            "x": 0.62,
                            "y": 0.4,
                            "w": 12,
                            "h": 0.5,
                            "text": "NPS 分布",
                            "size": 26,
                            "bold": True,
                            "name": "witty-title",
                        },
                        {
                            "kind": "chart",
                            "chart": "doughnut",
                            "x": 0.62,
                            "y": 1.2,
                            "w": 12,
                            "h": 5,
                            "categories": ["推荐", "被动", "贬损"],
                            "series": [{"name": "人数", "values": [12, 5, 3]}],
                            "colors": ["#01706C", "#F4C542", "#C0392B"],
                            "name": "witty-nps",
                        },
                    ],
                }
            ],
        }
        with tempfile.TemporaryDirectory() as tmp:
            path = str(Path(tmp) / "nps.pptx")
            pptx_render(path, deck=json.dumps(payload, ensure_ascii=False), preview=False)
            from pptx import Presentation
            from pptx.util import Pt

            slide = Presentation(path).slides[0]
            chart = next(shape.chart for shape in slide.shapes if getattr(shape, "has_chart", False))
            labels = chart.plots[0].data_labels
            self.assertTrue(labels.show_percentage)
            self.assertTrue(labels.show_value)
            self.assertEqual(labels.font.size, Pt(14))
            rgb = labels.font.color.rgb
            self.assertEqual((rgb[0], rgb[1], rgb[2]), (18, 63, 60))

    def test_edit_slide_keeps_existing_color(self) -> None:
        payload = {
            "title": "青绿封面",
            "theme": {"cover_bg": "#163A3A", "cover_ink": "#F3EFE6", "ink": "#1A2428", "accent": "#C45C26"},
            "slides": [
                {
                    "bg": "#163A3A",
                    "boxes": [
                        {
                            "kind": "text",
                            "x": 0.7,
                            "y": 2.4,
                            "w": 12,
                            "h": 1.2,
                            "text": "青绿封面",
                            "size": 32,
                            "color": "#F3EFE6",
                            "bold": True,
                            "name": "witty-title",
                        }
                    ],
                },
                {
                    "bg": "#F3EFE6",
                    "boxes": [
                        {
                            "kind": "text",
                            "x": 0.62,
                            "y": 0.5,
                            "w": 12,
                            "h": 0.7,
                            "text": "要点",
                            "size": 26,
                            "color": "#C45C26",
                            "bold": True,
                            "name": "witty-title",
                        },
                        {
                            "kind": "bullets",
                            "x": 0.62,
                            "y": 1.4,
                            "w": 12,
                            "h": 4,
                            "items": ["一条"],
                            "size": 18,
                            "color": "#1A2428",
                            "name": "witty-body",
                        },
                    ],
                },
            ],
        }
        with tempfile.TemporaryDirectory() as tmp:
            path = str(Path(tmp) / "keep.pptx")
            pptx_render(path, deck=json.dumps(payload, ensure_ascii=False), preview=False)
            pptx_edit_slide(path, 1, title="改过的封面")
            pptx_edit_slide(path, 2, title="改过的要点", bullets="新一条")
            from pptx import Presentation

            slides = list(Presentation(path).slides)
            cover = next(shape for shape in slides[0].shapes if shape.name == "witty-title")
            body = next(shape for shape in slides[1].shapes if shape.name == "witty-title")
            self.assertEqual(cover.text_frame.text, "改过的封面")
            self.assertEqual(tuple(cover.text_frame.paragraphs[0].runs[0].font.color.rgb), (0xF3, 0xEF, 0xE6))
            self.assertEqual(tuple(body.text_frame.paragraphs[0].runs[0].font.color.rgb), (0xC4, 0x5C, 0x26))
            self.assertIn("新一条", pptx_outline(path))

    def test_kind_macro_render_warns(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = str(Path(tmp) / "macro.pptx")
            text = pptx_render(path, deck=json.dumps(_sample_deck(), ensure_ascii=False), preview=False)
            self.assertIn("只用了 kind 宏", text)

    def test_legacy_create_still_edits(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = str(Path(tmp) / "old.pptx")
            pptx_create(path, "封面", "内网")
            pptx_add_slide(path, "要点", "一条\n两条")
            pptx_edit_slide(path, 2, title="修订", bullets="新一条")
            outline = pptx_outline(path)
            self.assertIn("封面", outline)
            self.assertIn("修订", outline)
            self.assertIn("新一条", outline)

    def test_lint_overflow_and_contrast(self) -> None:
        payload = {
            "title": "脏稿",
            "theme": {"bg": "#FFFFFF", "ink": "#EEEEEE", "accent": "#FFFFFF"},
            "slides": [
                {
                    "bg": "#FFFFFF",
                    "boxes": [
                        {
                            "kind": "text",
                            "x": 12.5,
                            "y": 0.4,
                            "w": 3,
                            "h": 0.5,
                            "text": "溢出",
                            "size": 16,
                            "color": "#EEEEEE",
                            "name": "witty-title",
                        },
                        {
                            "kind": "text",
                            "x": 12.0,
                            "y": 0.4,
                            "w": 2,
                            "h": 0.5,
                            "text": "叠上",
                            "size": 16,
                            "color": "#EEEEEE",
                        },
                    ],
                }
            ],
        }
        report = pptx_check(deck=json.dumps(payload, ensure_ascii=False))
        self.assertIn("overflow", report)
        self.assertIn("contrast", report)
        self.assertIn("overlap", report)
        issues = lint_deck(parse_deck(payload))
        codes = {item.code for item in issues}
        self.assertIn("overflow", codes)
        self.assertIn("contrast", codes)

    def test_replace_slide_and_stored_theme(self) -> None:
        payload = {
            "title": "青绿",
            "theme": {"bg": "#F3EFE6", "ink": "#1A2428", "accent": "#C45C26", "accent2": "#163A3A"},
            "slides": [
                {
                    "bg": "#163A3A",
                    "boxes": [
                        {
                            "kind": "rect",
                            "x": 0,
                            "y": 0,
                            "w": 0.2,
                            "h": 7.5,
                            "fill": "#C45C26",
                            "name": "witty-bar",
                        },
                        {
                            "kind": "text",
                            "x": 0.7,
                            "y": 2.4,
                            "w": 12,
                            "h": 1,
                            "text": "青绿",
                            "size": 32,
                            "color": "#F3EFE6",
                            "bold": True,
                            "name": "witty-title",
                        },
                    ],
                },
                {
                    "bg": "#F3EFE6",
                    "boxes": [
                        {
                            "kind": "text",
                            "x": 0.62,
                            "y": 0.5,
                            "w": 12,
                            "h": 0.6,
                            "text": "旧页",
                            "size": 26,
                            "color": "#C45C26",
                            "bold": True,
                            "name": "witty-title",
                        }
                    ],
                },
            ],
        }
        with tempfile.TemporaryDirectory() as tmp:
            path = str(Path(tmp) / "live.pptx")
            pptx_render(path, deck=json.dumps(payload, ensure_ascii=False), preview=False)
            from pptx import Presentation

            stored = load_stored_theme(Presentation(path))
            self.assertIsNotNone(stored)
            self.assertEqual(stored.accent, (0xC4, 0x5C, 0x26))
            pptx_replace_slide(
                path,
                2,
                json.dumps(
                    {
                        "bg": "#F3EFE6",
                        "boxes": [
                            {
                                "kind": "text",
                                "x": 0.62,
                                "y": 0.5,
                                "w": 12,
                                "h": 0.6,
                                "text": "新页",
                                "size": 26,
                                "color": "#C45C26",
                                "bold": True,
                                "name": "witty-title",
                            }
                        ],
                    },
                    ensure_ascii=False,
                ),
            )
            pres = Presentation(path)
            self.assertEqual(len(list(pres.slides)), 2)
            page2 = next(shape for shape in pres.slides[1].shapes if shape.name == "witty-title")
            self.assertIn("新页", page2.text_frame.text)
            bar = next(shape for shape in pres.slides[0].shapes if shape.name == "witty-bar")
            self.assertEqual(shape_fill_rgb(bar), (0xC4, 0x5C, 0x26))
            pptx_add_slide(path, "追加", "跟主题")
            added = Presentation(path).slides[2]
            title = next(shape for shape in added.shapes if shape.name == "witty-title")
            self.assertEqual(tuple(title.text_frame.paragraphs[0].runs[0].font.color.rgb), (0x1A, 0x24, 0x28))

    def test_edit_named_box_and_add_page(self) -> None:
        payload = {
            "title": "青绿",
            "theme": {"bg": "#F3EFE6", "ink": "#1A2428", "accent": "#C45C26"},
            "slides": [
                {
                    "bg": "#F3EFE6",
                    "boxes": [
                        {
                            "kind": "text",
                            "x": 0.7,
                            "y": 0.7,
                            "w": 12,
                            "h": 0.7,
                            "text": "旧标题",
                            "size": 28,
                            "color": "#C45C26",
                            "bold": True,
                            "name": "witty-title",
                        }
                    ],
                }
            ],
        }
        with tempfile.TemporaryDirectory() as tmp:
            path = str(Path(tmp) / "edit.pptx")
            pptx_render(path, deck=json.dumps(payload, ensure_ascii=False), preview=False)
            listed = pptx_list_boxes(path, 1)
            self.assertIn("witty-title", listed)
            pptx_edit_box(path, 1, "witty-title", text="新标题")
            self.assertIn("新标题", pptx_outline(path))
            pptx_add_page(
                path,
                json.dumps(
                    {
                        "bg": "#F3EFE6",
                        "boxes": [
                            {
                                "kind": "text",
                                "x": 0.7,
                                "y": 0.7,
                                "w": 12,
                                "h": 0.7,
                                "text": "插入页",
                                "size": 28,
                                "color": "#1A2428",
                                "bold": True,
                                "name": "witty-title",
                            }
                        ],
                    },
                    ensure_ascii=False,
                ),
                index=1,
            )
            from pptx import Presentation

            slides = list(Presentation(path).slides)
            self.assertEqual(len(slides), 2)
            first = next(shape for shape in slides[0].shapes if shape.name == "witty-title")
            self.assertIn("插入页", first.text_frame.text)


if __name__ == "__main__":
    unittest.main()
