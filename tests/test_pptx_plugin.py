from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from witty_agent.kernel_surface import KERNEL_TOOLS, is_kernel_tool, is_kernel_tool_module
from witty_agent.loop import READONLY_TOOLS
from witty_agent.plan_mode import MUTATING_TOOLS
from witty_agent.plugins import list_plugins, plugin_owns
from witty_agent.plugins.pptx import (
    pptx_add_picture,
    pptx_add_slide,
    pptx_create,
    pptx_edit_slide,
    pptx_outline,
)
from witty_agent.prompts import get_prompt
from witty_agent.system_prompt import build_system_prompt
from witty_agent.tools import list_tools
from witty_agent.tool_surface import select_advertised_names


class PptxPluginTests(unittest.TestCase):
    def test_business_tools_are_not_kernel(self) -> None:
        names = {item.name for item in list_tools()}
        for name in (
            "pptx_create",
            "pptx_add_slide",
            "pptx_edit_slide",
            "pptx_add_picture",
            "pptx_outline",
            "pptx_render",
            "pptx_from_html",
            "pptx_themes",
            "pptx_check",
            "pptx_replace_slide",
            "pptx_list_boxes",
            "pptx_edit_box",
            "pptx_add_page",
        ):
            self.assertIn(name, names)
            self.assertFalse(is_kernel_tool(name))
            self.assertNotIn(name, KERNEL_TOOLS)
        self.assertFalse(is_kernel_tool_module("witty_agent.plugins.pptx"))
        self.assertTrue(plugin_owns("pptx_create"))
        self.assertIn("pptx", {item["name"] for item in list_plugins()["plugins"]})
        self.assertIn("pptx_outline", READONLY_TOOLS)
        self.assertIn("pptx_check", READONLY_TOOLS)
        self.assertIn("pptx_list_boxes", READONLY_TOOLS)
        self.assertIn("pptx_create", MUTATING_TOOLS)
        self.assertIn("pptx_replace_slide", MUTATING_TOOLS)
        self.assertIn("pptx_edit_box", MUTATING_TOOLS)
        self.assertIn("pptx_add_page", MUTATING_TOOLS)
        self.assertNotIn("pptx_outline", MUTATING_TOOLS)
        self.assertNotIn("pptx_check", MUTATING_TOOLS)

    def test_offline_create_edit_and_missing_image(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = str(Path(tmp) / "deck.pptx")
            pptx_create(path, "封面", "内网")
            pptx_add_slide(path, "要点", "一条\n两条")
            pptx_edit_slide(path, 2, title="修订", bullets="新一条")
            outline = pptx_outline(path)
            self.assertIn("封面", outline)
            self.assertIn("修订", outline)
            self.assertIn("新一条", outline)
            from pptx import Presentation
            from pptx.util import Inches

            deck = Presentation(path)
            self.assertGreater(deck.slide_width, Inches(13))
            self.assertGreaterEqual(len(list(deck.slides[0].shapes)), 4)
            self.assertTrue(any(shape.name == "witty-footer" for shape in deck.slides[1].shapes))
            missing = Path(tmp) / "no.png"
            with self.assertRaises(FileNotFoundError) as caught:
                pptx_add_picture(path, str(missing), title="图")
            self.assertIn(str(missing), str(caught.exception))
            with self.assertRaises(ValueError) as bad:
                pptx_edit_slide(path, 99, title="无")
            self.assertIn("99", str(bad.exception))

    def test_surface_and_prompt(self) -> None:
        shown = select_advertised_names(
            "帮我改这个 pptx",
            sorted(KERNEL_TOOLS) + ["pptx_create", "pptx_outline", "pptx_edit_slide"],
        )
        self.assertIn("pptx_outline", shown)
        self.assertIn("pptx_edit_slide", shown)
        self.assertNotIn(
            "pptx_create",
            select_advertised_names("review the auth module", sorted(KERNEL_TOOLS) + ["pptx_create"]),
        )
        text = build_system_prompt(
            ".",
            tool_names=["pptx_create", "pptx_outline"],
            skills=[],
            context_files=[],
            prompt="做个PPT",
        )
        self.assertIn("pptx_create", text)
        self.assertIn("slides", text)
        self.assertIn("不要调用公网美化", text)


if __name__ == "__main__":
    unittest.main()
