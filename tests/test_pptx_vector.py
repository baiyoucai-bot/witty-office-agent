"""矢量层：原生形状盒子、描边、形状内文字、带箭头的连接线。

重点盯三件事：几何顶点算得对、成稿落的是 PowerPoint 原生形状（WPS 里能编辑）、
三个渲染端（成稿 / 预览 / 快照）认的是同一组盒子。
"""

from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path

from witty_agent.plugins.pptx_kit import parse_deck
from witty_agent.plugins.pptx_kit.flex import solve_layout
from witty_agent.plugins.pptx_kit.lint import lint_deck
from witty_agent.plugins.pptx_kit.preview import render_html
from witty_agent.plugins.pptx_kit.render import write_pptx
from witty_agent.plugins.pptx_kit.schema import SHAPE_KINDS
from witty_agent.plugins.pptx_kit.shapes import (
    POLY_KINDS,
    clip_path,
    normalize_point,
    polygon_points,
    text_inset,
)

EPS = 1e-6


def _deck(boxes: list[dict], **slide) -> object:
    payload = {"title": "矢量", "theme": "grid", "slides": [{"kind": "custom", "boxes": boxes, **slide}]}
    return parse_deck(json.dumps(payload, ensure_ascii=False))


class ShapeGeometryTest(unittest.TestCase):
    def test_points_stay_inside_unit_box(self) -> None:
        """所有顶点都得落在 0..1 之内，否则形状会戳出盒子外。"""
        for kind in sorted(POLY_KINDS):
            for point in ("right", "left", "up", "down"):
                for w, h in ((2.0, 1.0), (1.0, 2.0), (1.0, 1.0), (6.0, 0.4)):
                    for fx, fy in polygon_points(kind, w, h, point):
                        self.assertGreaterEqual(fx, -EPS, f"{kind}/{point} 顶点越界")
                        self.assertLessEqual(fx, 1 + EPS, f"{kind}/{point} 顶点越界")
                        self.assertGreaterEqual(fy, -EPS, f"{kind}/{point} 顶点越界")
                        self.assertLessEqual(fy, 1 + EPS, f"{kind}/{point} 顶点越界")

    def test_notch_scales_with_absolute_size(self) -> None:
        """扁盒子的缺口按高度取，不能随宽度一起拉长成一根细刺。"""
        wide = polygon_points("chevron", 8.0, 1.0)
        # 缺口深度 = 0.5 * 高 = 0.5 英寸，占宽度的 1/16
        self.assertAlmostEqual(wide[5][0], 0.5 / 8.0, places=6)
        square = polygon_points("chevron", 1.0, 1.0)
        self.assertAlmostEqual(square[5][0], 0.5, places=6)

    def test_left_point_mirrors_right(self) -> None:
        right = polygon_points("arrow", 2.0, 1.0, "right")
        left = polygon_points("arrow", 2.0, 1.0, "left")
        self.assertEqual([(round(1 - x, 6), round(y, 6)) for x, y in right], [(round(x, 6), round(y, 6)) for x, y in left])

    def test_only_arrow_honours_point(self) -> None:
        """成稿里只有 arrow 四个朝向都有对应的原生形状。chevron / pentagon 固定朝右，
        菱形三角根本没朝向——预览和快照必须跟着固定，否则自查图会画出导不出来的样子。"""
        for kind in ("diamond", "triangle", "chevron", "pentagon"):
            for point in ("left", "up", "down"):
                self.assertEqual(
                    polygon_points(kind, 2.0, 1.0, point),
                    polygon_points(kind, 2.0, 1.0, "right"),
                    f"{kind} 不该跟着 point 变形",
                )
        self.assertNotEqual(polygon_points("arrow", 2.0, 1.0, "up"), polygon_points("arrow", 2.0, 1.0, "right"))

    def test_non_polygon_kinds_have_no_points(self) -> None:
        for kind in ("rect", "round", "oval", "text"):
            self.assertEqual(polygon_points(kind, 1.0, 1.0), [])

    def test_clip_path_is_css_polygon(self) -> None:
        path = clip_path("diamond", 1.0, 1.0)
        self.assertTrue(path.startswith("polygon("))
        self.assertEqual(path.count("%"), 8)
        self.assertEqual(clip_path("oval", 1.0, 1.0), "")

    def test_normalize_point_falls_back(self) -> None:
        self.assertEqual(normalize_point("LEFT"), "left")
        self.assertEqual(normalize_point("斜的"), "right")
        self.assertEqual(normalize_point(""), "right")

    def test_pointy_kinds_inset_text(self) -> None:
        """尖角形状得给文字让位，矩形不用——不然字压到斜边上。"""
        self.assertGreater(text_inset("diamond"), 0)
        self.assertGreater(text_inset("chevron"), 0)
        self.assertEqual(text_inset("rect"), 0.0)


class VectorSchemaTest(unittest.TestCase):
    def test_new_kinds_survive_parsing(self) -> None:
        kinds = ["oval", "diamond", "triangle", "chevron", "pentagon", "arrow"]
        deck = _deck([{"kind": k, "x": 0.5 + i, "y": 1.5, "w": 0.8, "h": 0.8, "fill": "#185FA5"} for i, k in enumerate(kinds)])
        self.assertEqual([b.kind for b in deck.slides[0].boxes], kinds)

    def test_aliases_map_to_canonical_kinds(self) -> None:
        deck = _deck([
            {"kind": "circle", "x": 1, "y": 1, "w": 1, "h": 1, "fill": "#111"},
            {"kind": "rhombus", "x": 3, "y": 1, "w": 1, "h": 1, "fill": "#111"},
        ])
        self.assertEqual([b.kind for b in deck.slides[0].boxes], ["oval", "diamond"])

    def test_stroke_and_point_round_trip(self) -> None:
        deck = _deck([{"kind": "rect", "x": 1, "y": 1, "w": 2, "h": 1, "stroke": "#185FA5", "stroke_w": 1.5, "point": "down"}])
        box = deck.slides[0].boxes[0]
        self.assertEqual(box.stroke, "#185FA5")
        self.assertAlmostEqual(box.stroke_w, 1.5)
        self.assertEqual(box.point, "down")

    def test_line_kind_is_not_swallowed_by_line_chart(self) -> None:
        """`kind: "line"` 是分隔线。以前会被折线图的简写吞掉，整条线变成空图表。"""
        deck = _deck([{"kind": "line", "x": 1, "y": 3, "w": 4, "h": 0.03, "fill": "#ccc"}])
        self.assertEqual(deck.slides[0].boxes[0].kind, "line")

    def test_line_with_data_is_still_a_chart(self) -> None:
        """带了数据的 `kind: "line"` 仍按折线图走，老稿子不能因为上一条修复而失效。"""
        deck = _deck([{
            "kind": "line", "x": 1, "y": 3, "w": 4, "h": 2,
            "categories": ["一", "二"], "series": [{"name": "量", "values": [1, 2]}],
        }])
        box = deck.slides[0].boxes[0]
        self.assertEqual(box.kind, "chart")
        self.assertEqual(box.chart, "line")


class VectorFlexTest(unittest.TestCase):
    def test_labeled_shape_gets_natural_height(self) -> None:
        """带文字的形状能自己量高，不用作者给 h。"""
        boxes = solve_layout(
            {"dir": "col", "children": [{"kind": "chevron", "fill": "#111", "color": "#fff", "text": "数据归集", "size": 15}]},
            13.333,
            7.5,
        )
        self.assertEqual(len(boxes), 1)
        self.assertGreater(boxes[0]["h"], 0.2)

    def test_bare_shape_has_no_natural_height(self) -> None:
        """没文字的纯色形状没有内容高度，作者得自己给 h 或 flex。"""
        boxes = solve_layout(
            {"dir": "col", "h": 3.0, "children": [{"kind": "diamond", "fill": "#111"}]},
            13.333,
            7.5,
        )
        self.assertLess(boxes[0]["h"], 0.05)

    def test_stroke_only_node_still_paints(self) -> None:
        """只给描边没给填充的容器也要落成盒子——架构图的空心框就是这么写的。"""
        boxes = solve_layout(
            {"dir": "row", "h": 1.0, "children": [{"stroke": "#185FA5", "stroke_w": 1.0}]},
            13.333,
            7.5,
        )
        self.assertEqual(len(boxes), 1)
        self.assertEqual(boxes[0]["kind"], "rect")


class VectorLintTest(unittest.TestCase):
    def test_labeled_shape_is_checked_for_contrast(self) -> None:
        """形状里的字也得查对比度，不然白字配浅底能一路蒙混过去。"""
        deck = _deck([
            {"kind": "chevron", "x": 0.62, "y": 2.0, "w": 3.0, "h": 0.9,
             "fill": "#F2F2F2", "color": "#FFFFFF", "text": "看不清", "size": 15},
            {"kind": "text", "x": 0.62, "y": 5.0, "w": 6.0, "h": 0.6, "text": "垫底", "size": 16},
        ])
        self.assertTrue(any(i.code == "contrast" for i in lint_deck(deck)))

    def test_bare_shape_skips_text_checks(self) -> None:
        """没文字的装饰形状不该被当正文查字号。"""
        deck = _deck([
            {"kind": "arrow", "x": 3.0, "y": 2.0, "w": 0.5, "h": 0.4, "fill": "#185FA5", "size": 8},
            {"kind": "text", "x": 0.62, "y": 5.0, "w": 6.0, "h": 0.6, "text": "垫底", "size": 16},
        ])
        self.assertFalse(any(i.code == "tiny_text" for i in lint_deck(deck)))

    def test_labeled_shape_overflow_is_reported(self) -> None:
        deck = _deck([
            {"kind": "round", "x": 0.62, "y": 2.0, "w": 1.4, "h": 0.3,
             "fill": "#185FA5", "color": "#FFFFFF", "text": "这一段字远远塞不进这个矮盒子里去", "size": 16},
            {"kind": "text", "x": 0.62, "y": 5.0, "w": 6.0, "h": 0.6, "text": "垫底", "size": 16},
        ])
        self.assertTrue(any(i.code == "text_overflow" for i in lint_deck(deck)))

    def test_shape_carrying_its_own_text_is_not_an_empty_card(self) -> None:
        """自己带文字的卡片不算空腔——内容就在它身上，不在子盒子里。"""
        deck = _deck([
            {"kind": "round", "x": 0.62, "y": 1.6, "w": 4.0, "h": 2.2,
             "fill": "#185FA5", "color": "#FFFFFF", "text": "整改闭环", "size": 18},
            {"kind": "text", "x": 6.0, "y": 5.0, "w": 6.0, "h": 0.6, "text": "垫底", "size": 16},
        ])
        self.assertFalse(any(i.code == "slack" for i in lint_deck(deck)))


class VectorRenderTest(unittest.TestCase):
    def _shapes(self, boxes: list[dict]):
        import pptx

        deck = _deck(boxes)
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "vec.pptx"
            write_pptx(deck, str(path))
            pres = pptx.Presentation(str(path))
            return {s.name: s for s in pres.slides[0].shapes}

    def test_vector_kinds_become_native_autoshapes(self) -> None:
        """必须落成 PowerPoint 原生形状：这样在 WPS 里还能拖点改色，而不是一张死图。"""
        from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

        want = {
            "oval": MSO_SHAPE.OVAL,
            "diamond": MSO_SHAPE.DIAMOND,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "chevron": MSO_SHAPE.CHEVRON,
            "pentagon": MSO_SHAPE.PENTAGON,
        }
        shapes = self._shapes([
            {"kind": k, "name": k, "x": 0.5 + i * 1.2, "y": 1.5, "w": 1.0, "h": 0.8, "fill": "#185FA5"}
            for i, k in enumerate(want)
        ])
        for kind, enum in want.items():
            self.assertEqual(shapes[kind].shape_type, MSO_SHAPE_TYPE.AUTO_SHAPE, f"{kind} 不是原生形状")
            self.assertEqual(shapes[kind].auto_shape_type, enum, f"{kind} 形状类型不对")

    def test_arrow_point_picks_the_right_shape(self) -> None:
        from pptx.enum.shapes import MSO_SHAPE

        shapes = self._shapes([
            {"kind": "arrow", "name": f"a-{p}", "point": p, "x": 0.5 + i * 1.2, "y": 1.5, "w": 1.0, "h": 0.8, "fill": "#185FA5"}
            for i, p in enumerate(("right", "left", "up", "down"))
        ])
        self.assertEqual(shapes["a-right"].auto_shape_type, MSO_SHAPE.RIGHT_ARROW)
        self.assertEqual(shapes["a-left"].auto_shape_type, MSO_SHAPE.LEFT_ARROW)
        self.assertEqual(shapes["a-up"].auto_shape_type, MSO_SHAPE.UP_ARROW)
        self.assertEqual(shapes["a-down"].auto_shape_type, MSO_SHAPE.DOWN_ARROW)

    def test_shape_text_lands_in_the_shape(self) -> None:
        """字要写进形状本身，不是另摞一个文本框——否则在 WPS 里拖形状字不跟着走。"""
        shapes = self._shapes([
            {"kind": "chevron", "name": "step", "x": 0.62, "y": 2.0, "w": 2.4, "h": 0.9,
             "fill": "#185FA5", "color": "#FFFFFF", "text": "数据归集", "size": 15},
        ])
        self.assertEqual(shapes["step"].text_frame.text, "数据归集")

    def test_stroke_draws_a_real_outline(self) -> None:
        from pptx.util import Pt

        shapes = self._shapes([
            {"kind": "rect", "name": "framed", "x": 0.62, "y": 2.0, "w": 2.4, "h": 0.9,
             "fill": "#FFFFFF", "stroke": "#185FA5", "stroke_w": 1.5},
        ])
        line = shapes["framed"].line
        self.assertEqual(line.width, Pt(1.5))
        self.assertEqual(str(line.color.rgb), "185FA5")

    def test_pointed_line_becomes_a_connector_with_arrowhead(self) -> None:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.oxml.ns import qn

        shapes = self._shapes([
            {"kind": "line", "name": "link", "x": 2.0, "y": 3.0, "w": 1.2, "h": 0.4,
             "fill": "#185FA5", "point": "right", "stroke_w": 2},
        ])
        link = shapes["link"]
        self.assertEqual(link.shape_type, MSO_SHAPE_TYPE.LINE)
        self.assertIsNotNone(link.line._get_or_add_ln().find(qn("a:tailEnd")), "连接线没有箭头")

    def test_plain_line_stays_a_bar(self) -> None:
        """不给 point 的 line 还是一条色条，老稿子的分隔线不能变成箭头。"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        shapes = self._shapes([
            {"kind": "line", "name": "rule", "x": 0.62, "y": 3.0, "w": 4.0, "h": 0.03, "fill": "#cccccc"},
        ])
        self.assertEqual(shapes["rule"].shape_type, MSO_SHAPE_TYPE.AUTO_SHAPE)


class VectorPreviewTest(unittest.TestCase):
    def test_preview_clips_polygons_and_rounds_ovals(self) -> None:
        html = render_html(_deck([
            {"kind": "diamond", "x": 1.0, "y": 2.0, "w": 1.0, "h": 1.0, "fill": "#185FA5"},
            {"kind": "oval", "x": 3.0, "y": 2.0, "w": 1.0, "h": 1.0, "fill": "#185FA5"},
        ]))
        self.assertIn("clip-path:polygon(", html)
        self.assertIn("border-radius:50%", html)

    def test_preview_keeps_vector_attrs_for_round_trip(self) -> None:
        """预览要把 stroke / point 写回 data-*，否则 HTML 回读会把矢量信息丢光。"""
        html = render_html(_deck([
            {"kind": "arrow", "x": 1.0, "y": 2.0, "w": 1.0, "h": 0.5,
             "fill": "#185FA5", "point": "down", "stroke": "#0C447C", "stroke_w": 1.0},
        ]))
        self.assertIn("data-point='down'", html)
        self.assertIn("data-stroke='#0C447C'", html)
        self.assertIn("data-stroke-w='1.0'", html)


if __name__ == "__main__":
    unittest.main()
